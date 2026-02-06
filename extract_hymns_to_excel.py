#!/usr/bin/env python3
"""
Extract all hymns from PowerPoint files and create an Excel report.
Lists hymn number, title, slide location, and slide count.
"""

import os
import re
from pptx import Presentation
import openpyxl
from openpyxl.styles import Font, Alignment, PatternFill
from openpyxl.utils import get_column_letter

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# Directories to search
SEARCH_DIRS = [
    os.path.join(BASE_DIR, "OneDrive_2026-02-05", "Holy Communion Services - Slides", "Malayalam HCS"),
    os.path.join(BASE_DIR, "OneDrive_2026-02-05", "Holy Communion Services - Slides", "English HCS"),
]


def find_all_pptx_files():
    """Find all PowerPoint files in the search directories."""
    pptx_files = []
    for d in SEARCH_DIRS:
        if os.path.isdir(d):
            for root, dirs, files in os.walk(d):
                for f in files:
                    if f.endswith(".pptx") and not f.startswith("~$"):
                        pptx_files.append(os.path.join(root, f))
    return pptx_files


def extract_hymn_number(text):
    """Extract hymn number from text like 'Hymn No. 123' or 'Hymn 123'."""
    # Look for "Hymn No. XXX" or "Hymn XXX"
    match = re.search(r'Hymn\s*(?:No\.?\s*)?(\d{1,3})\b', text, re.IGNORECASE)
    if match:
        return match.group(1)
    return None


def extract_hymn_title(text, hymn_num):
    """Extract hymn title from text, removing the hymn number part."""
    # Remove "Hymn No. XXX" part
    text = re.sub(r'Hymn\s*(?:No\.?\s*)?\d{1,3}\b', '', text, flags=re.IGNORECASE)
    # Clean up
    text = text.strip()
    # Remove leading punctuation like : or -
    text = re.sub(r'^[:\-\s]+', '', text)
    # Take only the first line or first 100 characters as title
    lines = text.split('\n')
    if lines:
        text = lines[0].strip()
    if len(text) > 100:
        text = text[:100]
    return text.strip()


def analyze_pptx_file(pptx_path):
    """
    Analyze a PowerPoint file and extract hymn information.
    Returns a list of hymn dictionaries with:
    - hymn_number: str
    - title: str
    - file_name: str
    - file_path: str
    - title_slide: int (1-based)
    - content_slides: list of int (1-based)
    - total_slides: int
    """
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"  ⚠ Could not open {pptx_path}: {e}")
        return []
    
    file_name = os.path.basename(pptx_path)
    file_path = pptx_path
    
    hymns = []
    current_hymn = None
    
    for slide_idx, slide in enumerate(prs.slides, start=1):
        # Extract all text from slide
        all_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += " " + shape.text_frame.text.strip()
        
        # Check if this slide contains a hymn number
        hymn_num = extract_hymn_number(all_text)
        
        if hymn_num:
            # This is likely a title slide for a new hymn
            if current_hymn and current_hymn['content_slides']:
                # Save the previous hymn
                hymns.append(current_hymn)
            
            # Start a new hymn
            title = extract_hymn_title(all_text, hymn_num)
            current_hymn = {
                'hymn_number': hymn_num,
                'title': title,
                'file_name': file_name,
                'file_path': file_path,
                'title_slide': slide_idx,
                'content_slides': [],
                'total_slides': 0
            }
        elif current_hymn:
            # This is likely a content slide for the current hymn
            # Check if it has Malayalam or English text content
            if len(all_text.strip()) > 10:  # Has meaningful content
                # Check if it's not a generic slide (like offertory, message, etc.)
                generic_keywords = ['message', 'offertory', 'confession', 'holy communion', 
                                   'thanksgiving', 'closing', 'opening']
                is_generic = any(keyword in all_text.lower() for keyword in generic_keywords)
                
                # Only add if it's not a known generic slide and has substantial text
                if not is_generic:
                    current_hymn['content_slides'].append(slide_idx)
    
    # Don't forget the last hymn
    if current_hymn and current_hymn['content_slides']:
        hymns.append(current_hymn)
    
    # Calculate total slides for each hymn
    for hymn in hymns:
        hymn['total_slides'] = len(hymn['content_slides'])
    
    return hymns


def clean_text_for_excel(text):
    """Remove illegal characters for Excel cells."""
    if not text:
        return ""
    # Remove control characters and other illegal characters
    cleaned = ''.join(char for char in text if ord(char) >= 32 or char in '\t\n\r')
    # Limit length to avoid issues
    if len(cleaned) > 1000:
        cleaned = cleaned[:1000] + "..."
    return cleaned


def create_excel_report(all_hymns, output_file):
    """Create an Excel file with hymn information."""
    # Sort by hymn number (as integer)
    all_hymns.sort(key=lambda x: int(x['hymn_number']))
    
    # Create workbook
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Hymns Report"
    
    # Define headers
    headers = ['Hymn Number', 'Title', 'File Name', 'Title Slide', 'Content Slides', 'Total Slides']
    ws.append(headers)
    
    # Style headers
    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF", size=12)
    
    for col in range(1, len(headers) + 1):
        cell = ws.cell(row=1, column=col)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal='center', vertical='center')
    
    # Add data
    for hymn in all_hymns:
        # Format content slides as a range or list
        if hymn['content_slides']:
            content_slides_str = f"{min(hymn['content_slides'])}-{max(hymn['content_slides'])}"
        else:
            content_slides_str = "N/A"
        
        row_data = [
            int(hymn['hymn_number']),
            clean_text_for_excel(hymn['title']),
            clean_text_for_excel(hymn['file_name']),
            hymn['title_slide'],
            content_slides_str,
            hymn['total_slides']
        ]
        ws.append(row_data)
    
    # Auto-adjust column widths
    for col in range(1, len(headers) + 1):
        column_letter = get_column_letter(col)
        max_length = 0
        
        for cell in ws[column_letter]:
            try:
                if len(str(cell.value)) > max_length:
                    max_length = len(str(cell.value))
            except:
                pass
        
        adjusted_width = min(max_length + 2, 50)  # Cap at 50
        ws.column_dimensions[column_letter].width = adjusted_width
    
    # Center align numeric columns
    for row in ws.iter_rows(min_row=2, max_row=ws.max_row):
        row[0].alignment = Alignment(horizontal='center')  # Hymn Number
        row[3].alignment = Alignment(horizontal='center')  # Title Slide
        row[4].alignment = Alignment(horizontal='center')  # Content Slides
        row[5].alignment = Alignment(horizontal='center')  # Total Slides
    
    # Save workbook
    wb.save(output_file)
    print(f"\n✓ Excel report saved to: {output_file}")


def main():
    print("=" * 70)
    print("Extracting Hymn Information from PowerPoint Files")
    print("=" * 70)
    
    # Find all PowerPoint files
    print("\n🔍 Searching for PowerPoint files...")
    pptx_files = find_all_pptx_files()
    print(f"  Found {len(pptx_files)} PowerPoint files")
    
    # Extract hymn information from each file
    all_hymns = []
    for pptx_file in pptx_files:
        print(f"\n📄 Analyzing: {os.path.basename(pptx_file)}")
        hymns = analyze_pptx_file(pptx_file)
        print(f"  Found {len(hymns)} hymns")
        all_hymns.extend(hymns)
    
    print(f"\n📊 Total hymns found across all files: {len(all_hymns)}")
    
    # Create Excel report
    output_file = os.path.join(BASE_DIR, "hymns_report.xlsx")
    create_excel_report(all_hymns, output_file)
    
    # Print summary
    print("\n" + "=" * 70)
    print("Summary:")
    print("=" * 70)
    print(f"Total PowerPoint files analyzed: {len(pptx_files)}")
    print(f"Total hymns found: {len(all_hymns)}")
    print(f"Hymn numbers range: {min(int(h['hymn_number']) for h in all_hymns)} - {max(int(h['hymn_number']) for h in all_hymns)}")
    print(f"Output file: {output_file}")
    print("=" * 70)


if __name__ == "__main__":
    main()
