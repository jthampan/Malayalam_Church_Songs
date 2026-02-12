#!/usr/bin/env python3
"""
Extract Malayalam hymns from PowerPoint files and create an Excel report.
Lists hymn number, title (from Manglish or Malayalam content), slide location, and slide count.
"""

import os
import re
import json
from pptx import Presentation
import openpyxl
from openpyxl.styles import Font, Alignment, PatternFill
from openpyxl.utils import get_column_letter

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PARENT_DIR = os.path.dirname(BASE_DIR)

# Only search Malayalam directories
MALAYALAM_SEARCH_DIRS = [
    os.path.join(PARENT_DIR, "onedrive_git_local", "Holy Communion Services - Slides", "Malayalam HCS"),
]


def find_all_pptx_files():
    """Find all PowerPoint files in Malayalam directories."""
    pptx_files = []
    for d in MALAYALAM_SEARCH_DIRS:
        if os.path.isdir(d):
            for root, dirs, files in os.walk(d):
                for f in files:
                    if f.endswith(".pptx") and not f.startswith("~$"):
                        pptx_files.append(os.path.join(root, f))
    return pptx_files


def extract_hymn_number_and_title(text):
    """Extract hymn number and/or title from a single text block."""
    hymn_num = None
    title = None
    
    # Look for explicit hymn number patterns
    patterns = [
        r'Hymn\s*No\.?\s*(\d{1,3})\b',
        r'Song\s*No\.?\s*(\d{1,3})\b',
        r'Hymn\s*[-–]\s*(\d{1,3})\b',
    ]
    for pattern in patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            hymn_num = match.group(1)
            break
    
    # Look for title patterns (e.g., "Holy Communion Hymn – Title Text")
    title_patterns = [
        r'(?:Holy\s+Communion\s+)?Hymn\s*[-–]\s*([A-Za-z][A-Za-z\s]+?)(?:\s*\d|$)',
        r'(?:Opening|Closing|Offertory|Communion)\s+Hymn\s*[-–]\s*([A-Za-z][A-Za-z\s]+?)(?:\s*\d|$)',
    ]
    for pattern in title_patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            title = match.group(1).strip()
            # Clean up the title
            if title and len(title) > 5:
                break
    
    return hymn_num, title


def extract_hymn_number(text):
    """Extract hymn number from a single text block (backward compatibility)."""
    hymn_num, _ = extract_hymn_number_and_title(text)
    return hymn_num


def extract_hymn_number_from_slide(slide):
    """Extract hymn number and/or title from a slide header."""
    hymn_num = None
    title = None
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        
        num, ttl = extract_hymn_number_and_title(text)
        if num:
            hymn_num = num
        if ttl:
            title = ttl
        
        # If we found something, return it
        if hymn_num or title:
            return hymn_num, title
    
    return None, None


def clean_text_for_excel(text):
    """Remove illegal characters for Excel cells."""
    if not text:
        return ""
    # Remove control characters and other illegal characters
    cleaned = ''.join(char for char in text if ord(char) >= 32 or char in '\t\n\r')
    # Limit length to avoid issues
    if len(cleaned) > 200:
        cleaned = cleaned[:200] + "..."
    return cleaned


def extract_title_from_content(slide):
    """Extract title from slide content (Manglish only, first 3 words from lyrics)."""
    content_lines = []
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                
                # Skip if it's just a hymn number reference
                if re.match(r'^\d{1,3}$', text.strip()):
                    continue
                
                # Check for Malayalam Unicode characters (we want Manglish, not Malayalam script)
                malayalam_chars = sum(1 for c in text if '\u0D00' <= c <= '\u0D7F')
                if malayalam_chars > 0:  # Reject ANY Malayalam Unicode
                    continue
                
                # Skip text with ANY weird/corrupted characters
                # Only allow: letters, numbers, spaces, common punctuation, and whitespace
                allowed_special = set("-'!?,. \n\r\t\x0b")  # Added newline, tab, vertical tab
                has_bad_chars = False
                for c in text:
                    if not (c.isalnum() or c in allowed_special):
                        has_bad_chars = True
                        break
                
                if has_bad_chars:
                    continue
                if has_bad_chars:
                    continue
                
                # Check if text looks like Manglish (mix of English alphabet with Malayalam pronunciation)
                # Good Manglish should have mostly ASCII letters with some diacritics
                total_alpha_chars = sum(1 for c in text if c.isalpha())
                ascii_letter_count = sum(1 for c in text if c.isascii() and c.isalpha())
                if total_alpha_chars > 0 and ascii_letter_count < total_alpha_chars * 0.9:  # Less than 90% ASCII letters
                    continue
                
                # Must have reasonable vowels (to avoid gibberish)
                vowels = sum(1 for c in text.lower() if c in 'aeiou')
                if ascii_letter_count > 0 and vowels < ascii_letter_count * 0.2:  # Less than 20% vowels
                    continue
                
                # Skip pure English text (short common English words)
                words = text.split()
                if len(words) >= 2:
                    common_english = ['ride', 'on', 'the', 'and', 'or', 'but', 'in', 'to', 'for', 'of', 'with', 'at', 'by', 'from']
                    english_word_count = sum(1 for w in words[:3] if w.lower() in common_english)
                    if english_word_count >= 2:  # At least 2 common English words in first 3
                        continue
                
                # Remove leading hymn numbers from text
                text = re.sub(r'^\d{1,3}\s+', '', text)
                    
                # Skip common labels, section headers, and metadata
                skip_patterns = [
                    r'^(hymn|offertory|communion|confession|thanksgiving|opening|closing)',
                    r'^holy\s+communion',  # Any "Holy Communion" text
                    r'^(message|prayer|scripture|reading)',
                    r'^\d+\s*$',  # Just numbers
                    r'^slide\s+\d+',
                    r'^page\s+\d+',
                    r'^uen\s*[-–]\s*s\d+',  # UEN codes
                    r'^\w+\$\w+\s+cm',  # Encoded text like "B«nŠb» cm"
                    r'^song\s+no',  # Skip "Song No. 331" etc
                    r'^theme:',
                    r'^\d+\s+[a-z]+\s+\d{4}',  # Dates
                    r'order\s+of\s+worship',
                    r'sacred\s+music',
                    r'choir\s+dedication',
                    r'^dedication\s*[-–]',  # Dedication labels
                    r'^easter\s+sunday',  # Easter event descriptions
                    r'^good\s+friday',
                    r'^palm\s+sunday',
                    r'^maundy\s+thursday',
                    r'^she\s+(only|came)',  # English descriptive text
                    r'^he\s+(turned|came)',
                    r'^ride\s+on',  # English hymn titles
                ]
                
                should_skip = False
                for pattern in skip_patterns:
                    if re.search(pattern, text.lower()):
                        should_skip = True
                        break
                
                if should_skip:
                    continue
                
                # Check if text looks like Manglish (mix of English alphabet with Malayalam pronunciation)
                # Good Manglish should have mostly ASCII letters with some diacritics
                ascii_letter_count = sum(1 for c in text if c.isascii() and c.isalpha())
                if total_alpha_chars > 0 and ascii_letter_count < total_alpha_chars * 0.5:  # Less than 50% ASCII letters
                    continue
                
                # This looks like actual Manglish song content - add it
                if len(text) > 3:  # At least 4 characters
                    content_lines.append(text)
                    
                    # Get first meaningful line
                    if len(content_lines) >= 1:
                        break
        
        if len(content_lines) >= 1:
            break
    
    # Extract first 3 words from the content as title
    if content_lines:
        full_text = content_lines[0]
        # Split into words and take first 3
        words = full_text.split()
        if len(words) >= 3:
            return ' '.join(words[:3])
        else:
            return full_text
    
    return ""


def analyze_pptx_file(pptx_path):
    """
    Analyze a PowerPoint file and extract Malayalam hymn information.
    Returns a list of hymn dictionaries.
    """
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"  ⚠ Could not open {pptx_path}: {e}")
        return []
    
    file_name = os.path.basename(pptx_path)
    file_path = pptx_path
    
    hymns = []
    current_hymn = None
    last_hymn_num = None
    slides_list = list(prs.slides)
    
    for slide_idx, slide in enumerate(slides_list, start=1):
        # Extract all text from slide
        all_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += " " + shape.text_frame.text.strip()
        
        # Skip slides with very little content
        if len(all_text.strip()) < 5:
            if current_hymn and len(current_hymn['content_slides']) > 0:
                # Empty slide might signal end of hymn
                hymns.append(current_hymn)
                current_hymn = None
                last_hymn_num = None
            continue
        
        # Skip slides that are clearly NOT hymns (prayers, readings, agenda, etc.)
        non_hymn_patterns = [
            r'\bprayer\b',
            r'\bresponse\b',
            r'\b[LC]\s*[-–]\s*for\b',  # Leader/Congregation prayers (L - For, C - We)
            r'\bleader\b',
            r'\bcongregation\b',
            r'\bdedication\b.*\d+\s+of\s+\d+',  # "Dedication: 4 of 11"
            r'\bintercessory\s+prayer\b',
            r'\bthanksgiving\s+prayer\b',
            r'\breading\b',
            r'\bscripture\b',
            r'hymns?\s+list',  # "Hymns List" index/agenda pages
            r'b/?a:',  # "B/A: 363" in agenda
            r'youtube\.be',  # Links in agenda pages
            r'theme:',  # "Theme: The Resurrection..."
        ]
        
        is_non_hymn = any(re.search(pattern, all_text, re.IGNORECASE) for pattern in non_hymn_patterns)
        if is_non_hymn:
            # This is a prayer/reading slide, not a hymn - skip it
            if current_hymn and len(current_hymn['content_slides']) > 0:
                hymns.append(current_hymn)
                current_hymn = None
                last_hymn_num = None
            continue
        
        # Look for hymn number and/or title on this slide
        hymn_num, header_title = extract_hymn_number_from_slide(slide)
        
        # Get content from this slide
        content_text = extract_title_from_content(slide)
        
        # Check if this is a section divider (has keywords but minimal actual content)
        section_keywords = ['opening', 'offertory', 'communion', 'confession', 'thanksgiving', 'closing', 'message']
        has_section_keyword = any(keyword in all_text.lower() for keyword in section_keywords)
        
        # Check if slide has "Hymn" or "Keerthanam" explicitly
        has_hymn_label = bool(re.search(r'\bhymn\b|\bkeerthanam\b', all_text, re.IGNORECASE))
        
        # Determine if this is a start of a new hymn
        is_new_hymn = False
        
        if hymn_num or header_title:
            # We found a hymn number or title
            if not current_hymn:
                # No current hymn, start a new one
                is_new_hymn = True
            elif hymn_num and current_hymn['hymn_number'] and hymn_num != current_hymn['hymn_number']:
                # Different hymn number, save current and start new
                if current_hymn['content_slides']:
                    hymns.append(current_hymn)
                is_new_hymn = True
            elif header_title and not hymn_num and not current_hymn['hymn_number']:
                # New title-only hymn following another title-only hymn
                if current_hymn['content_slides']:
                    hymns.append(current_hymn)
                is_new_hymn = True
            # else: same hymn, continue with current hymn
        
        if is_new_hymn:
            # Try to get title: prefer header title, then content title
            title = header_title if header_title else content_text
            
            # If no good title, try next few slides (up to 5 slides ahead)
            if not title or len(title) < 5:
                for next_idx in range(slide_idx, min(slide_idx + 5, len(slides_list))):
                    if next_idx < len(slides_list):
                        next_slide = slides_list[next_idx]
                        next_title = extract_title_from_content(next_slide)
                        if next_title and len(next_title) > 5:
                            title = next_title
                            break
            
            # Start a new hymn (hymn_number can be None for title-only hymns)
            current_hymn = {
                'hymn_number': hymn_num,  # Can be None
                'title': title,
                'file_name': file_name,
                'file_path': file_path,
                'title_slide': slide_idx,
                'content_slides': [],
                'total_slides': 0,
                'first_content_slide': None
            }
            last_hymn_num = hymn_num
            
            # Add this slide as content if it has meaningful content
            if content_text and len(content_text) > 5:
                current_hymn['content_slides'].append(slide_idx)
                if not current_hymn['first_content_slide']:
                    current_hymn['first_content_slide'] = slide_idx
                    # Use this as title if we don't have one yet
                    if not current_hymn['title']:
                        current_hymn['title'] = content_text
        
        elif current_hymn:
            # We're tracking a hymn, add this slide if it has content
            if content_text and len(content_text) > 5:
                # Check if this looks like it might be a different section
                # If it has a section keyword AND a new hymn number/title, it's probably new
                if has_section_keyword and (hymn_num or header_title) and (not current_hymn['hymn_number'] or hymn_num != current_hymn['hymn_number']):
                    # This is a new section/hymn
                    if current_hymn['content_slides']:
                        hymns.append(current_hymn)
                    
                    # Try to get title: prefer header title, then content
                    title = header_title if header_title else content_text
                    if not title or len(title) < 5:
                        for next_idx in range(slide_idx, min(slide_idx + 3, len(slides_list))):
                            if next_idx < len(slides_list):
                                next_slide = slides_list[next_idx]
                                next_title = extract_title_from_content(next_slide)
                                if next_title and len(next_title) > 5:
                                    title = next_title
                                    break
                    
                    current_hymn = {
                        'hymn_number': hymn_num,
                        'title': title,
                        'file_name': file_name,
                        'file_path': file_path,
                        'title_slide': slide_idx,
                        'content_slides': [],
                        'total_slides': 0,
                        'first_content_slide': None
                    }
                    last_hymn_num = hymn_num
                    if content_text and len(content_text) > 5:
                        current_hymn['content_slides'].append(slide_idx)
                        current_hymn['first_content_slide'] = slide_idx
                        if not current_hymn['title']:
                            current_hymn['title'] = content_text
                else:
                    # Continue with current hymn
                    current_hymn['content_slides'].append(slide_idx)
                    
                    # Update title if we don't have one yet
                    if not current_hymn['title']:
                        current_hymn['title'] = content_text
                        current_hymn['first_content_slide'] = slide_idx
            else:
                # Slide has no meaningful content, might be end of hymn
                # Check next slide behavior
                pass
    
    # Don't forget the last hymn
    if current_hymn and current_hymn['content_slides']:
        hymns.append(current_hymn)
    
    # Calculate total slides for each hymn
    for hymn in hymns:
        hymn['total_slides'] = len(hymn['content_slides'])
    
    # Remove duplicates (same hymn number/title in same file)
    seen = set()
    unique_hymns = []
    for hymn in hymns:
        # Use both hymn_number and title for uniqueness (handle None values)
        key = (hymn['hymn_number'], hymn.get('title', ''), hymn['file_name'])
        if key not in seen:
            seen.add(key)
            unique_hymns.append(hymn)
    
    return unique_hymns


def create_excel_report(all_hymns, output_file):
    """Create an Excel file with hymn information from Malayalam files."""
    # Load KK hymn mapping JSON
    kk_json_path = os.path.join(BASE_DIR, "kk_hymn_mapping.json")
    kk_hymns = {}
    if os.path.exists(kk_json_path):
        with open(kk_json_path, 'r', encoding='utf-8') as f:
            kk_hymns = json.load(f)
        print(f"\n📋 Loaded {len(kk_hymns)} hymns from kk_hymn_mapping.json")
    else:
        print(f"\n⚠ Warning: kk_hymn_mapping.json not found at {kk_json_path}")
    
    # Add KK hymns to all_hymns list
    kk_file_path = os.path.join(PARENT_DIR, "onedrive_git_local", "Holy Communion Services - Slides", "Malayalam HCS", "Hymns_malayalam_KK.pptx")
    for hymn_num, title in kk_hymns.items():
        all_hymns.append({
            'hymn_number': hymn_num,
            'title': title,
            'file_name': 'Hymns_malayalam_KK.pptx',
            'file_path': kk_file_path,
            'title_slide': None,
            'content_slides': [],
            'total_slides': 0,
            'first_content_slide': None
        })
    
    # Sort: hymns with numbers first (by number), then title-only hymns (alphabetically)
    def sort_key(x):
        if x['hymn_number']:
            return (0, int(x['hymn_number']), x['file_name'])
        else:
            return (1, x.get('title', ''), x['file_name'])
    
    all_hymns.sort(key=sort_key)
    
    # Create workbook
    wb = openpyxl.Workbook()
    
    # ============= TAB 1: Sort by Hymn Number =============
    ws1 = wb.active
    ws1.title = "Sort by Hymn Number"
    
    # Define headers
    headers = ['Hymn Number', 'Title', 'Slide Name']
    ws1.append(headers)
    
    # Style headers
    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF", size=12)
    
    for col in range(1, len(headers) + 1):
        cell = ws1.cell(row=1, column=col)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal='center', vertical='center')
    
    # Add data sorted by hymn number
    for hymn in all_hymns:
        row_data = [
            int(hymn['hymn_number']) if hymn['hymn_number'] else '',
            clean_text_for_excel(hymn['title']),
            clean_text_for_excel(hymn['file_name'])
        ]
        ws1.append(row_data)
    
    # Auto-adjust column widths
    for col in range(1, len(headers) + 1):
        column_letter = get_column_letter(col)
        max_length = 0
        
        for cell in ws1[column_letter]:
            try:
                if len(str(cell.value)) > max_length:
                    max_length = len(str(cell.value))
            except:
                pass
        
        adjusted_width = min(max_length + 2, 60)
        ws1.column_dimensions[column_letter].width = adjusted_width
    
    # Center align hymn number column
    for row in ws1.iter_rows(min_row=2, max_row=ws1.max_row):
        row[0].alignment = Alignment(horizontal='center')
    
    # ============= TAB 2: Sort by Filename =============
    ws2 = wb.create_sheet(title="Sort by Filename")
    
    # Headers
    ws2.append(headers)
    for col in range(1, len(headers) + 1):
        cell = ws2.cell(row=1, column=col)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal='center', vertical='center')
    
    # Sort by filename
    all_hymns_by_file = sorted(all_hymns, key=lambda x: (x['file_name'], int(x['hymn_number']) if x['hymn_number'] else 9999))
    
    # Add data sorted by filename
    for hymn in all_hymns_by_file:
        row_data = [
            int(hymn['hymn_number']) if hymn['hymn_number'] else '',
            clean_text_for_excel(hymn['title']),
            clean_text_for_excel(hymn['file_name'])
        ]
        ws2.append(row_data)
    
    # Auto-adjust column widths
    for col in range(1, len(headers) + 1):
        column_letter = get_column_letter(col)
        max_length = 0
        
        for cell in ws2[column_letter]:
            try:
                if len(str(cell.value)) > max_length:
                    max_length = len(str(cell.value))
            except:
                pass
        
        adjusted_width = min(max_length + 2, 60)
        ws2.column_dimensions[column_letter].width = adjusted_width
    
    # Center align hymn number column
    for row in ws2.iter_rows(min_row=2, max_row=ws2.max_row):
        row[0].alignment = Alignment(horizontal='center')
    
    # Save workbook
    wb.save(output_file)
    print(f"\n✓ Excel report saved to: {output_file}")
    print(f"  Tab 1: Sort by Hymn Number ({len(all_hymns)} entries)")
    print(f"  Tab 2: Sort by Filename ({len(all_hymns_by_file)} entries)")


def main():
    print("=" * 70)
    print("Extracting Malayalam Hymn Information from PowerPoint Files")
    print("(Including both Malayalam script and Manglish)")
    print("=" * 70)
    
    # Find all PowerPoint files
    print("\n🔍 Searching for Malayalam PowerPoint files...")
    pptx_files = find_all_pptx_files()
    print(f"  Found {len(pptx_files)} PowerPoint files")
    
    # Extract hymn information from each file
    all_hymns = []
    for pptx_file in pptx_files:
        print(f"\n📄 Analyzing: {os.path.basename(pptx_file)}")
        hymns = analyze_pptx_file(pptx_file)
        print(f"  Found {len(hymns)} hymns")
        all_hymns.extend(hymns)
    
    print(f"\n📊 Total hymns found across all files: {len(all_hymns)}")
    
    if not all_hymns:
        print("\n⚠ No hymns found!")
        return
    
    # Create Excel report (will deduplicate inside)
    output_file = os.path.join(BASE_DIR, "malayalam_hymns_report.xlsx")
    create_excel_report(all_hymns, output_file)
    
    # Count unique hymn numbers and title-only hymns for summary
    hymns_with_numbers = [h for h in all_hymns if h['hymn_number']]
    hymns_without_numbers = [h for h in all_hymns if not h['hymn_number']]
    unique_hymn_nums = len(set(h['hymn_number'] for h in hymns_with_numbers))
    
    # Print summary
    print("\n" + "=" * 70)
    print("Summary:")
    print("=" * 70)
    print(f"Total PowerPoint files analyzed: {len(pptx_files)}")
    print(f"Total hymn entries found: {len(all_hymns)}")
    print(f"  - Hymns with numbers: {len(hymns_with_numbers)} ({unique_hymn_nums} unique)")
    print(f"  - Title-only hymns: {len(hymns_without_numbers)}")
    if hymns_with_numbers:
        print(f"Hymn numbers range: {min(int(h['hymn_number']) for h in hymns_with_numbers)} - {max(int(h['hymn_number']) for h in hymns_with_numbers)}")
    print(f"Output file: {output_file}")
    print(f"Note: Excel report shows ALL instances (same hymn may appear in multiple files)")
    print("=" * 70)


if __name__ == "__main__":
    main()
