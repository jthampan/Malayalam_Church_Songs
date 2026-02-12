#!/usr/bin/env python3
"""
Malayalam Church Songs PPT Generator
====================================
Searches for Malayalam songs by hymn number in existing PPT files, then generates
a new presentation using the template from "4 Jan 2026" Malayalam HCS PPT.

Sections supported:
  - Opening Song
  - ThanksGiving Prayers  
  - Offertory (with QR code)
  - Message (title slide only)
  - Confession
  - Holy Communion (with Holy Communion image)
  - Closing Hymn

Usage:
    python3 generate_malayalam_hcs_ppt.py --batch songs.txt "Output Name.pptx"
    
    songs.txt format:
        hymn_num|label|title_hint
        # Lines starting with # are comments
        # Use 'Message' label for message-only title slide

Images are extracted and saved to the 'images/' folder.
"""

import os
import sys
import re
import json
from kk_hymn_search import find_hymn_in_kk_pptx
from copy import deepcopy
from datetime import datetime
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
import subprocess
import shutil


# ═══════════════════════════════════════════════════════════════════════════════
# CONFIGURATION
# ═══════════════════════════════════════════════════════════════════════════════

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PARENT_DIR = os.path.dirname(BASE_DIR)
IMAGES_DIR = os.path.join(PARENT_DIR, "images")

# Create images directory if it doesn't exist
os.makedirs(IMAGES_DIR, exist_ok=True)

# Load KK hymn mapping
KK_HYMN_MAPPING = {}
kk_mapping_path = os.path.join(BASE_DIR, "kk_hymn_mapping.json")
if os.path.exists(kk_mapping_path):
    with open(kk_mapping_path, 'r', encoding='utf-8') as f:
        KK_HYMN_MAPPING = json.load(f)

# Path to bundled/local hymn files folder
# NOTE: For exe builds, this folder is bundled at BUILD TIME by build_exe.py
ONEDRIVE_GIT_LOCAL = os.path.join(PARENT_DIR, "onedrive_git_local")

# ═══════════════════════════════════════════════════════════════════════════════
# FILE PATHS
# ═══════════════════════════════════════════════════════════════════════════════

TEMPLATE_PPT = os.path.join(
    ONEDRIVE_GIT_LOCAL,
    "Holy Communion Services - Slides",
    "Malayalam HCS",
    "2026- Mal",
    "4 Jan 2026.pptx",
)

# Hymns PPT - look in onedrive_git_local first, then fall back to BASE_DIR
HYMNS_PPT_LOCATIONS = [
    os.path.join(ONEDRIVE_GIT_LOCAL, "Holy Communion Services - Slides", "Malayalam HCS", "Hymns_malayalam_KK.pptx"),
    os.path.join(BASE_DIR, "Hymns_malayalam_KK.pptx"),
]
HYMNS_PPT = next((path for path in HYMNS_PPT_LOCATIONS if os.path.exists(path)), HYMNS_PPT_LOCATIONS[-1])
PDF_FILE = os.path.join(BASE_DIR, "Kristeeya Keerthanagal.pdf")

# Image paths
HOLY_COMMUNION_IMAGE = os.path.join(IMAGES_DIR, "holy_communion.jpg")

# Directories to search for existing PPTs
SEARCH_DIRS = [
    ONEDRIVE_GIT_LOCAL,
    BASE_DIR,
]

def resolve_image_path(filename):
    """Find an image in the source folder first, then fallback to packaged images."""
    candidates = [
        os.path.join(os.getcwd(), "images", filename),
        os.path.join(BASE_DIR, "images", filename),
    ]
    for path in candidates:
        if os.path.exists(path):
            return path
    return candidates[0]

def find_template_ppt():
    """Find the Malayalam template PPT in the selected source folder or fallback paths."""
    template_name = "4 Jan 2026.pptx"
    # Search in: user-provided folder, PARENT_DIR, onedrive_git_local
    search_roots = [os.getcwd(), PARENT_DIR, ONEDRIVE_GIT_LOCAL]

    # Prefer the current working directory (set by the GUI to the source folder)
    for root in search_roots:
        if not root or not os.path.isdir(root):
            continue
        for dirpath, _, filenames in os.walk(root):
            if template_name in filenames:
                return os.path.join(dirpath, template_name)

    # Final fallback to the static path if it exists
    if os.path.exists(TEMPLATE_PPT):
        return TEMPLATE_PPT

    return None

# ─── Font / layout settings (from template analysis) ─────────────────────────
TITLE_FONT = "Gabriola"
TITLE_SIZE = Pt(60)

HEADER_FONT = "Segoe UI"
ROMANIZED_FONT = "Segoe UI"
ROMANIZED_SIZE = Pt(23)
MALAYALAM_FONT = "Noto Sans Malayalam"
MALAYALAM_SIZE = Pt(22)
MALAYALAM_COLOR = RGBColor(0x00, 0x00, 0x00)

FOOTER_FONT = "Segoe UI"
FOOTER_SIZE = Pt(14)

# Title bar styling (from template)
TITLE_BAR_HEIGHT = Emu(486000)
TITLE_BAR_COLOR = RGBColor(232, 211, 211)  # #E8D3D3
TITLE_BAR_FONT_SIZE = Pt(24)

# Maximum lines per content slide
MAX_LINES_PER_SLIDE = 8

# Holy Communion image position (from template)
HC_IMAGE_LEFT = Inches(1.673)
HC_IMAGE_TOP = Inches(0.772)
HC_IMAGE_WIDTH = Inches(6.654)
HC_IMAGE_HEIGHT = Inches(4.437)


# ═══════════════════════════════════════════════════════════════════════════════
# IMAGE EXTRACTION AND MANAGEMENT
# ═══════════════════════════════════════════════════════════════════════════════

def ensure_holy_communion_image():
    """Ensure the Holy Communion image exists in the images folder."""
    target_image = resolve_image_path("holy_communion.jpg")
    if os.path.exists(target_image):
        return True
    
    # Try to extract from 1 Feb 2026.pptx
    source_pptx = os.path.join(
        PARENT_DIR,
        "onedrive_git_local",
        "Holy Communion Services - Slides",
        "Malayalam HCS",
        "2026- Mal",
        "1 Feb 2026.pptx"
    )
    
    if not os.path.exists(source_pptx):
        print(f"  ⚠ Warning: Cannot find source for Holy Communion image")
        return False
    
    try:
        prs = Presentation(source_pptx)
        # Slide 23 (index 22) has the Holy Communion image
        slide = prs.slides[22]
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                os.makedirs(os.path.dirname(target_image), exist_ok=True)
                with open(target_image, "wb") as f:
                    f.write(img.blob)
                print(f"  ✓ Extracted Holy Communion image to {target_image}")
                return True
    except Exception as e:
        print(f"  ⚠ Could not extract Holy Communion image: {e}")
    
    return False


def extract_qr_code_from_slide(slide):
    """Extract QR code image from a slide if present."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # QR codes are typically positioned on the right side
            if shape.left.inches > 5:  # Right half of slide
                return shape.image.blob, shape.image.ext
    return None, None


# ═══════════════════════════════════════════════════════════════════════════════
# PDF EXTRACTION FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════

def extract_song_text_from_pdf_by_page(pdf_path, target_hymn_num):
    """
    Extract a specific song from the PDF by searching page by page.
    Returns the raw text of the song.
    """
    import pypdfium2 as pdfium

    pdf = pdfium.PdfDocument(pdf_path)
    target = str(target_hymn_num)

    # Find pages containing this hymn number
    found_pages = []
    for i in range(8, len(pdf)):  # Skip index pages
        page = pdf[i]
        textpage = page.get_textpage()
        text = textpage.get_text_bounded()

        if re.search(r"(?:^|\s)" + re.escape(target) + r"(?:\s*\(|$|\s)", text, re.MULTILINE):
            found_pages.append((i, text))

    if not found_pages:
        return None

    # Extract text from found pages
    result_lines = []
    collecting = False

    for page_idx, page_text in found_pages:
        lines = page_text.split("\n")
        for line in lines:
            stripped = line.strip()
            if not stripped:
                if collecting:
                    result_lines.append("")
                continue

            hymn_match = re.search(r"(\d{1,3})\s*(?:\(\d{1,3}\))?\s*$", stripped)
            if hymn_match:
                num = hymn_match.group(1)
                if num == target and not collecting:
                    collecting = True
                    result_lines.append(stripped)
                    continue
                elif num != target and collecting and int(num) > 0:
                    collecting = False
                    break

            if collecting:
                result_lines.append(stripped)

        if not collecting and result_lines:
            break

    return "\n".join(result_lines) if result_lines else None


# ═══════════════════════════════════════════════════════════════════════════════
# PPT SEARCH FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════

# Language-specific search directories
# These will be the base directories - the script will search recursively
def get_search_dirs(language="Malayalam"):
    """
    Get search directories based on current working directory and language.
    
    Search logic:
    1. If user provides a path (cwd != BASE_DIR), search for language-specific folder there
    2. Always include onedrive_git_local as fallback
    
    Args:
        language: "Malayalam" or "English" - determines which HCS folder to look for
    """
    # Use current working directory as base (set by GUI or default to script location)
    cwd = os.getcwd()
    
    # Determine language folder name
    if language.lower() == "english":
        lang_folder = "English HCS"
    else:
        lang_folder = "Malayalam HCS"
    
    search_dirs = []
    
    # Check if user provided a custom path (cwd is different from script location)
    user_provided_path = cwd != BASE_DIR and cwd != PARENT_DIR
    
    if user_provided_path:
        # User provided a path - search for language-specific folder
        
        # Check if cwd already IS the language folder
        if os.path.isdir(cwd) and os.path.basename(cwd).lower() == lang_folder.lower():
            search_dirs.append(cwd)
        else:
            # Check if cwd contains the language folder as immediate child
            immediate_child = os.path.join(cwd, lang_folder)
            if os.path.isdir(immediate_child):
                search_dirs.append(immediate_child)
            else:
                # Search recursively for the language folder (e.g., user gave "Holy Communion Services - Slides")
                for dirpath, dirnames, _ in os.walk(cwd):
                    for d in dirnames:
                        if d.lower() == lang_folder.lower():
                            search_dirs.append(os.path.join(dirpath, d))
                    # Stop after finding first match (don't search too deep)
                    if search_dirs:
                        break
                
                # If language folder not found, use the user path as-is (search everything)
                if not search_dirs and os.path.isdir(cwd):
                    search_dirs.append(cwd)
    
    # Always add onedrive_git_local as fallback (from exe bundle or downloaded)
    onedrive_path = os.path.join(ONEDRIVE_GIT_LOCAL, "Holy Communion Services - Slides", lang_folder)
    if os.path.isdir(onedrive_path) and onedrive_path not in search_dirs:
        search_dirs.append(onedrive_path)
    elif os.path.isdir(ONEDRIVE_GIT_LOCAL):
        # Search for language folder under onedrive_git_local
        for dirpath, dirnames, _ in os.walk(ONEDRIVE_GIT_LOCAL):
            for d in dirnames:
                if d.lower() == lang_folder.lower():
                    candidate = os.path.join(dirpath, d)
                    if candidate not in search_dirs:
                        search_dirs.append(candidate)
            if any(lang_folder.lower() in p.lower() for p in search_dirs):
                break
    
    # Final fallback to BASE_DIR if nothing found
    if not search_dirs:
        search_dirs = [BASE_DIR]
    
    return search_dirs


# Initialize with default Malayalam directories (will be updated when script runs)
MALAYALAM_SEARCH_DIRS = get_search_dirs()

def find_all_pptx_files(search_dirs):
    """Recursively find all .pptx files in Malayalam directories."""
    # Use Malayalam directories
    dirs_to_search = get_search_dirs()
    
    pptx_files = []
    for d in dirs_to_search:
        if os.path.isdir(d):
            for root, dirs, files in os.walk(d):
                for f in files:
                    if f.endswith(".pptx") and not f.startswith("~$"):
                        pptx_files.append(os.path.join(root, f))
    return pptx_files


def find_song_slide_indices_in_pptx(pptx_path, target_hymn_num="", song_title_hint=""):
    """
    Find the slide indices for a specific hymn in a PPTX file.
    
    Search logic:
    - If only hymn number provided: search by "Hymn No XXX"
    - If only song title provided: search by song title words
    - If both provided: search by either (hymn number OR song title)
    
    Args:
        pptx_path: Path to the PPTX file
        target_hymn_num: The hymn number to search for (e.g., "171") - optional
        song_title_hint: Song title to search for - optional
    
    Returns (title_slide_idx, [content_slide_indices], extracted_title) or (None, [], "").
    """
    target = str(target_hymn_num) if target_hymn_num else ""
    try:
        prs = Presentation(pptx_path)
    except Exception:
        return None, [], ""

    # Need at least one search criteria
    if not target and not song_title_hint:
        return None, [], ""

    title_slide_idx = None
    content_indices = []
    extracted_title = ""
    collecting = False
    
    # Prepare title search words
    title_words = []
    if song_title_hint and len(song_title_hint) > 2:
        title_words = [w for w in song_title_hint.split()[:3] if len(w) > 2]  # First 3 words, min 3 chars

    def slide_has_lyrics(text):
        has_malayalam = bool(re.search(r"[\u0D00-\u0D7F]", text))
        ascii_letters = sum(1 for c in text if c.isalpha() and c.isascii())
        return has_malayalam or ascii_letters >= 30
    
    def is_image_only_slide(slide):
        """Check if slide is primarily just an image without meaningful text content."""
        # Count text content
        total_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                total_text += shape.text_frame.text.strip()
        
        # Remove common labels like title bar text
        meaningful_text = total_text
        for label in ["Holy Communion Hymn", "Opening Hymn", "Thanksgiving", "Confession", 
                      "Closing Hymn", "Offertory", "Communion"]:
            meaningful_text = meaningful_text.replace(label, "")
        
        # Remove slide numbers like "30", "31"
        meaningful_text = re.sub(r'\b\d{1,3}\b', '', meaningful_text)
        meaningful_text = re.sub(r'\d+\s*:\s*\d+\s+of\s+\d+', '', meaningful_text)  # Footer
        
        # Count actual text content (letters)
        letter_count = sum(1 for c in meaningful_text if c.isalpha())
        
        # If less than 30 letters of actual content, it's likely just an image slide
        return letter_count < 30

    def build_hymn_pattern(num):
        # Pattern: "Hymn No 171" or "Song No 171"
        return r"(?:Hymn\s*(?:No\.?\s*)?[-–]?\s*|Song\s*No\.?\s*)" + re.escape(num) + r"(?:\s|\b)"
    
    def find_hymn_in_text(target_num, text):
        """Find hymn number in text using regular service PPT patterns."""
        if not target_num:
            return False
        
        hymn_pattern = build_hymn_pattern(target_num)
        return bool(re.search(hymn_pattern, text, re.IGNORECASE))

    for i, slide in enumerate(prs.slides):
        all_text = ""

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                all_text += " " + text

        # Skip summary/index slides that list multiple hymns
        hymn_count = len(re.findall(r'\b\d{2,3}\b', all_text))
        if hymn_count > 4:
            continue
        
        # Check for hymn number match (if hymn number provided)
        hymn_match = False
        if target:
            hymn_match = find_hymn_in_text(target, all_text)
        
        # Check for song title match (if title provided)
        title_match = False
        if title_words:
            # Match if at least 2 title words are found (or 1 if only 1 word provided)
            matches = sum(1 for word in title_words if word.lower() in all_text.lower())
            title_match = matches >= min(2, len(title_words))

        # Determine if this slide is a match
        is_match = False
        if target and song_title_hint:
            # Both provided: either one matches
            is_match = hymn_match or title_match
        elif target:
            # Only hymn number: must match hymn
            is_match = hymn_match
        else:
            # Only title: must match title
            is_match = title_match

        if is_match and not collecting:
            # Found start of our song - this is the title slide
            # Store the hymn number found on this slide (if any) for tracking
            hymn_found_on_title = re.search(r"Hymn\s*(?:No\.?\s*)?[-–]?\s*(\d+)", all_text, re.IGNORECASE)
            found_hymn_num = hymn_found_on_title.group(1) if hymn_found_on_title else ""
            # Also store what section we're in (from title slide)
            found_section = ""
            for sec in ["Opening", "Confession", "Offertory", "Thanksgiving", "Communion", "Closing", "Dedication", "B/A"]:
                if sec.lower() in all_text.lower():
                    found_section = sec
                    break
            
            title_slide_idx = i
            collecting = True
            
            # Extract title from the title slide (find largest Malayalam/Manglish text)
            max_text = ""
            for shape in prs.slides[i].shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    # Skip footer patterns and slide numbers
                    if re.search(r'^\d+\s*:\s*\d+\s+of\s+\d+', text):  # Footer like "30 : 31 of 106"
                        continue
                    if re.match(r'^\d{1,3}$', text):  # Just a slide number
                        continue
                    if 'Hymn' in text and len(text) < 30:  # Skip "Hymn No 313" headers
                        continue
                    if any(sec in text for sec in ["Opening", "Confession", "Offertory", "Thanksgiving", "Communion", "Closing"]) and len(text) < 40:
                        continue
                    # Check if text has Malayalam or Manglish content
                    has_malayalam = bool(re.search(r"[\u0D00-\u0D7F]", text))
                    has_manglish = len([c for c in text if c.isalpha() and c.isascii()]) > 10
                    if (has_malayalam or has_manglish) and len(text) > len(max_text):
                        max_text = text
            if max_text:
                # Take only the first line as the title, limited to first 3 words
                first_line = max_text.split('\n')[0].strip()
                words = first_line.split()
                extracted_title = ' '.join(words[:3]) if len(words) > 3 else first_line
            
            # Only add title slide to content if it has lyrics and isn't image-only
            if slide_has_lyrics(all_text) and not is_image_only_slide(slide):
                content_indices.append(i)
            continue

        if collecting:
            # Check if this is a section title slide (no lyrics, just section name)
            # These indicate end of current song
            section_title_patterns = [
                r"^[\s]*Thanksgiving\s*Prayers[\s]*$",
                r"^[\s]*Offertory[\s]*$",
                r"^[\s]*Confession[\s]*$",
                r"^[\s]*Dedication[\s]*$",
                r"^[\s]*Opening\s*Hymn[\s]*$",
                r"^[\s]*Closing\s*Hymn[\s]*$",
                r"^[\s]*Communion[\s]*$",
            ]
            # Check if slide text is ONLY a section title (very short text)
            clean_text = all_text.strip()
            if len(clean_text) < 50:  # Short slide - likely a section header
                for pattern in section_title_patterns:
                    if re.search(pattern, clean_text, re.IGNORECASE):
                        # This is a section header slide - stop collecting
                        break
                else:
                    # No section pattern matched, continue checking
                    pass
                if re.search(r"Thanksgiving|Offertory|Confession|Dedication", clean_text, re.IGNORECASE):
                    if not re.search(r"Hymn\s*(?:No\.?\s*)?\d+", clean_text, re.IGNORECASE):
                        # Section header without hymn number - stop
                        break
            
            # Check if this is a NEW section title slide with the same or different hymn
            # (e.g., "Holy Communion - Hymn No 171" when we started at "Confession - Hymn No 171")
            for sec in ["Opening", "Confession", "Offertory", "Thanksgiving", "Communion", "Closing", "Dedication"]:
                if sec.lower() in all_text.lower() and "Hymn" in all_text:
                    # This slide has a section label + Hymn - it's a title slide for a new section
                    if found_section and sec.lower() != found_section.lower():
                        # Different section - stop collecting (even if same hymn number)
                        break
            else:
                # No section break found, continue checking
                pass
            
            # If we broke out of the section check loop, stop collecting
            if sec.lower() in all_text.lower() and "Hymn" in all_text and found_section and sec.lower() != found_section.lower():
                break
            
            # Check if this slide still belongs to our song
            has_our_hymn = False
            if target:
                hymn_pattern = build_hymn_pattern(target)
                has_our_hymn = bool(re.search(hymn_pattern, all_text, re.IGNORECASE))
            elif found_hymn_num:
                # If we found a hymn number on the title slide, use that
                hymn_pattern = build_hymn_pattern(found_hymn_num)
                has_our_hymn = bool(re.search(hymn_pattern, all_text, re.IGNORECASE))
            
            has_our_title = False
            if title_words:
                matches = sum(1 for word in title_words if word.lower() in all_text.lower())
                has_our_title = matches >= min(2, len(title_words))
            
            # Check if this is a different hymn's slide
            check_hymn = target if target else found_hymn_num
            if check_hymn:
                other_hymn = re.search(r"Hymn\s*(?:No\.?\s*)?[-–]?\s*(\d+)", all_text, re.IGNORECASE)
                if other_hymn:
                    found_num = other_hymn.group(1)
                    if found_num != check_hymn:
                        # Different hymn number found - stop collecting
                        break
            else:
                # No hymn number to check - look for any new "Hymn No" pattern on a slide
                # that doesn't have our title words (indicates start of new song)
                if re.search(r"Hymn\s*(?:No\.?\s*)?\d+", all_text, re.IGNORECASE):
                    if not has_our_title:
                        break
            
            # Check for section keywords that indicate a new section (with colon = label)
            section_keywords = ("Message", "Confession:", "Offertory:", "Dedication:", 
                               "Opening:", "Closing:", "Communion:", "B/A:")
            if any(kw in all_text for kw in section_keywords):
                # If this slide has section keywords but not our reference, stop
                if not has_our_hymn and not has_our_title:
                    break
            
            # CRITICAL: Check if this slide starts a completely NEW song
            # Look for indicators that this is a new song:
            # 1. Title slide pattern with section + "Hymn No XXX" (where XXX != our hymn)
            # 2. A slide with a completely new song title in large text (not part of current hymn)
            
            # Check if this looks like a title slide for a DIFFERENT song
            # Title slides typically have: Section name + "Hymn No XXX" or just large song title
            is_likely_new_title_slide = False
            
            # Check for "Section - Hymn No XXX" pattern where XXX is different
            for sec in ["Opening", "Confession", "Offertory", "Thanksgiving", "Communion", "Closing", "Dedication"]:
                if sec in all_text and "Hymn" in all_text:
                    # This has section + Hymn pattern
                    other_hymn_match = re.search(r"Hymn\s*(?:No\.?\s*)?(\d+)", all_text, re.IGNORECASE)
                    if other_hymn_match:
                        found_num = other_hymn_match.group(1)
                        check_hymn = target if target else found_hymn_num
                        if check_hymn and found_num != check_hymn:
                            # Different hymn number = new song
                            is_likely_new_title_slide = True
                            break
            
            # Check if this slide has minimal text (< 100 chars) suggesting it's a title/section slide
            if len(all_text.strip()) < 100:
                # Short text - might be a new title slide
                # If it has a song title pattern (capitalized words) but no lyrics
                if any(word[0].isupper() for word in all_text.split() if len(word) > 3):
                    # Has capitalized words - could be title
                    if not has_our_hymn and not has_our_title:
                        # And doesn't have our hymn reference
                        is_likely_new_title_slide = True
            
            if is_likely_new_title_slide:
                break
            
            # Skip image-only slides (e.g., Holy Communion intro images)
            if is_image_only_slide(slide):
                continue
            
            # This slide belongs to our song - add as content slide
            content_indices.append(i)

    # If no title was extracted from title slide, try to extract from first content slide
    if not extracted_title and content_indices:
        first_content_idx = content_indices[0]
        max_text = ""
        for shape in prs.slides[first_content_idx].shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                # Skip footer patterns and slide numbers
                if re.search(r'^\d+\s*:\s*\d+\s+of\s+\d+', text):  # Footer like "30 : 31 of 106"
                    continue
                if re.match(r'^\d{1,3}$', text):  # Just a slide number
                    continue
                if 'Hymn' in text and len(text) < 30:  # Skip "Hymn No 313" headers
                    continue
                if any(sec in text for sec in ["Opening", "Confession", "Offertory", "Thanksgiving", "Communion", "Closing"]) and len(text) < 50:
                    continue
                # Check if text has Malayalam or Manglish content
                has_malayalam = bool(re.search(r"[\u0D00-\u0D7F]", text))
                has_manglish = len([c for c in text if c.isalpha() and c.isascii()]) > 10
                if (has_malayalam or has_manglish) and len(text) > len(max_text):
                    max_text = text
        if max_text:
            # Take only the first line as the title, limited to first 3 words
            first_line = max_text.split('\n')[0].strip()
            words = first_line.split()
            extracted_title = ' '.join(words[:3]) if len(words) > 3 else first_line

    return title_slide_idx, content_indices, extracted_title


# Global tracking of used slide ranges to prevent duplicates
USED_SLIDE_RANGES = {}

def find_best_song_source(hymn_num, song_name):
    """
    Search all PPT files and find the best source for a Malayalam song.
    Returns the source with the most content slides.
    Tracks used slide ranges to prevent adding the same hymn content twice.
    
    Search priority:
    1. First search all non-KK PPT files (actual service presentations)
    2. Only if not found, search KK.pptx files (hymn book as last resort)
    
    Args:
        hymn_num: The hymn number to search for
        song_name: Optional song title hint
    
    Returns: (pptx_path, title_idx, content_indices, extracted_title) or (None, None, [], "")
    """
    pptx_files = find_all_pptx_files(SEARCH_DIRS)
    
    # Separate KK files from regular service PPT files
    kk_files = [pf for pf in pptx_files if "KK" in os.path.basename(pf).upper() or "Kristeeya" in os.path.basename(pf)]
    regular_files = [pf for pf in pptx_files if pf not in kk_files]
    
    best_source = None
    best_count = 0
    best_title_idx = None
    best_content = []
    best_extracted_title = ""
    
    # First search regular service PPT files
    for pf in regular_files:
        t_idx, c_indices, extracted_title = find_song_slide_indices_in_pptx(pf, hymn_num, song_name)
        if c_indices and len(c_indices) > best_count:
            # Check if these slides were already used BY A DIFFERENT HYMN NUMBER
            # (Same hymn can be reused multiple times in one service)
            slide_key = f"{pf}:{min(c_indices)}-{max(c_indices)}"
            if slide_key not in USED_SLIDE_RANGES or USED_SLIDE_RANGES[slide_key] == hymn_num:
                best_source = pf
                best_count = len(c_indices)
                best_title_idx = t_idx
                best_content = c_indices
                best_extracted_title = extracted_title
    
    # If not found in regular files, search KK files as last resort
    if not best_source:
        for pf in kk_files:
            # Use specialized KK search function for KK.pptx files
            t_idx, c_indices, extracted_title = find_hymn_in_kk_pptx(pf, hymn_num)
            
            if c_indices and len(c_indices) > best_count:
                slide_key = f"{pf}:{min(c_indices)}-{max(c_indices)}"
                if slide_key not in USED_SLIDE_RANGES or USED_SLIDE_RANGES[slide_key] == hymn_num:
                    best_source = pf
                    best_count = len(c_indices)
                    best_title_idx = t_idx
                    best_content = c_indices
                    # Use title from kk_hymn_mapping.json instead of extracted title
                    best_extracted_title = KK_HYMN_MAPPING.get(str(hymn_num), extracted_title)
    
    # Mark these slides as used if we found something
    if best_source and best_content:
        slide_key = f"{best_source}:{min(best_content)}-{max(best_content)}"
        USED_SLIDE_RANGES[slide_key] = hymn_num
    
    return best_source, best_title_idx, best_content, best_extracted_title


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE CREATION FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════

def create_title_bar(slide, prs, title_text):
    """Create the pink title bar at the top of a content slide."""
    # Rectangle background
    title_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Emu(0), Emu(0),
        prs.slide_width, TITLE_BAR_HEIGHT
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = TITLE_BAR_COLOR
    title_bar.line.fill.background()
    
    # Text on bar
    title_textbox = slide.shapes.add_textbox(
        Emu(0), Emu(100000),
        prs.slide_width, TITLE_BAR_HEIGHT
    )
    title_tf = title_textbox.text_frame
    title_tf.word_wrap = False
    title_p = title_tf.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_run = title_p.add_run()
    title_run.text = title_text
    title_run.font.name = "Segoe UI"
    title_run.font.size = TITLE_BAR_FONT_SIZE
    title_run.font.bold = False
    title_run.font.color.rgb = RGBColor(0, 0, 0)


def add_title_slide(prs, layout, song_label, hymn_num, song_name=""):
    """Add a title/intro slide for a song section."""
    slide = prs.slides.add_slide(layout)

    tb = slide.shapes.add_textbox(
        Emu(12031), Emu(1679968), Emu(9156031), Emu(1783563)
    )
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = song_label
    run.font.name = TITLE_FONT
    run.font.size = TITLE_SIZE
    run.font.bold = True

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    if hymn_num:
        # Has hymn number
        if song_name:
            run2.text = f"Hymn No. {hymn_num} - {song_name}"
        else:
            run2.text = f"Hymn No. {hymn_num}"
    else:
        # No hymn number - show just title
        if song_name:
            run2.text = f"Hymn - {song_name}"
        else:
            run2.text = "Hymn"
    run2.font.name = TITLE_FONT
    run2.font.size = TITLE_SIZE
    run2.font.bold = True

    return slide


def add_message_slide(prs, layout):
    """Add a Message title slide (no hymn content)."""
    slide = prs.slides.add_slide(layout)

    tb = slide.shapes.add_textbox(
        Emu(12031), Emu(1679968), Emu(9156031), Emu(1783563)
    )
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Message"
    run.font.name = TITLE_FONT
    run.font.size = TITLE_SIZE
    run.font.bold = True

    return slide


def add_holy_communion_intro_slide(prs, blank_layout, hymn_num, song_name=""):
    """Add a Holy Communion intro slide - blank slide with title bar and image."""
    slide = prs.slides.add_slide(blank_layout)
    
    # Remove title placeholders
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph = shape.placeholder_format
            if ph.type in [1, 2, 3]:
                shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    
    # Create the pink title bar with hymn info
    if song_name:
        title_text = f"Holy Communion - {song_name}"
    else:
        title_text = f"Holy Communion - Hymn No: {hymn_num}"
    create_title_bar(slide, prs, title_text)

    # Add the Holy Communion image below the title bar
    hc_image_path = resolve_image_path("holy_communion.jpg")
    if os.path.exists(hc_image_path):
        slide.shapes.add_picture(
            hc_image_path,
            HC_IMAGE_LEFT, HC_IMAGE_TOP,
            HC_IMAGE_WIDTH, HC_IMAGE_HEIGHT
        )

    return slide


def is_malayalam_text(text):
    """Check if text contains Malayalam Unicode characters or corrupted Malayalam."""
    # Malayalam Unicode range: U+0D00 to U+0D7F
    malayalam_count = sum(1 for char in text if '\u0D00' <= char <= '\u0D7F')
    if malayalam_count > len(text) * 0.2:
        return True
    
    # Check for corrupted/encoded Malayalam by looking at special characters
    # Corrupted Malayalam typically has backslashes, special chars mixed with ASCII
    special_chars = sum(1 for char in text if ord(char) >= 128 or char in '\\')
    total_chars = len(text.strip())
    
    # If more than 10% special chars/non-ASCII, likely Malayalam
    if total_chars > 50 and special_chars > total_chars * 0.1:
        return True
    
    return False


def clone_slides_from_source(source_pptx_path, slide_indices, target_prs,
                              song_label, hymn_num, slide_num_start,
                              blank_layout, song_name="", is_offertory=False):
    """
    Clone slides from source PPT, preserving text and images.
    For Offertory slides with overlapping text, splits into Manglish and Malayalam slides.
    Returns the number of slides added.
    """
    source_prs = Presentation(source_pptx_path)
    source_slides = list(source_prs.slides)
    added = 0
    
    # Check if source is a KK hymn file
    is_kk_file = "KK" in os.path.basename(source_pptx_path).upper() or "Kristeeya" in os.path.basename(source_pptx_path)
    
    # Build title text for content slides (just Hymn No and Title, no section label)
    if song_name:
        slide_title_text = f"Hymn No {hymn_num} - {song_name}"
    else:
        slide_title_text = f"Hymn No {hymn_num}"
    
    # QR code settings for Offertory slides
    QR_LEFT = Inches(6.97)       # Position on right side
    QR_TOP = Inches(1.08)        # Vertical position
    QR_WIDTH = Inches(3.04)      # QR code width
    QR_HEIGHT = Inches(3.13)     # QR code height
    TEXT_MAX_RIGHT = Inches(6.5) # Text should not extend beyond this
    TEXT_LEFT_MARGIN = Inches(0.16)
    TEXT_WIDTH = Inches(6.0)     # Width for text when QR code is present
    
    # Check if QR code file exists
    qr_code_path = None
    if is_offertory:
        potential_path = resolve_image_path("qr_code.png")
        if os.path.exists(potential_path):
            qr_code_path = potential_path

    for idx in slide_indices:
        if idx >= len(source_slides):
            continue

        src_slide = source_slides[idx]
        
        # Clone the exact slide preserving all formatting
        new_slide = clone_slide_exact(src_slide, target_prs, blank_layout)

        # If this is from a KK hymn file, remove the KK header and decorative shapes
        if is_kk_file:
            remove_kk_hymn_header(new_slide)
            remove_kk_decorative_shapes(new_slide)

        # Always remove QR code and UEN from source slides (in case source was from offertory)
        remove_qr_and_uen(new_slide)
        remove_footer_text(new_slide)

        # Update title bar text based on which section this hymn is being added to
        update_title_bar_text(new_slide, song_label, hymn_num, song_name)

        # Only add QR code if this is an Offertory slide
        if is_offertory and qr_code_path:
            add_qr_code_to_slide(new_slide, qr_code_path, QR_LEFT, QR_TOP, QR_WIDTH, QR_HEIGHT)

        # Update title bar color to match template
        update_title_bar_color(new_slide)
        added += 1

    return added


def split_text_for_slides(text):
    """Split text into two parts, preferring stanza breaks."""
    lines = text.splitlines()
    if len(lines) <= 1:
        return text, text
    
    midpoint = len(lines) // 2
    blank_indices = [i for i, line in enumerate(lines) if not line.strip()]
    if blank_indices:
        split_at = min(blank_indices, key=lambda i: abs(i - midpoint)) + 1
        if split_at <= 0 or split_at >= len(lines):
            split_at = midpoint
    else:
        split_at = midpoint
    
    first = "\n".join(lines[:split_at]).strip()
    second = "\n".join(lines[split_at:]).strip()
    return first, second


def apply_text_frame_settings(source_tf, target_tf):
    """Copy text frame settings that affect layout/spacing."""
    try:
        target_tf.word_wrap = source_tf.word_wrap
    except Exception:
        pass
    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom", "vertical_anchor", "auto_size"):
        try:
            setattr(target_tf, attr, getattr(source_tf, attr))
        except Exception:
            pass


def apply_font_settings(source_font, target_font):
    """Copy font settings when present."""
    try:
        if source_font.name:
            target_font.name = source_font.name
    except Exception:
        pass
    try:
        if source_font.size:
            target_font.size = source_font.size
    except Exception:
        pass
    try:
        if source_font.bold is not None:
            target_font.bold = source_font.bold
    except Exception:
        pass
    try:
        if source_font.italic is not None:
            target_font.italic = source_font.italic
    except Exception:
        pass
    try:
        if source_font.color and source_font.color.type is not None:
            target_font.color.rgb = source_font.color.rgb
    except Exception:
        pass


def apply_paragraph_settings(source_para, target_para):
    """Copy paragraph spacing settings to prevent overlap."""
    try:
        target_para.alignment = source_para.alignment
    except Exception:
        pass
    for attr in ("space_before", "space_after", "line_spacing", "level"):
        try:
            setattr(target_para, attr, getattr(source_para, attr))
        except Exception:
            pass
    try:
        apply_font_settings(source_para.font, target_para.font)
    except Exception:
        pass


def copy_text_shape_to_slide(source_shape, target_slide, left_margin, max_width, text_override=None):
    """Copy a text shape to target slide with adjusted position and width."""
    # Adjust position to fit alongside QR code
    shape_left = left_margin
    shape_width = max_width
    
    new_shape = target_slide.shapes.add_textbox(
        shape_left, source_shape.top, shape_width, source_shape.height
    )
    new_tf = new_shape.text_frame
    apply_text_frame_settings(source_shape.text_frame, new_tf)
    new_tf.word_wrap = True

    if text_override is not None:
        p = new_tf.paragraphs[0]
        apply_paragraph_settings(source_shape.text_frame.paragraphs[0], p)
        run = p.add_run()
        run.text = text_override
        
        if source_shape.text_frame.paragraphs and source_shape.text_frame.paragraphs[0].runs:
            src_run = source_shape.text_frame.paragraphs[0].runs[0]
            apply_font_settings(src_run.font, run.font)
        return

    for pi, para in enumerate(source_shape.text_frame.paragraphs):
        if pi == 0:
            new_para = new_tf.paragraphs[0]
        else:
            new_para = new_tf.add_paragraph()
        
        apply_paragraph_settings(para, new_para)
        
        for ri, run in enumerate(para.runs):
            if ri == 0 and pi == 0 and len(new_para.runs) > 0:
                new_run = new_para.runs[0]
            else:
                new_run = new_para.add_run()
            
            new_run.text = run.text
            
            if run.font.name:
                new_run.font.name = run.font.name
            if run.font.size:
                new_run.font.size = run.font.size
            if run.font.bold is not None:
                new_run.font.bold = run.font.bold
            if run.font.italic is not None:
                new_run.font.italic = run.font.italic
            
            try:
                if run.font.color and run.font.color.type is not None:
                    new_run.font.color.rgb = run.font.color.rgb
            except:
                pass
            
            # Apply paragraph font defaults when run font is missing
            if new_run.font.name is None and para.font.name:
                new_run.font.name = para.font.name
            if new_run.font.size is None and para.font.size:
                new_run.font.size = para.font.size
            if new_run.font.bold is None and para.font.bold is not None:
                new_run.font.bold = para.font.bold
            if new_run.font.italic is None and para.font.italic is not None:
                new_run.font.italic = para.font.italic
            try:
                if new_run.font.color is None and para.font.color and para.font.color.type is not None:
                    new_run.font.color.rgb = para.font.color.rgb
            except Exception:
                pass


def add_qr_code_to_slide(slide, qr_code_path, left, top, width, height):
    """Add QR code and UEN text to a slide."""
    try:
        slide.shapes.add_picture(
            qr_code_path,
            left, top,
            width, height
        )
        
        # Add UEN text below QR code
        uen_left = left
        uen_top = top + height + Inches(0.1)
        uen_width = width
        uen_height = Inches(0.3)
        
        uen_box = slide.shapes.add_textbox(uen_left, uen_top, uen_width, uen_height)
        uen_tf = uen_box.text_frame
        uen_tf.word_wrap = False
        uen_p = uen_tf.paragraphs[0]
        uen_p.alignment = PP_ALIGN.CENTER
        
        uen_run = uen_p.add_run()
        uen_run.text = "UEN - S86CC0315K"
        uen_run.font.name = "Arial"
        uen_run.font.size = Pt(10)
        uen_run.font.bold = False
    except Exception as e:
        print(f"    Warning: Could not add QR code: {e}")


def remove_qr_and_uen(slide):
    """Remove existing QR images and UEN text from a slide."""
    for shape in list(slide.shapes):
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.left.inches > 5:
                sp = shape._element
                sp.getparent().remove(sp)
                continue
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "UEN" in text:
                    sp = shape._element
                    sp.getparent().remove(sp)
        except Exception:
            continue


def remove_footer_text(slide):
    """Remove footer text like "B/A: 1 of 2" or "Communion 2: 1 of 7"."""
    footer_pattern = re.compile(r".+?:\s*\d+\s+of\s+\d+", re.IGNORECASE)
    for shape in list(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        if footer_pattern.fullmatch(text) or footer_pattern.search(text):
            sp = shape._element
            sp.getparent().remove(sp)


def remove_kk_hymn_header(slide):
    """Remove the KK hymn header/title from the slide (top text with 'Hymn No: X' pattern)."""
    for shape in list(slide.shapes):
        if not shape.has_text_frame:
            continue
        # Check if shape is at the top of the slide
        if shape.top > Emu(1000000):  # Not in header area (< 1 inch from top)
            continue
        
        text = shape.text_frame.text.strip()
        if not text:
            continue
        
        # Check if this looks like a KK hymn header (contains "Hymn No:" or just the hymn number)
        if re.search(r'Hymn\s*No\s*:?\s*\d+', text, re.IGNORECASE):
            sp = shape._element
            sp.getparent().remove(sp)
            continue
        
        # Also remove standalone hymn numbers at the top
        if re.match(r'^\d{1,3}$', text):
            sp = shape._element
            sp.getparent().remove(sp)


def remove_kk_decorative_shapes(slide):
    """Remove decorative shapes from KK slides (checkmarks, arrows, etc)."""
    slide_width = slide.width if hasattr(slide, 'width') else Emu(9144000)
    slide_height = slide.height if hasattr(slide, 'height') else Emu(6858000)
    
    for shape in list(slide.shapes):
        # Skip text frames and pictures (actual content)
        if shape.has_text_frame:
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue
        
        # Remove shapes that are clearly decorative (not the main content area)
        # Check if shape is in corners or edges (not center content area)
        center_x = slide_width // 2
        center_y = slide_height // 2
        shape_center_x = shape.left + (shape.width // 2)
        shape_center_y = shape.top + (shape.height // 2)
        
        # If shape center is far from slide center, it's likely decorative
        # Content is usually in center 60% of slide
        horizontal_margin = slide_width * 0.2  # 20% margins on each side
        vertical_margin = slide_height * 0.2
        
        is_outside_content = (
            shape_center_x < horizontal_margin or 
            shape_center_x > (slide_width - horizontal_margin) or
            shape_center_y < vertical_margin or
            shape_center_y > (slide_height - vertical_margin)
        )
        
        # Also remove any shape that's not very large (decorative elements are smaller than content)
        is_small = shape.width < (slide_width * 0.4) and shape.height < (slide_height * 0.4)
        
        if is_outside_content or is_small:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
            except:
                pass


def clone_slide_exact(src_slide, target_prs, blank_layout):
    """Clone a slide preserving original shapes and formatting."""
    new_slide = target_prs.slides.add_slide(blank_layout)

    # Remove all default shapes/placeholders
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Copy shapes
    for shape in src_slide.shapes:
        el = deepcopy(shape._element)
        new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')

    # Copy relationships (images, etc.) - need to import media parts to target presentation
    for rel in src_slide.part.rels.values():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        
        if rel.is_external:
            # External relationships (e.g., hyperlinks) use target_ref string
            new_slide.part.rels._add_relationship(
                rel.reltype,
                rel.target_ref,
                rel.rId,
            )
        elif "image" in rel.reltype:
            # Image relationships - import the image part to target presentation
            try:
                source_part = rel._target
                # Import image/media part to target presentation
                # Wrap blob bytes in BytesIO to provide file-like interface
                image_stream = BytesIO(source_part.blob)
                image_part = target_prs.part.package.get_or_add_image_part(image_stream)
                # Add relationship from slide to the imported image
                new_slide.part.relate_to(image_part, rel.reltype, rel.rId)
            except Exception as e:
                # Skip relationships we can't copy
                print(f"    ⚠ Warning: Could not copy image: {e}")
                continue

    return new_slide


def update_title_bar_text(slide, label, hymn_num, song_name=""):
    """Update the title bar text to match the section label and hymn number."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Find title bar by checking if it's at the top and contains hymn-related text
        if shape.top > Emu(500000):  # Not in title bar area
            continue
        
        text = shape.text_frame.text.strip()
        # Check if this is a title bar (contains "Hymn" or section name)
        if "Hymn" not in text and "Offertory" not in text and "Opening" not in text and "Confession" not in text and "Communion" not in text and "Closing" not in text and "ThanksGiving" not in text:
            continue
        
        try:
            # Update the text to show only section label
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            # Show only the section label (e.g., "Opening", "ThanksGiving", "Confession")
            run.text = label
            run.font.name = "Calibri"
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 0, 0)
        except Exception as e:
            pass
        break  # Only update the first matching title bar


def update_title_bar_color(slide, color=TITLE_BAR_COLOR):
    """Update the existing title bar fill color to match the template."""
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        if shape.top > Emu(100000):
            continue
        if shape.height > Emu(700000):
            continue
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
        except Exception:
            pass


def adjust_text_overlaps(slide, padding=Inches(0.08)):
    """Shift text boxes down if they overlap to avoid collisions."""
    text_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        # Skip title bar and footer
        if "Hymn No" in text or "Holy Communion" in text:
            continue
        if "Communion" in text and "of" in text:
            continue
        text_shapes.append(shape)
    
    text_shapes.sort(key=lambda s: s.top)
    prev_bottom = None
    for shape in text_shapes:
        if prev_bottom is None:
            prev_bottom = shape.top + shape.height
            continue
        if shape.top < prev_bottom + padding:
            shape.top = prev_bottom + padding
        prev_bottom = shape.top + shape.height


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION PROCESSING FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════

def process_opening_song(prs, title_layout, blank_layout, hymn_num, song_name, slide_counter):
    """Process Opening Song section."""
    label = "Opening"
    
    # Find the best source (most slides) across all PPT files, prioritizing "Opening" section in KK.pptx
    best_pf, t_idx, c_indices, extracted_title = find_best_song_source(hymn_num, song_name)
    
    if best_pf and c_indices:
        display_title = extracted_title if extracted_title else song_name
        print(f"  ✓ Found in PPT: {os.path.basename(best_pf)} ({len(c_indices)} content slides)")
        add_title_slide(prs, title_layout, label, hymn_num, display_title)
        slide_counter += 1
        num_added = clone_slides_from_source(
            best_pf, c_indices, prs, label, hymn_num, slide_counter,
            blank_layout, display_title
        )
        slide_counter += num_added
        return slide_counter, display_title
    
    # Not found - add title only
    print(f"  ⚠ Song not found - adding title slide only")
    add_title_slide(prs, title_layout, label, hymn_num, song_name if song_name else "(Song not found)")
    return slide_counter + 1, song_name


def process_thanksgiving_prayers(prs, title_layout, blank_layout, hymn_num, song_name, slide_counter):
    """Process ThanksGiving Prayers section (shown as B/A on summary, ThanksGiving on slides)."""
    label = "ThanksGiving"
    
    best_pf, t_idx, c_indices, extracted_title = find_best_song_source(hymn_num, song_name, )
    
    if best_pf and c_indices:
        display_title = extracted_title if extracted_title else song_name
        print(f"  ✓ Found in PPT: {os.path.basename(best_pf)} ({len(c_indices)} content slides)")
        add_title_slide(prs, title_layout, label, hymn_num, display_title)
        slide_counter += 1
        num_added = clone_slides_from_source(
            best_pf, c_indices, prs, label, hymn_num, slide_counter,
            blank_layout, display_title
        )
        slide_counter += num_added
        return slide_counter, display_title
    
    print(f"  ⚠ Song not found - adding title slide only")
    add_title_slide(prs, title_layout, label, hymn_num, song_name if song_name else "(Song not found)")
    return slide_counter + 1, song_name


def process_offertory(prs, title_layout, blank_layout, hymn_num, song_name, slide_counter):
    """Process Offertory section (with QR code extraction)."""
    label = "Offertory"
    
    best_pf, t_idx, c_indices, extracted_title = find_best_song_source(hymn_num, song_name, )
    
    if best_pf and c_indices:
        display_title = extracted_title if extracted_title else song_name
        print(f"  ✓ Found in PPT: {os.path.basename(best_pf)} ({len(c_indices)} content slides)")
        add_title_slide(prs, title_layout, label, hymn_num, display_title)
        slide_counter += 1
        num_added = clone_slides_from_source(
            best_pf, c_indices, prs, label, hymn_num, slide_counter,
            blank_layout, display_title, is_offertory=True
        )
        slide_counter += num_added
        return slide_counter, display_title
    
    print(f"  ⚠ Song not found - adding title slide only")
    add_title_slide(prs, title_layout, label, hymn_num, song_name if song_name else "(Song not found)")
    return slide_counter + 1, song_name


def process_message(prs, title_layout, slide_counter):
    """Process Message section (title slide only)."""
    add_message_slide(prs, title_layout)
    print(f"  ✓ Added Message title slide")
    return slide_counter + 1


def process_confession(prs, title_layout, blank_layout, hymn_num, song_name, slide_counter):
    """Process Confession section."""
    label = "Confession"
    
    best_pf, t_idx, c_indices, extracted_title = find_best_song_source(hymn_num, song_name, )
    
    if best_pf and c_indices:
        display_title = extracted_title if extracted_title else song_name
        print(f"  ✓ Found in PPT: {os.path.basename(best_pf)} ({len(c_indices)} content slides)")
        add_title_slide(prs, title_layout, label, hymn_num, display_title)
        slide_counter += 1
        num_added = clone_slides_from_source(
            best_pf, c_indices, prs, label, hymn_num, slide_counter,
            blank_layout, display_title
        )
        slide_counter += num_added
        return slide_counter, display_title
    
    print(f"  ⚠ Song not found - adding title slide only")
    add_title_slide(prs, title_layout, label, hymn_num, song_name if song_name else "(Song not found)")
    return slide_counter + 1, song_name


def process_holy_communion(prs, title_layout, blank_layout, hymn_num, song_name, slide_counter):
    """Process Holy Communion section - every song gets Holy Communion intro slide with image."""
    label = "Communion"
    
    # Ensure Holy Communion image exists
    ensure_holy_communion_image()
    
    best_pf, t_idx, c_indices, extracted_title = find_best_song_source(hymn_num, song_name, )
    
    if best_pf and c_indices:
        display_title = extracted_title if extracted_title else song_name
        print(f"  ✓ Found in PPT: {os.path.basename(best_pf)} ({len(c_indices)} content slides)")
        
        # Add Holy Communion intro slide with image
        add_holy_communion_intro_slide(prs, blank_layout, hymn_num, display_title)
        slide_counter += 1
        
        num_added = clone_slides_from_source(
            best_pf, c_indices, prs, label, hymn_num, slide_counter,
            blank_layout, display_title
        )
        slide_counter += num_added
        return slide_counter, display_title
    
    print(f"  ⚠ Song not found - adding title slide only")
    add_holy_communion_intro_slide(prs, blank_layout, hymn_num, song_name if song_name else "(Song not found)")
    return slide_counter + 1, song_name


def process_closing_hymn(prs, title_layout, blank_layout, hymn_num, song_name, slide_counter):
    """Process Closing Hymn section."""
    label = "Closing"
    
    best_pf, t_idx, c_indices, extracted_title = find_best_song_source(hymn_num, song_name, )
    
    if best_pf and c_indices:
        display_title = extracted_title if extracted_title else song_name
        print(f"  ✓ Found in PPT: {os.path.basename(best_pf)} ({len(c_indices)} content slides)")
        add_title_slide(prs, title_layout, label, hymn_num, display_title)
        slide_counter += 1
        num_added = clone_slides_from_source(
            best_pf, c_indices, prs, label, hymn_num, slide_counter,
            blank_layout, display_title
        )
        slide_counter += num_added
        return slide_counter, display_title
    
    print(f"  ⚠ Song not found - adding title slide only")
    add_title_slide(prs, title_layout, label, hymn_num, song_name if song_name else "(Song not found)")
    return slide_counter + 1, song_name


def process_generic_song(prs, title_layout, blank_layout, label, hymn_num, song_name, slide_counter):
    """Process any generic song section."""
    best_pf, t_idx, c_indices, extracted_title = find_best_song_source(hymn_num, song_name, )
    
    if best_pf and c_indices:
        display_title = extracted_title if extracted_title else song_name
        print(f"  ✓ Found in PPT: {os.path.basename(best_pf)} ({len(c_indices)} content slides)")
        add_title_slide(prs, title_layout, label, hymn_num, display_title)
        slide_counter += 1
        num_added = clone_slides_from_source(
            best_pf, c_indices, prs, label, hymn_num, slide_counter,
            blank_layout, display_title
        )
        slide_counter += num_added
        return slide_counter, display_title
    
    print(f"  ⚠ Song not found - adding title slide only")
    add_title_slide(prs, title_layout, label, hymn_num, song_name if song_name else "(Song not found)")
    return slide_counter + 1, song_name


# ═══════════════════════════════════════════════════════════════════════════════
# SUMMARY SLIDE GENERATION
# ═══════════════════════════════════════════════════════════════════════════════

def create_summary_slide(prs, title_layout, song_list, service_date=None):
    """
    Create a summary slide listing all songs in the service.
    
    Note: ThanksGiving is displayed as "B/A" on the summary slide only.
    Communion songs are grouped together with the label shown once.
    Uses a textbox with formatted text for simplicity.
    """
    summary_slide = prs.slides.add_slide(title_layout)
    
    # Update the date in the footer if provided
    if service_date:
        import re
        # Find and update the footer text in the slide layout shapes
        for shape in summary_slide.slide_layout.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                # Look for the footer with date pattern "DD Month YYYY"
                date_pattern = r'\d{1,2}\s+(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}'
                if re.search(date_pattern, text) and "Holy Communion Service" in text:
                    # This is the footer - update the date
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if re.search(date_pattern, run.text):
                                run.text = re.sub(date_pattern, service_date, run.text)
                                break
    
    # Create a textbox positioned below the title
    # Slide is approximately 9144000 x 6858000 EMU
    left = Emu(800000)      # Left margin
    top = Emu(550000)       # Below "Mar Thoma Syrian Church, Singapore" title
    width = Emu(7500000)    # Wide enough for content
    height = Emu(5500000)   # Tall enough for all lines
    
    textbox = summary_slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    # Fixed column width - all numbers start at same position
    LABEL_COLUMN_WIDTH = 15  # Total chars for label column (including colon)
    
    # Build formatted content
    last_label = None
    first_para = True
    
    for song in song_list:
        label_lower = song['label'].lower()
        is_communion = label_lower in ("communion", "holy communion")
        is_thanksgiving = label_lower in ("thanksgiving", "thanksgiving prayers", "b/a")
        
        if first_para:
            p = tf.paragraphs[0]
            first_para = False
        else:
            p = tf.add_paragraph()
        
        # Set paragraph spacing
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        
        if is_communion and last_label == "communion":
            # Continuation of communion - indent with fixed width spaces
            run = p.add_run()
            indent = " " * LABEL_COLUMN_WIDTH
            run.text = f"{indent}{song['hymn_num']} {song.get('title_hint', '')}"
            run.font.name = "Arial"
            run.font.size = Pt(14)
            run.font.bold = False
        else:
            # Show label with fixed-width padding
            display_label = "B/A" if is_thanksgiving else song['label']
            label_with_colon = f"{display_label}:"
            
            # Pad label to fixed width
            spaces_needed = LABEL_COLUMN_WIDTH - len(label_with_colon)
            padded_label = label_with_colon + " " * spaces_needed
            
            # Single run with padded label + hymn info
            run = p.add_run()
            if song['hymn_num']:
                run.text = f"{padded_label}{song['hymn_num']} {song.get('title_hint', '')}"
            else:
                run.text = padded_label
            run.font.name = "Arial"
            run.font.size = Pt(14)
            
            # Make label portion bold by using separate runs
            p.clear()
            run1 = p.add_run()
            run1.text = label_with_colon
            run1.font.name = "Arial"
            run1.font.size = Pt(14)
            run1.font.bold = True
            
            run2 = p.add_run()
            padding = " " * spaces_needed
            if song['hymn_num']:
                run2.text = f"{padding}{song['hymn_num']} {song.get('title_hint', '')}"
            else:
                # If no hymn number, show just the title if available
                title = song.get('title_hint', '')
                run2.text = f"{padding}{title}" if title else ""
            run2.font.name = "Arial"
            run2.font.size = Pt(14)
            run2.font.bold = False
        
        last_label = "communion" if is_communion else label_lower
    
    return summary_slide


def update_summary_slide_from_slides(prs, song_list):
    """Extract first line of Manglish from each hymn section and update summary."""
    if len(prs.slides) < 2:
        return
    
    # Scan slides to find first Manglish line for each section
    hymn_titles = {}  # {hymn_num: first_line}
    fallback_titles = {}  # {hymn_num: first_line}
    current_hymn = None
    
    for slide_idx in range(1, len(prs.slides)):
        slide = prs.slides[slide_idx]
        
        # First pass: check if this slide has a section header (not a footer)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            
            # Skip lines that look like footers (contain " of " pattern)
            if " of " in text and ":" in text:
                # print(f"    [SKIP] Footer detected: {text[:50]}")
                continue
            
            # Look for section title pattern with hymn number
            # Patterns: "Opening Hymn No 231", "ThanksGiving Hymn- 236", "Communion Hymn- 242"
            section_match = re.search(r"(Opening|Thanksgiving|ThanksGiving|Offertory|Confession|Communion|Closing|Dedication|B/A)\s+(?:Hymn.*?\s+)?(\d+)", text)
            if section_match:
                hymn_num = section_match.group(2)
                # Accept hymn numbers with 2 or more digits (excludes single-digit footer numbers like "6")
                if len(hymn_num) >= 2:
                    current_hymn = hymn_num
                    if current_hymn not in hymn_titles:
                        hymn_titles[current_hymn] = ""
                else:
                    # Skip single-digit numbers like "6" from corrupted footers
                    pass
        
        # Second pass: if we're tracking a hymn, look for first line of Manglish
        if current_hymn and hymn_titles[current_hymn] == "":
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text or len(text) < 5:
                    continue
                
                # Skip section headers and footers
                if re.search(r"(Opening|Thanksgiving|ThanksGiving|Offertory|Confession|Communion|Closing|B/A|Dedication)", text, re.IGNORECASE):
                    continue
                if re.match(r"^\d+\s*:\s*\d+\s+of\s+\d+", text):  # Footer like "1: 1 of 6"
                    continue
                if re.match(r"^(Trinity|Message|Holy|Mar Thoma)", text):
                    continue
                if "Hymn" in text:
                    continue
                
                # Get first line (handle both \n and \x0b vertical tab)
                first_line = text.replace('\x0b', '\n').split('\n')[0].strip()
                
                if len(first_line) < 5 or len(first_line) > 120:
                    continue
                
                # Check if it's Manglish (mostly Latin characters, no Hindi/Malayalam)
                basic_latin = sum(1 for c in first_line if (c >= 'a' and c <= 'z') or (c >= 'A' and c <= 'Z'))
                total_alpha = sum(1 for c in first_line if c.isalpha())
                malayalam_chars = sum(1 for c in first_line if ord(c) >= 0x0D00 and ord(c) <= 0x0D7F)
                hindi_chars = sum(1 for c in first_line if ord(c) >= 0x0900 and ord(c) <= 0x097F)
                # Check for other non-Latin chars (corrupted Malayalam, etc)
                non_ascii = sum(1 for c in first_line if ord(c) > 127)
                # Count special/weird characters
                special_chars = sum(1 for c in first_line if c in '³²¹{}[]<>|\\©®™')
                
                if total_alpha < 10:
                    continue
                if malayalam_chars > 3 or hindi_chars > 0:
                    continue
                # If >20% of text is non-ASCII, it's likely not Manglish
                if len(first_line) > 0 and non_ascii / len(first_line) > 0.20:
                    continue
                # If line has special garbage characters, skip it (no tolerance)
                if special_chars > 0:
                    continue
                
                # Extract only first 2-3 words for summary
                words = first_line.split()
                title_words = ' '.join(words[:3]) if len(words) >= 3 else ' '.join(words[:2])
                
                hymn_titles[current_hymn] = title_words
                break  # Found it for this hymn, move to next slide

        # Fallback: store a readable line even if Manglish rules fail
        if current_hymn and hymn_titles[current_hymn] == "" and current_hymn not in fallback_titles:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text or len(text) < 5:
                    continue
                if re.search(r"(Opening|Thanksgiving|ThanksGiving|Offertory|Confession|Communion|Closing|Dedication|B/A)", text, re.IGNORECASE):
                    continue
                if re.match(r"^\d+\s*:\s*\d+\s+of\s+\d+", text):
                    continue
                if re.match(r"^(Trinity|Message|Holy|Mar Thoma)", text):
                    continue
                if "Hymn" in text:
                    continue
                first_line = text.replace('\x0b', '\n').split('\n')[0].strip()
                if len(first_line) < 5:
                    continue
                words = first_line.split()
                title_words = ' '.join(words[:3]) if len(words) >= 3 else ' '.join(words[:2])
                fallback_titles[current_hymn] = title_words
                break
    
    # Update song_list with extracted titles (only if not already set)
    for song in song_list:
        hymn_num = str(song.get('hymn_num', ''))
        # Skip if title_hint already exists (was set during processing)
        if song.get('title_hint'):
            continue
        # Try to extract from slides
        if hymn_num in hymn_titles and hymn_titles[hymn_num]:
            song['title_hint'] = hymn_titles[hymn_num]
        elif hymn_num in fallback_titles and fallback_titles[hymn_num]:
            song['title_hint'] = fallback_titles[hymn_num]
    
    # Now rebuild summary with titles
    summary_slide = prs.slides[0]
    for shape in summary_slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if "Mar Thoma" in text or "Holy Communion" in text:
            continue
        if shape.width > Inches(4):  # Content area
            tf = shape.text_frame
            tf.clear()
            
            LABEL_COLUMN_WIDTH = 15
            first_para = True
            last_label = None
            communion_counter = 0
            
            for song in song_list:
                label = song.get('label', '')
                label_lower = label.lower()
                is_communion = label_lower in ('communion', 'holy communion')
                
                if is_communion:
                    communion_counter += 1
                    if communion_counter == 1:
                        display_label = "Communion"
                        label_with_colon = f"{display_label}:"
                    else:
                        display_label = ""
                        label_with_colon = ""
                else:
                    communion_counter = 0
                    if label_lower in ('b/a', 'thanksgiving', 'thanksgiving prayers'):
                        display_label = "B/A"
                    else:
                        display_label = label
                    label_with_colon = f"{display_label}:"
                
                if not first_para:
                    p = tf.add_paragraph()
                else:
                    p = tf.paragraphs[0]
                    first_para = False
                
                p.space_before = Pt(2)
                p.space_after = Pt(2)
                
                if label_with_colon:
                    spaces_needed = max(1, LABEL_COLUMN_WIDTH - len(label_with_colon))
                else:
                    spaces_needed = LABEL_COLUMN_WIDTH
                
                p.clear()
                run1 = p.add_run()
                run1.text = label_with_colon
                run1.font.name = "Arial"
                run1.font.size = Pt(14)
                run1.font.bold = True
                
                run2 = p.add_run()
                padding = " " * spaces_needed
                if song.get('hymn_num'):
                    run2.text = f"{padding}{song['hymn_num']} {song.get('title_hint', '')}"
                else:
                    # If no hymn number, show just the title if available
                    title = song.get('title_hint', '')
                    run2.text = f"{padding}{title}" if title else ""
                run2.font.name = "Arial"
                run2.font.size = Pt(14)
                run2.font.bold = False
            
            break


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN GENERATION LOGIC
# ═══════════════════════════════════════════════════════════════════════════════

def generate_presentation(song_list, output_filename=None, service_date=None):
    """
    Generate a PPT presentation from a list of songs.
    
    song_list: list of dicts with keys:
        - hymn_num: int or str (the hymn number, or empty for Message)
        - label: str (Opening, ThanksGiving, Offertory, Message, Confession, Communion, Closing)
        - title_hint: str (optional song title hint)
    service_date: str (optional service date in format "DD Month YYYY", e.g., "04 January 2026")
    """
    # Reset used slide ranges for this generation
    global USED_SLIDE_RANGES
    USED_SLIDE_RANGES = {}
    
    # Refresh search directories based on current working directory
    global MALAYALAM_SEARCH_DIRS
    MALAYALAM_SEARCH_DIRS = get_search_dirs()
    
    if output_filename is None:
        today = datetime.now().strftime("%d %b %Y")
        output_filename = f"{today} - Generated Malayalam HCS.pptx"

    output_path = os.path.join(BASE_DIR, output_filename)
    
    print(f"  Language: Malayalam")
    print(f"  Search directories: {get_search_dirs()}")

    def normalize_service_date(date_text):
        if not date_text:
            return ""
        formats = [
            "%d %b %Y",
            "%d %B %Y",
            "%d-%b-%Y",
            "%d-%B-%Y",
            "%d/%m/%Y",
            "%d-%m-%Y",
        ]
        for fmt in formats:
            try:
                parsed = datetime.strptime(date_text, fmt)
                return parsed.strftime("%d %B %Y")
            except ValueError:
                continue
        return date_text

    # Load template and create presentation
    template_path = find_template_ppt()
    if not template_path:
        raise FileNotFoundError(
            "Malayalam Template PPT not found.\n\n"
            "Please ensure the source folder contains '4 Jan 2026.pptx' under the Malayalam HCS folder."
        )

    print(f"  Template: {template_path}")
    prs = Presentation(template_path)

    # Update date in slide masters/layouts if provided
    normalized_date = normalize_service_date(service_date)
    if normalized_date:
        date_pattern = re.compile(r"\d{1,2}\s+(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}")
        for master in prs.slide_masters:
            for shape in master.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if date_pattern.search(run.text):
                            run.text = date_pattern.sub(normalized_date, run.text)
        for layout in prs.slide_layouts:
            for shape in layout.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if date_pattern.search(run.text):
                            run.text = date_pattern.sub(normalized_date, run.text)

    # Remove all existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # Get layouts
    title_layout = None
    blank_layout = None
    
    for layout in prs.slide_layouts:
        if layout.name == "Malayalam HC - 03 Jan 2021":
            title_layout = layout
        elif layout.name == "1_Blank" and blank_layout is None:
            blank_layout = layout

    if title_layout is None:
        title_layout = prs.slide_layouts[0]
    if blank_layout is None:
        blank_layout = prs.slide_layouts[-1]
    
    print(f"  Using title layout: '{title_layout.name}'")
    print(f"  Using content layout: '{blank_layout.name}'")

    slide_counter = 1

    # Create summary slide
    print("\n📋 Creating summary slide...")
    create_summary_slide(prs, title_layout, song_list, normalized_date)
    slide_counter += 1

    # Process each song section
    print("\n🎵 Processing songs...\n")

    for song_info in song_list:
        hymn_num = str(song_info["hymn_num"]) if song_info["hymn_num"] else ""
        label = song_info["label"]
        title_hint = song_info.get("title_hint", "")

        print(f"── {label}: Hymn No {hymn_num} {title_hint} ──" if hymn_num else f"── {label} ──")

        # Route to appropriate processor based on label
        label_lower = label.lower()
        
        # Each processor now returns (slide_counter, extracted_title)
        if label_lower == "opening":
            slide_counter, extracted_title = process_opening_song(prs, title_layout, blank_layout, hymn_num, title_hint, slide_counter, )
            if extracted_title:
                song_info["title_hint"] = extracted_title
        elif label_lower in ("thanksgiving", "thanksgiving prayers", "b/a"):
            slide_counter, extracted_title = process_thanksgiving_prayers(prs, title_layout, blank_layout, hymn_num, title_hint, slide_counter, )
            if extracted_title:
                song_info["title_hint"] = extracted_title
        elif label_lower == "offertory":
            slide_counter, extracted_title = process_offertory(prs, title_layout, blank_layout, hymn_num, title_hint, slide_counter, )
            if extracted_title:
                song_info["title_hint"] = extracted_title
        elif label_lower == "message":
            slide_counter = process_message(prs, title_layout, slide_counter)
        elif label_lower == "confession":
            slide_counter, extracted_title = process_confession(prs, title_layout, blank_layout, hymn_num, title_hint, slide_counter, )
            if extracted_title:
                song_info["title_hint"] = extracted_title
        elif label_lower in ("communion", "holy communion"):
            slide_counter, extracted_title = process_holy_communion(prs, title_layout, blank_layout, hymn_num, title_hint, slide_counter, )
            if extracted_title:
                song_info["title_hint"] = extracted_title
        elif label_lower == "closing":
            slide_counter, extracted_title = process_closing_hymn(prs, title_layout, blank_layout, hymn_num, title_hint, slide_counter, )
            if extracted_title:
                song_info["title_hint"] = extracted_title
        elif label_lower == "dedication":
            slide_counter, extracted_title = process_generic_song(prs, title_layout, blank_layout, "Dedication", hymn_num, title_hint, slide_counter, )
            if extracted_title:
                song_info["title_hint"] = extracted_title
        else:
            slide_counter, extracted_title = process_generic_song(prs, title_layout, blank_layout, label, hymn_num, title_hint, slide_counter, )
            if extracted_title:
                song_info["title_hint"] = extracted_title

    # Now update the summary slide with extracted titles from the created slides
    update_summary_slide_from_slides(prs, song_list)

    # Save
    prs.save(output_path)
    print(f"\n{'═' * 60}")
    print(f"✅ Presentation saved: {output_path}")
    print(f"   Total slides: {slide_counter - 1}")
    print(f"{'═' * 60}")

    return output_path


# ═══════════════════════════════════════════════════════════════════════════════
# USER INTERFACE
# ═══════════════════════════════════════════════════════════════════════════════

def get_song_list_from_user():
    """Interactive prompt to get the song list from the user."""
    print("╔══════════════════════════════════════════════════════════╗")
    print("║       Church Songs PPT Generator                         ║")
    print("║       (Malayalam Holy Communion Service)                 ║")
    print("╚══════════════════════════════════════════════════════════╝")
    print()
    print("Enter songs one by one. For each song, provide:")
    print("  1. The hymn number (e.g., 91, or leave empty for Message)")
    print("  2. The label (Opening, ThanksGiving, Offertory, Message, Confession, Communion, Closing)")
    print("  3. Optionally, a title hint")
    print()
    print("Type 'done' when finished, or 'example' to see a sample.")
    print()

    songs = []
    while True:
        try:
            entry = input(f"Song {len(songs) + 1} (hymn_num label [title_hint]) or 'done': ").strip()
        except (EOFError, KeyboardInterrupt):
            print()
            break

        if entry.lower() == "done":
            break

        if entry.lower() == "example":
            print("  Example: 91 Opening Vaazthin vaazhthin")
            print("  Example: 110 ThanksGiving")
            print("  Example: 420 Offertory")
            print("  Example: Message")
            print("  Example: 211 Confession")
            print("  Example: 313 Communion")
            print("  Example: 427 Closing")
            continue

        if not entry:
            continue

        # Special case: Message (no hymn number)
        if entry.lower() == "message":
            songs.append({"hymn_num": "", "label": "Message", "title_hint": ""})
            print(f"  ✓ Added: Message")
            continue

        parts = entry.split(None, 2)
        if len(parts) < 2:
            print("  ⚠ Please provide at least: hymn_number label (or just 'Message')")
            continue

        hymn_num = parts[0]
        label = parts[1]
        title_hint = parts[2] if len(parts) > 2 else ""

        try:
            int(hymn_num)
        except ValueError:
            print(f"  ⚠ '{hymn_num}' is not a valid hymn number")
            continue

        songs.append({
            "hymn_num": hymn_num,
            "label": label.capitalize(),
            "title_hint": title_hint,
        })
        print(f"  ✓ Added: {label} - Hymn No {hymn_num} {title_hint}")

    return songs


def main():
    """Main entry point."""
    if len(sys.argv) > 1 and sys.argv[1] == "--batch":
        songs = []
        language = "Malayalam"  # Default language
        service_date = None  # Service date from batch file
        if len(sys.argv) > 2:
            batch_file = sys.argv[2]
            print(f"Reading songs from: {batch_file}")
            with open(batch_file, "r") as f:
                for line in f:
                    line = line.strip()
                    if not line:
                        continue
                    # Parse language directive
                    if line.startswith("# Language:"):
                        language = line.split(":", 1)[1].strip()
                        print(f"  Language set to: {language}")
                        continue
                    # Parse date directive
                    if line.lower().startswith("# date:") or line.lower().startswith("date:"):
                        service_date = line.split(":", 1)[1].strip()
                        print(f"  Service date set to: {service_date}")
                        continue
                    if line.startswith("#"):
                        continue
                    parts = line.split("|")
                    if len(parts) >= 2:
                        songs.append({
                            "hymn_num": parts[0].strip(),
                            "label": parts[1].strip(),
                            "title_hint": parts[2].strip() if len(parts) > 2 else "",
                        })
                    elif parts[0].strip().lower() == "message":
                        songs.append({"hymn_num": "", "label": "Message", "title_hint": ""})
        else:
            print("Usage: python3 generate_hcs_ppt.py --batch songs.txt [output.pptx]")
            return
    else:
        songs = get_song_list_from_user()
        language = "Malayalam"  # Default for interactive mode

    if not songs:
        print("No songs specified. Exiting.")
        return

    print(f"\n📝 Song list ({len(songs)} songs):")
    for i, s in enumerate(songs, 1):
        if s['hymn_num']:
            print(f"  {i}. {s['label']}: Hymn No {s['hymn_num']} {s['title_hint']}")
        else:
            print(f"  {i}. {s['label']}")

    output_name = None
    is_batch = len(sys.argv) > 1 and sys.argv[1] == "--batch"
    
    if not is_batch:
        try:
            custom_name = input("\nOutput filename (press Enter for default): ").strip()
            if custom_name:
                output_name = custom_name
        except (EOFError, KeyboardInterrupt):
            pass
    
    if len(sys.argv) > 3:
        output_name = sys.argv[3]

    # Pass service_date if it was set from batch file
    if 'service_date' in locals() and service_date:
        generate_presentation(songs, output_name, service_date)
    else:
        generate_presentation(songs, output_name)


if __name__ == "__main__":
    main()
