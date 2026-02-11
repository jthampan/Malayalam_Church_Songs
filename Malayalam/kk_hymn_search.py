#!/usr/bin/env python3
"""
KK.pptx-specific hymn search logic
Handles special format of KK hymn book with corner numbers and footer patterns
"""

import re
from pptx import Presentation


def find_hymn_in_kk_pptx(pptx_path, hymn_number):
    """
    Search for a hymn in KK.pptx format
    
    Args:
        pptx_path: Path to the KK.pptx file
        hymn_number: String hymn number to search for (e.g., "8", "143")
    
    Returns:
        tuple: (title_slide_index, content_slide_indices, extracted_title)
               Returns (None, [], "") if not found
    """
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error opening {pptx_path}: {e}")
        return None, [], ""
    
    target_hymn_num = str(hymn_number)
    collecting = False
    content_indices = []
    title_slide_idx = None
    footer_hymn_num_at_start = None
    last_slide_count = None
    extracted_title = ""
    
    # Get slide width to determine position thresholds
    slide_width = prs.slide_width if hasattr(prs, 'slide_width') else 9144000
    right_corner_threshold = slide_width * 0.85  # Skip rightmost 15%
    
    for i, slide in enumerate(prs.slides):
        all_text = ""
        shape_texts = []
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                all_text += " " + text
                if text:
                    shape_texts.append(text)
        
        # Check if this slide has our target hymn number
        # ONLY check shapes on LEFT or CENTER of slide, NOT right corner
        hymn_match = False
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                
                # Skip shapes on the FAR RIGHT corner (rightmost 15%)
                if hasattr(shape, 'left') and shape.left > right_corner_threshold:
                    continue
                
                # Check for patterns like "Holy Communion – 143" or "Section - 143"
                if '–' in text or '-' in text or '—' in text:
                    dash_match = re.search(r'[–\-—]\s*(\d+)', text)
                    if dash_match and dash_match.group(1) == target_hymn_num:
                        hymn_match = True
                
                if hymn_match:
                    break
                
                # Check for standalone number at start/middle (not right end)
                words = text.split()
                mid_point = len(words) // 2
                first_half = words[:mid_point + 1] if mid_point > 0 else words[:1]
                
                for word in first_half:
                    clean_word = word.strip('.,;:!?"\'()[]{}–-')
                    if clean_word == target_hymn_num or clean_word == f"{target_hymn_num}.":
                        hymn_match = True
                        break
                
                if hymn_match:
                    break
        
        # Start collecting when we find the hymn
        if hymn_match and not collecting:
            title_slide_idx = i
            collecting = True
            
            # Extract title from first slide (largest non-footer text)
            max_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if len(text) > len(max_text) and not re.match(r'^\d+\.?$', text) and 'of' not in text.lower():
                        max_text = text
            extracted_title = max_text
            
            # Extract "X of Y" from first slide to know total slides
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    slide_counter = re.search(r'[-–:]\s*(\d+)\s+of\s+(\d+)', text, re.IGNORECASE)
                    if slide_counter:
                        last_slide_count = int(slide_counter.group(2))
                    hymn_in_footer = re.search(r'Hymn\s*#?\s*(\d+)', text, re.IGNORECASE)
                    if hymn_in_footer:
                        footer_hymn_num_at_start = hymn_in_footer.group(1)
            
            content_indices.append(i)
            continue
        
        # While collecting, check footer for boundaries
        if collecting:
            current_slide_num = None
            current_total = None
            current_footer_hymn = None
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    slide_counter = re.search(r'[-–:]\s*(\d+)\s+of\s+(\d+)', text, re.IGNORECASE)
                    if slide_counter:
                        current_slide_num = int(slide_counter.group(1))
                        current_total = int(slide_counter.group(2))
                    hymn_in_footer = re.search(r'Hymn\s*#?\s*(\d+)', text, re.IGNORECASE)
                    if hymn_in_footer:
                        current_footer_hymn = hymn_in_footer.group(1)
            
            # STOP if: hymn number in footer changed
            if current_footer_hymn and footer_hymn_num_at_start:
                if current_footer_hymn != footer_hymn_num_at_start:
                    break
            
            # STOP if: completed expected slides and now seeing "1 of X" again
            if last_slide_count and current_slide_num == 1 and len(content_indices) >= last_slide_count:
                break
            
            content_indices.append(i)
    
    if not content_indices:
        return None, [], ""
    
    return title_slide_idx, content_indices, extracted_title


if __name__ == "__main__":
    # Scan and list all hymns in KK.pptx
    import sys
    import os
    
    # Find KK.pptx
    kk_path = None
    search_root = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'onedrive_git_local')
    
    if not os.path.exists(search_root):
        search_root = os.path.join(os.path.dirname(__file__), '..', 'onedrive_git_local')
    
    for root, dirs, files in os.walk(search_root):
        for f in files:
            if 'KK' in f.upper() and f.endswith('.pptx'):
                kk_path = os.path.join(root, f)
                break
        if kk_path:
            break
    
    if not kk_path:
        print("KK.pptx not found!")
        sys.exit(1)
    
    print(f"Scanning: {os.path.basename(kk_path)}")
    print("=" * 80)
    
    try:
        prs = Presentation(kk_path)
        slide_width = prs.slide_width if hasattr(prs, 'slide_width') else 9144000
        right_corner_threshold = slide_width * 0.85
        
        hymns_found = {}
        
        for i, slide in enumerate(prs.slides):
            # Look for hymn numbers in left/center area only
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    
                    # Skip right corner shapes
                    if hasattr(shape, 'left') and shape.left > right_corner_threshold:
                        continue
                    
                    # Look for standalone numbers or "Section - Number" patterns
                    # Check for dash patterns
                    dash_match = re.search(r'[–\-—]\s*(\d+)', text)
                    if dash_match:
                        hymn_num = dash_match.group(1)
                        if hymn_num not in hymns_found:
                            hymns_found[hymn_num] = i
                        continue
                    
                    # Check for standalone number at start
                    words = text.split()
                    if words:
                        mid_point = len(words) // 2
                        first_half = words[:mid_point + 1] if mid_point > 0 else words[:1]
                        
                        for word in first_half:
                            clean_word = word.strip('.,;:!?"\'()[]{}–-')
                            if clean_word.isdigit() and len(clean_word) <= 3:
                                if clean_word not in hymns_found:
                                    hymns_found[clean_word] = i
        
        # Sort by hymn number
        sorted_hymns = sorted(hymns_found.items(), key=lambda x: int(x[0]))
        
        print(f"\nFound {len(sorted_hymns)} hymn numbers:\n")
        
        for hymn_num, slide_idx in sorted_hymns:
            # Get title using our search function
            title_idx, content_idx, title = find_hymn_in_kk_pptx(kk_path, hymn_num)
            
            if content_idx:
                title_preview = title[:60] + "..." if len(title) > 60 else title
                print(f"Hymn {hymn_num:>3}: {len(content_idx):>2} slides at index {title_idx:>3} | {title_preview}")
            else:
                print(f"Hymn {hymn_num:>3}: Not found (detected at slide {slide_idx})")
        
        print("\n" + "=" * 80)
        
    except Exception as e:
        print(f"Error scanning KK.pptx: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
