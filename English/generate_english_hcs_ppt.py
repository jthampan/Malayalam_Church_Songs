#!/usr/bin/env python3
"""
English Church Songs PPT Generator
===================================
Searches for English songs by hymn number in existing PPT files, then generates
a new presentation using the template from "8 Feb 2026" English HCS PPT.

Sections supported:
  - Opening Song
  - ThanksGiving Prayers  
  - Offertory (with QR code)
  - Message (title slide only)
  - Confession
  - Holy Communion (with Holy Communion image)
  - Closing Hymn

Usage:
    python3 generate_english_hcs_ppt.py --batch songs.txt "Output Name.pptx"
    
    songs.txt format:
        hymn_num|label|title_hint
        # Lines starting with # are comments
        # Use 'Message' label for message-only title slide

Images are extracted and saved to the 'images/' folder.
"""

import os
import sys
import re
from copy import deepcopy
from datetime import datetime
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
import subprocess
import shutil


# ═══════════════════════════════════════════════════════════════════════════════
# CONFIGURATION
# ═══════════════════════════════════════════════════════════════════════════════

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PARENT_DIR = os.path.dirname(BASE_DIR)
IMAGES_DIR = os.path.join(PARENT_DIR, "images")

# Create images directory if it doesn't exist
os.makedirs(IMAGES_DIR, exist_ok=True)

# Path to bundled/local hymn files folder
ONEDRIVE_GIT_LOCAL = os.path.join(PARENT_DIR, "onedrive_git_local")

# ═══════════════════════════════════════════════════════════════════════════════
# FILE PATHS
# ═══════════════════════════════════════════════════════════════════════════════

TEMPLATE_PPT = os.path.join(
    ONEDRIVE_GIT_LOCAL,
    "Holy Communion Services - Slides",
    "English HCS",
    "2026 - Eng HCS",
    "8 Feb 2026.pptx",
)

# Image paths
HOLY_COMMUNION_IMAGE = os.path.join(IMAGES_DIR, "holy_communion.jpg")

# Directories to search for existing PPTs
SEARCH_DIRS = [
    ONEDRIVE_GIT_LOCAL,
    BASE_DIR,
]

def resolve_image_path(filename):
    """Find an image in the source folder first, then fallback to packaged images."""
    candidates = [
        os.path.join(os.getcwd(), "images", filename),
        os.path.join(BASE_DIR, "images", filename),
        os.path.join(PARENT_DIR, "images", filename),
    ]
    for path in candidates:
        if os.path.exists(path):
            return path
    return candidates[0]

def find_template_ppt():
    """Find the English template PPT in the selected source folder or fallback paths."""
    template_name = "8 Feb 2026.pptx"
    # Search in: user-provided folder, PARENT_DIR, onedrive_git_local
    search_roots = [os.getcwd(), PARENT_DIR, ONEDRIVE_GIT_LOCAL]

    # Prefer the current working directory (set by the GUI to the source folder)
    for root in search_roots:
        if not root or not os.path.isdir(root):
            continue
        for dirpath, _, filenames in os.walk(root):
            if template_name in filenames:
                return os.path.join(dirpath, template_name)

    # Final fallback to the static path if it exists
    if os.path.exists(TEMPLATE_PPT):
        return TEMPLATE_PPT

    return None

# ─── Font / layout settings (from template analysis) ─────────────────────────
TITLE_FONT = "Gabriola"
TITLE_SIZE = Pt(60)

HEADER_FONT = "Segoe UI"
CONTENT_FONT = "Arial"
CONTENT_SIZE = Pt(22)
CONTENT_COLOR = RGBColor(0x00, 0x00, 0x00)

FOOTER_FONT = "Segoe UI"
FOOTER_SIZE = Pt(14)

# Title bar styling (from template)
TITLE_BAR_HEIGHT = Emu(486000)
TITLE_BAR_COLOR = RGBColor(232, 211, 211)  # #E8D3D3
TITLE_BAR_FONT_SIZE = Pt(24)

# Maximum lines per content slide
MAX_LINES_PER_SLIDE = 8

# Holy Communion image position (from template)
HC_IMAGE_LEFT = Inches(1.673)
HC_IMAGE_TOP = Inches(0.772)
HC_IMAGE_WIDTH = Inches(6.654)
HC_IMAGE_HEIGHT = Inches(4.437)


# ═══════════════════════════════════════════════════════════════════════════════
# IMAGE EXTRACTION AND MANAGEMENT
# ═══════════════════════════════════════════════════════════════════════════════

def ensure_holy_communion_image():
    """Ensure the Holy Communion image exists in the images folder."""
    target_image = resolve_image_path("holy_communion.jpg")
    if os.path.exists(target_image):
        return True
    
    # Try to extract from a recent English HCS PPT
    source_pptx = os.path.join(
        PARENT_DIR,
        "onedrive_git_local",
        "Holy Communion Services - Slides",
        "English HCS",
        "2026 - Eng HCS",
        "8 Feb 2026.pptx"
    )
    
    if not os.path.exists(source_pptx):
        print(f"  ⚠ Warning: Cannot find source for Holy Communion image")
        return False
    
    try:
        prs = Presentation(source_pptx)
        # Search for a slide with Holy Communion image
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = shape.image
                    os.makedirs(os.path.dirname(target_image), exist_ok=True)
                    with open(target_image, "wb") as f:
                        f.write(img.blob)
                    print(f"  ✓ Extracted Holy Communion image to {target_image}")
                    return True
    except Exception as e:
        print(f"  ⚠ Could not extract Holy Communion image: {e}")
    
    return False


# ═══════════════════════════════════════════════════════════════════════════════
# PPT SEARCH FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════

def get_search_dirs(language="English"):
    """
    Get search directories based on current working directory and language.
    
    Search logic:
    1. If user provides a path (cwd != BASE_DIR), search for language-specific folder there
    2. Always include onedrive_git_local as fallback
    
    Args:
        language: "Malayalam" or "English" - determines which HCS folder to look for
    """
    # Use current working directory as base (set by GUI or default to script location)
    cwd = os.getcwd()
    
    # Determine language folder name
    if language.lower() == "english":
        lang_folder = "English HCS"
    else:
        lang_folder = "Malayalam HCS"
    
    search_dirs = []
    
    # Check if user provided a custom path (cwd is different from script location)
    user_provided_path = cwd != BASE_DIR and cwd != PARENT_DIR
    
    if user_provided_path:
        # User provided a path - search for language-specific folder
        
        # Check if cwd already IS the language folder
        if os.path.isdir(cwd) and os.path.basename(cwd).lower() == lang_folder.lower():
            search_dirs.append(cwd)
        else:
            # Check if cwd contains the language folder as immediate child
            immediate_child = os.path.join(cwd, lang_folder)
            if os.path.isdir(immediate_child):
                search_dirs.append(immediate_child)
            else:
                # Search recursively for the language folder (e.g., user gave "Holy Communion Services - Slides")
                for dirpath, dirnames, _ in os.walk(cwd):
                    for d in dirnames:
                        if d.lower() == lang_folder.lower():
                            search_dirs.append(os.path.join(dirpath, d))
                    # Stop after finding first match (don't search too deep)
                    if search_dirs:
                        break
                
                # If language folder not found, use the user path as-is (search everything)
                if not search_dirs and os.path.isdir(cwd):
                    search_dirs.append(cwd)
    
    # Always add onedrive_git_local as fallback (from exe bundle or downloaded)
    onedrive_path = os.path.join(ONEDRIVE_GIT_LOCAL, "Holy Communion Services - Slides", lang_folder)
    if os.path.isdir(onedrive_path) and onedrive_path not in search_dirs:
        search_dirs.append(onedrive_path)
    elif os.path.isdir(ONEDRIVE_GIT_LOCAL):
        # Search for language folder under onedrive_git_local
        for dirpath, dirnames, _ in os.walk(ONEDRIVE_GIT_LOCAL):
            for d in dirnames:
                if d.lower() == lang_folder.lower():
                    candidate = os.path.join(dirpath, d)
                    if candidate not in search_dirs:
                        search_dirs.append(candidate)
            if any(lang_folder.lower() in p.lower() for p in search_dirs):
                break
    
    # Final fallback to BASE_DIR if nothing found
    if not search_dirs:
        search_dirs = [BASE_DIR]
    
    return search_dirs


# Initialize with default English directories
ENGLISH_SEARCH_DIRS = get_search_dirs("English")

def find_all_pptx_files(search_dirs):
    """Recursively find all .pptx files in English directories."""
    # Search only in onedrive_git_local directory
    onedrive_path = os.path.join(ONEDRIVE_GIT_LOCAL, "Holy Communion Services - Slides", "English HCS")
    
    pptx_files = []
    if os.path.isdir(onedrive_path):
        for root, dirs, files in os.walk(onedrive_path):
            for f in files:
                if f.endswith(".pptx") and not f.startswith("~$"):
                    pptx_files.append(os.path.join(root, f))
    return pptx_files


def is_title_slide(all_text, slide=None):
    """
    Detect if a slide is a title slide based on text content and visual elements.
    
    Title slides typically have:
    1. Very short text (< 150 chars)
    2. Section label keywords (Opening, Midnight, Hymn, Communion, etc.)
    3. 1-2 lines (section label + song title)
    4. No verse structure (no commas, periods, or multiple substantial lines)
    5. Large background image (covers most of the slide)
    6. Colored box/overlay (for section titles)
    
    Args:
        all_text: Combined text from all shapes on the slide
        slide: Optional slide object for visual analysis
        
    Returns:
        bool: True if this appears to be a title slide, False otherwise
    """
    clean_text = all_text.strip()
    
    # If slide object provided, check for visual indicators first
    if slide is not None:
        # Check for large background images (common in title slides)
        has_large_background = False
        has_colored_box = False
        
        for shape in slide.shapes:
            # Check for pictures (background images)
            if hasattr(shape, 'shape_type') and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                # Check if image is large (covers most of slide - likely a background)
                if hasattr(shape, 'width') and hasattr(shape, 'height'):
                    # Background images are typically > 8 inches wide and > 5 inches tall
                    if shape.width > 8000000 and shape.height > 5000000:  # In EMUs (1 inch = 914400 EMUs)
                        has_large_background = True
            
            # Check for colored rounded rectangles (title box overlays)
            if hasattr(shape, 'shape_type') and shape.shape_type == 1:  # MSO_SHAPE_TYPE.AUTO_SHAPE
                if hasattr(shape, 'auto_shape_type') and shape.auto_shape_type == 5:  # ROUNDED_RECTANGLE
                    # Title boxes are typically medium-large rectangles
                    if hasattr(shape, 'width') and hasattr(shape, 'height'):
                        if shape.width > 6000000 and shape.height > 2000000:  # Large colored box
                            has_colored_box = True
        
        # If slide has large background + colored box, it's very likely a title slide
        if has_large_background and has_colored_box:
            return True
        
        # If has large background + short text with title indicators, likely a title slide
        if has_large_background and len(clean_text) < 150:
            title_indicators = ["Hymn", "Opening", "Closing", "Midnight", "Confession", 
                               "Thanksgiving", "Offertory", "Dedication", "Song", "Communion"]
            if any(indicator in all_text for indicator in title_indicators):
                return True
    
    # Special check for "Holy Communion – [Song Title]" pattern with short text
    if re.search(r"Holy\s+Communion\s*[–-]\s*[A-Za-z]", all_text, re.IGNORECASE):
        if len(clean_text) < 150:
            # Count lines - title slides have 1-2 lines, verses have many
            lines = [line.strip() for line in all_text.split('\n') if len(line.strip()) > 10]
            # If only 1-2 lines, it's a title slide (even if title has commas like "Wonderful, Merciful")
            if len(lines) <= 2:
                return True
    
    # Fallback to text-based detection
    # Title slides are short
    if len(clean_text) >= 150:
        return False
    
    # Check for title slide indicators
    title_indicators = ["Hymn", "Opening", "Closing", "Midnight", "Confession", 
                       "Thanksgiving", "Offertory", "Dedication", "Song", "Communion"]
    has_indicator = any(indicator in all_text for indicator in title_indicators)
    
    if not has_indicator:
        return False
    
    # Check structure - title slides have 1-2 lines with no punctuation
    lines = [line.strip() for line in all_text.split('\n') if len(line.strip()) > 5]
    
    # Count verse structure indicators
    has_commas = sum(1 for line in lines if ',' in line)
    has_periods = sum(1 for line in lines if line.endswith(('.', '!', '?')))
    
    # Title slides typically have:
    # - 1-2 lines (section name + optional song title)
    # - No punctuation (no commas or periods)
    if len(lines) <= 2 and has_commas == 0 and has_periods == 0:
        return True
    
    return False


def find_song_slide_indices_in_pptx(pptx_path, target_hymn_num="", song_title_hint=""):
    """
    Find the slide indices for a specific hymn in a PPTX file.
    
    Search logic:
    - If only hymn number provided: search by "Hymn No XXX" or "Song No XXX"
    - If only song title provided: search by song title words
    - If both provided: search by either (hymn number OR song title)
    
    Args:
        pptx_path: Path to the PPTX file
        target_hymn_num: The hymn number to search for (e.g., "171") - optional
        song_title_hint: Song title to search for - optional
    
    Returns (title_slide_idx, [content_slide_indices], extracted_title) or (None, [], "").
    """
    target = str(target_hymn_num) if target_hymn_num else ""
    try:
        prs = Presentation(pptx_path)
    except Exception:
        return None, [], ""

    # Need at least one search criteria
    if not target and not song_title_hint:
        return None, [], ""

    title_slide_idx = None
    content_indices = []
    extracted_title = ""
    collecting = False
    
    # Prepare title search - normalize and get first 3-4 consecutive words
    title_search_phrase = ""
    if song_title_hint and len(song_title_hint) > 2:
        # Normalize: lowercase, remove special chars, keep only alphanumeric and spaces
        normalized_title = re.sub(r'[^a-z0-9\s]', '', song_title_hint.lower())
        # Get first 3-4 words (at least 3 chars each) for consecutive matching
        words = [w for w in normalized_title.split() if len(w) > 2]
        if words:
            # Use first 3 or 4 words (whichever gives us more)
            num_words = min(4, len(words)) if len(words) >= 4 else min(3, len(words))
            title_search_phrase = ' '.join(words[:num_words])

    def slide_has_lyrics(text):
        # For English, check for sufficient text content (lyrics typically have more text)
        ascii_letters = sum(1 for c in text if c.isalpha() and c.isascii())
        return ascii_letters >= 30
    
    def is_image_only_slide(slide):
        """Check if slide is primarily just an image without meaningful text content."""
        # Count text content
        total_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                total_text += shape.text_frame.text.strip()
        
        # Remove common labels like title bar text
        meaningful_text = total_text
        for label in ["Holy Communion Hymn", "Opening Hymn", "Thanksgiving", "Confession", 
                      "Closing Hymn", "Offertory", "Communion"]:
            meaningful_text = meaningful_text.replace(label, "")
        
        # Remove slide numbers like "30", "31"
        meaningful_text = re.sub(r'\b\d{1,3}\b', '', meaningful_text)
        meaningful_text = re.sub(r'\d+\s*:\s*\d+\s+of\s+\d+', '', meaningful_text)  # Footer
        
        # Count actual text content (letters)
        letter_count = sum(1 for c in meaningful_text if c.isalpha())
        
        # If less than 30 letters of actual content, it's likely just an image slide
        return letter_count < 30

    def build_hymn_pattern(num):
        # More comprehensive pattern to match various hymn number formats
        # Matches: "Hymn No 171", "Song No. 633", "Hymn – 171", "Song no: 40", "(36)", etc.
        escaped_num = re.escape(num)
        patterns = [
            # Standard formats with "No" or "No."
            r"(?:Hymn|Song)\s*[Nn]o\.?\s*[\:\-–]?\s*\(?" + escaped_num + r"\)?(?:\s|\b|,)",
            # Formats with just dash: "Hymn – 171" or "Song – 633"
            r"(?:Hymn|Song)\s*[-–]\s*" + escaped_num + r"(?:\s|\b|,)",
            # Formats in context: "Confession – Song No. 633"
            r"(?:Confession|Offertory|Communion)\s*[-–]\s*(?:Song|Hymn)\s*[Nn]o\.?\s*" + escaped_num + r"(?:\s|\b|,)",
            # Parentheses format at end of title: "A Christian home (36)"
            r"\(" + escaped_num + r"\)",
        ]
        # Combine all patterns with OR
        return "(" + "|".join(patterns) + ")"
    
    def find_hymn_in_text(target_num, text):
        """Find hymn number in text using regular service PPT patterns."""
        if not target_num:
            return False
        
        hymn_pattern = build_hymn_pattern(target_num)
        return bool(re.search(hymn_pattern, text, re.IGNORECASE))

    for i, slide in enumerate(prs.slides):
        all_text = ""
        first_lines = []  # Collect first line from all text boxes

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                all_text += " " + text
                # Capture first line from each text box
                if text:
                    first_lines.append(text.split('\n')[0].strip())

        # Skip summary/index slides that list multiple hymns
        hymn_count = len(re.findall(r'\b\d{2,3}\b', all_text))
        if hymn_count > 4:
            continue
        
        # Check for hymn number match (if hymn number provided)
        hymn_match = False
        if target:
            hymn_match = find_hymn_in_text(target, all_text)
        
        # Check for song title match (if title provided) - search all first lines
        title_match = False
        if title_search_phrase:
            # Check first line of each text box
            for first_line in first_lines:
                normalized_first_line = re.sub(r'[^a-z0-9\s]', '', first_line.lower())
                if title_search_phrase in normalized_first_line:
                    title_match = True
                    break

        # Determine if this slide is a match
        is_match = False
        if target and song_title_hint:
            # Both provided: either one matches
            is_match = hymn_match or title_match
        elif target:
            # Only hymn number: must match hymn
            is_match = hymn_match
        else:
            # Only title: must match title
            is_match = title_match

        if is_match and not collecting:
            # Found start of our song - this is the title slide
            # Store the hymn number found on this slide (if any) for tracking
            # Use multiple patterns to extract hymn number (including parentheses format)
            found_hymn_num = ""
            hymn_patterns = [
                r"(?:Hymn|Song)\s*(?:[Nn]o\.?:?\s*)?[-–]?\s*\(?(\d+)\)?",
                r"\((\d+)\)",  # Parentheses format like "(36)"
            ]
            for hp in hymn_patterns:
                hymn_found_on_title = re.search(hp, all_text, re.IGNORECASE)
                if hymn_found_on_title:
                    found_hymn_num = hymn_found_on_title.group(1)
                    break
            
            # If searching by hymn number only, ensure we found the correct hymn number
            if target and not song_title_hint:
                # Searching by hymn number only - verify we found the right number
                if found_hymn_num != target:
                    # This slide doesn't have our hymn number, skip it
                    continue
            # Also store what section we're in (from title slide)
            found_section = ""
            for sec in ["Opening", "Confession", "Offertory", "Thanksgiving", "Communion", "Closing", "Dedication", "B/A"]:
                if sec.lower() in all_text.lower():
                    found_section = sec
                    break
            
            title_slide_idx = i
            collecting = True
            
            # Extract title from the title slide (find largest English text)
            max_text = ""
            for shape in prs.slides[i].shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    # Skip footer patterns and slide numbers
                    if re.search(r'^\d+\s*:\s*\d+\s+of\s+\d+', text):  # Footer like "30 : 31 of 106"
                        continue
                    if re.match(r'^\d{1,3}$', text):  # Just a slide number
                        continue
                    if any(kw in text for kw in ['Hymn', 'Song']) and len(text) < 30:  # Skip "Hymn No 313" headers
                        continue
                    if any(sec in text for sec in ["Opening", "Confession", "Offertory", "Thanksgiving", "Communion", "Closing"]) and len(text) < 40:
                        continue
                    # Check if text has English content
                    has_english = len([c for c in text if c.isalpha() and c.isascii()]) > 10
                    if has_english and len(text) > len(max_text):
                        max_text = text
            if max_text:
                # Take only the first line as the title, limited to first 3 words
                first_line = max_text.split('\n')[0].strip()
                words = first_line.split()
                extracted_title = ' '.join(words[:3]) if len(words) > 3 else first_line
            
            # Check if title slide has lyrics (compact single-slide format)
            # If it has significant content, add it as a content slide too
            if slide_has_lyrics(all_text) and not is_image_only_slide(slide):
                content_indices.append(i)
            
            # Don't add title slide again in the collecting loop
            continue

        if collecting:
            clean_text = all_text.strip()
            
            # Skip very short slides (likely section dividers) but don't stop collecting
            if len(clean_text) < 20:
                continue
            
            # Check if this is a NEW section (different from where we started)
            # e.g., "Holy Communion - Hymn No 42" when we started at "Confession - Hymn No 42"
            # BUT: Only stop if this slide doesn't have our hymn number
            # (same hymn can appear in multiple sections, e.g., Confession and Communion)
            is_different_section = False
            for sec in ["Opening", "Confession", "Offertory", "Thanksgiving", "Communion", "Closing", "Dedication"]:
                if sec.lower() in all_text.lower() and ("Hymn" in all_text or "Song" in all_text):
                    if found_section and sec.lower() != found_section.lower():
                        # Different section detected
                        is_different_section = True
                        break
            
            # Check if this slide has our hymn number (recompute since we need it here)
            has_our_hymn = False
            if target:
                hymn_pattern = build_hymn_pattern(target)
                has_our_hymn = bool(re.search(hymn_pattern, all_text, re.IGNORECASE))
            elif found_hymn_num:
                hymn_pattern = build_hymn_pattern(found_hymn_num)
                has_our_hymn = bool(re.search(hymn_pattern, all_text, re.IGNORECASE))
            
            # Only stop for different section if this slide doesn't have our hymn
            if is_different_section and not has_our_hymn:
                break
            
            has_our_title = False
            if title_search_phrase:
                # Check first line of all text boxes
                slide_first_lines = []
                for shape in prs.slides[i].shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        slide_first_lines.append(shape.text_frame.text.strip().split('\n')[0].strip())
                
                for first_line in slide_first_lines:
                    normalized_first_line = re.sub(r'[^a-z0-9\s]', '', first_line.lower())
                    if title_search_phrase in normalized_first_line:
                        has_our_title = True
                        break
            
            # Check if this slide has a DIFFERENT hymn number
            check_hymn = target if target else found_hymn_num
            if check_hymn:
                other_hymn = re.search(r"(?:Hymn|Song)\s*(?:[Nn]o\.?:?\s*)?[-–]?\s*\(?(\d+)\)?", all_text, re.IGNORECASE)
                if other_hymn:
                    found_num = other_hymn.group(1)
                    if found_num != check_hymn:
                        # Different hymn number found - stop collecting
                        break
            else:
                # No hymn number - check for any new "Hymn No" pattern without our title
                if re.search(r"(?:Hymn|Song)\s*(?:[Nn]o\.?:?\s*)?\(?\d+\)?", all_text, re.IGNORECASE):
                    if not has_our_title:
                        break
            
            # Use comprehensive title slide detection (visual + text analysis)
            # Even if slide has our hymn number, check if it's a title-only slide (no lyrics)
            if is_title_slide(all_text, slide):
                if not has_our_hymn:
                    # This is a title slide for a different song - stop collecting
                    break
                else:
                    # Has our hymn number - check if it's a title-only transition slide
                    # Title-only slides typically have: header + slide number only (< 100 chars total)
                    # Content slides have: header + slide number + section marker + lyrics (> 100 chars)
                    if len(all_text.strip()) < 100:
                        # This is a title-only transition slide - skip but continue
                        continue
            
            # Skip image-only slides (e.g., Holy Communion intro images with no lyrics)
            if is_image_only_slide(slide):
                continue
            
            # This slide belongs to our song - add as content slide
            content_indices.append(i)

    # If no title was extracted from title slide, try to extract from first content slide
    if not extracted_title and content_indices:
        first_content_idx = content_indices[0]
        max_text = ""
        for shape in prs.slides[first_content_idx].shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                # Skip footer patterns and slide numbers
                if re.search(r'^\d+\s*:\s*\d+\s+of\s+\d+', text):  # Footer like "30 : 31 of 106"
                    continue
                if re.match(r'^\d{1,3}$', text):  # Just a slide number
                    continue
                if any(kw in text for kw in ['Hymn', 'Song']) and len(text) < 30:  # Skip "Hymn No 313" headers
                    continue
                if any(sec in text for sec in ["Opening", "Confession", "Offertory", "Thanksgiving", "Communion", "Closing"]) and len(text) < 50:
                    continue
                # Check if text has English content
                has_english = len([c for c in text if c.isalpha() and c.isascii()]) > 10
                if has_english and len(text) > len(max_text):
                    max_text = text
        if max_text:
            # Take only the first line as the title
            first_line = max_text.split('\n')[0].strip()
            # Remove any leading punctuation or special chars
            first_line = re.sub(r'^[^\w\s]+', '', first_line)
            extracted_title = first_line if len(first_line) < 60 else first_line[:57] + "..."
    
    # If still no title, try to get from the slide BEFORE the first content slide (the title slide)
    if not extracted_title and title_slide_idx is not None and content_indices:
        # The title slide might have the song title in a different shape
        for shape in prs.slides[title_slide_idx].shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                # Skip very short text
                if len(text) < 10:
                    continue
                # Skip just hymn/song number references
                if re.match(r'^(Opening|Closing|Confession|Offertory|Thanksgiving|Communion)\s*(Hymn|Song)?\s*$', text, re.IGNORECASE):
                    continue
                if re.match(r'^(Hymn|Song)\s*[Nn]o\.?\s*:?\s*\d+\s*$', text, re.IGNORECASE):
                    continue
                # If we find text that's NOT just section/hymn labels, use it
                has_alpha = len([c for c in text if c.isalpha()]) > 5
                if has_alpha:
                    first_line = text.split('\n')[0].strip()
                    extracted_title = first_line if len(first_line) < 60 else first_line[:57] + "..."
                    break

    return title_slide_idx, content_indices, extracted_title


# Global tracking of used slide ranges to prevent duplicates
USED_SLIDE_RANGES = {}

# Global variable to store template presentation for cloning
TEMPLATE_REFERENCE_PRS = None

def find_best_song_source(hymn_num, song_name):
    """
    Search all PPT files and find the best source for an English song.
    Returns the source with the most content slides.
    Tracks used slide ranges to prevent adding the same hymn content twice.
    
    Args:
        hymn_num: The hymn number to search for
        song_name: Optional song title hint
    
    Returns: (pptx_path, title_idx, content_indices, extracted_title) or (None, None, [], "")
    """
    pptx_files = find_all_pptx_files(SEARCH_DIRS)
    
    best_source = None
    best_count = 0
    best_title_idx = None
    best_content = []
    best_extracted_title = ""
    
    # Search all English service PPT files
    for pf in pptx_files:
        t_idx, c_indices, extracted_title = find_song_slide_indices_in_pptx(pf, hymn_num, song_name)
        if c_indices and len(c_indices) > best_count:
            # Check if these slides were already used BY A DIFFERENT HYMN NUMBER
            # (Same hymn can be reused multiple times in one service)
            slide_key = f"{pf}:{min(c_indices)}-{max(c_indices)}"
            if slide_key not in USED_SLIDE_RANGES or USED_SLIDE_RANGES[slide_key] == hymn_num:
                best_source = pf
                best_count = len(c_indices)
                best_title_idx = t_idx
                best_content = c_indices
                best_extracted_title = extracted_title
    
    # Mark these slides as used if we found something
    if best_source and best_content:
        slide_key = f"{best_source}:{min(best_content)}-{max(best_content)}"
        USED_SLIDE_RANGES[slide_key] = hymn_num
    
    return best_source, best_title_idx, best_content, best_extracted_title


# ═══════════════════════════════════════════════════════════════════════════════
# COMMON HELPER FUNCTIONS FOR SLIDE FORMATTING
# ═══════════════════════════════════════════════════════════════════════════════

def add_header_footer_to_title_slide(slide, prs, service_date=""):
    """Add standard header and footer text to a title slide."""
    # Add header text at the TOP
    header_textbox = slide.shapes.add_textbox(
        Emu(0), Emu(150000),
        prs.slide_width, Emu(400000)
    )
    header_tf = header_textbox.text_frame
    header_tf.word_wrap = False
    header_p = header_tf.paragraphs[0]
    header_p.alignment = PP_ALIGN.CENTER
    header_run = header_p.add_run()
    header_run.text = "Mar Thoma Syrian Church, Singapore"
    header_run.font.name = "Gill Sans MT"
    header_run.font.size = Pt(18)
    header_run.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add footer text at the BOTTOM
    footer_textbox = slide.shapes.add_textbox(
        Emu(0), Emu(4700000),
        prs.slide_width, Emu(400000)
    )
    footer_tf = footer_textbox.text_frame
    footer_tf.word_wrap = False
    footer_p = footer_tf.paragraphs[0]
    footer_p.alignment = PP_ALIGN.CENTER
    footer_run = footer_p.add_run()
    # If no service date provided, use today's date
    if not service_date:
        service_date = datetime.now().strftime("%d %B %Y")
    footer_run.text = f"English Holy Communion Service – {service_date}"
    footer_run.font.name = "Gill Sans MT"
    footer_run.font.size = Pt(18)
    footer_run.font.color.rgb = RGBColor(255, 255, 255)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE CREATION FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════

def create_title_bar(slide, prs, title_text):
    """Create the title bar at the top of a content slide."""
    # Rectangle background with the maroon/burgundy color from image
    title_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Emu(0), Emu(0),
        prs.slide_width, TITLE_BAR_HEIGHT
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(144, 85, 98)  # Maroon color from title box
    title_bar.line.fill.background()
    
    # Text on bar
    title_textbox = slide.shapes.add_textbox(
        Emu(0), Emu(0),
        prs.slide_width, TITLE_BAR_HEIGHT
    )
    title_tf = title_textbox.text_frame
    title_tf.word_wrap = False
    title_tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertically center the text
    title_p = title_tf.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_run = title_p.add_run()
    title_run.text = title_text
    title_run.font.name = "Segoe UI"
    title_run.font.size = Pt(28)
    title_run.font.bold = True  # Make it bold
    title_run.font.color.rgb = RGBColor(255, 255, 255)  # White text


def clone_slide(prs, template_slide_index, template_prs):
    """Clone a slide from the template presentation to preserve its layout and background."""
    from pptx.oxml.xmlchemy import OxmlElement
    import copy
    
    template_slide = template_prs.slides[template_slide_index]
    
    # Use the template slide's layout directly
    # Import the layout from template to target presentation if not already there
    template_layout = template_slide.slide_layout
    
    # Try to find if this exact layout already exists in target prs
    matching_layout = None
    for layout in prs.slide_layouts:
        if layout._element == template_layout._element:
            matching_layout = layout
            break
    
    # If layout not found, use the template's layout (it should work across presentations from same master)
    target_layout = matching_layout if matching_layout else template_layout
    
    # Add new slide with the layout
    new_slide = prs.slides.add_slide(target_layout)
    
    # Copy all shapes from template slide, EXCEPT empty placeholders with default text
    for shape in template_slide.shapes:
        # Skip placeholder shapes that are empty or contain default "Click to add" text
        if shape.is_placeholder and shape.has_text_frame:
            try:
                text = shape.text_frame.text.strip()
                # Skip if empty or contains "click to add"
                if not text or "click to add" in text.lower():
                    continue  # Skip this placeholder
            except:
                pass  # If there's any error checking, include the shape
        
        el = shape.element
        newel = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    return new_slide


def add_title_slide(prs, layout, song_label, hymn_num, song_name="", template_prs=None, service_date="", show_box=True):
    """Add a title/intro slide for a song section using background images.
    
    Args:
        show_box: If True, adds a semi-transparent colored box behind the text.
                  If False (for Opening), text is directly on background.
    """
    global TEMPLATE_REFERENCE_PRS
    
    # Use blank layout to have full control
    blank_layout = prs.slide_layouts[11] if len(prs.slide_layouts) > 11 else prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Remove any placeholder shapes
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.is_placeholder:
            shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    
    # Add background image (full slide)
    bg_image_path = resolve_image_path("english_title_bg.png")
    if os.path.exists(bg_image_path):
        slide.shapes.add_picture(
            bg_image_path,
            Emu(0), Emu(0),
            prs.slide_width, prs.slide_height
        )
    
    # Add standard header and footer using helper function
    add_header_footer_to_title_slide(slide, prs, service_date)
    
    # Add colored box for non-Opening sections
    if show_box:
        box_left = Emu(1127800)
        box_top = Emu(1238150)
        box_width = Emu(6892935)
        box_height = Emu(2257244)
        
        colored_box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            box_left, box_top,
            box_width, box_height
        )
        colored_box.fill.solid()
        colored_box.fill.fore_color.rgb = RGBColor(130, 76, 88)  # Lighter maroon/burgundy color
        
        # Apply transparency via XML (python-pptx transparency property doesn't always work)
        try:
            from pptx.oxml.xmlchemy import OxmlElement
            spPr = colored_box._element.spPr
            solidFill = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
            if solidFill is not None:
                srgbClr = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                if srgbClr is not None:
                    # Remove any existing alpha
                    for alpha in srgbClr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}alpha'):
                        srgbClr.remove(alpha)
                    # Add alpha (70986 = 70.986% opacity)
                    alpha = OxmlElement('a:alpha')
                    alpha.set('val', '70986')
                    srgbClr.append(alpha)
        except Exception as e:
            print(f"    Warning: Could not apply transparency: {e}")
        
        # Remove border completely
        try:
            ln = colored_box._element.spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
            if ln is not None:
                noFill = OxmlElement('a:noFill')
                ln.clear()
                ln.append(noFill)
        except:
            colored_box.line.fill.background()
    
    # Add text in the center - positioned higher to align with box top
    text_left = Emu(1239143)
    text_top = Emu(1338150)  # Closer to box top (100000 EMUs padding from box top)
    text_width = Emu(6661177)
    text_height = Emu(2157244)  # Adjusted height to fit within box
    
    textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # First line: Section label (e.g., "Opening")
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1.level = 0
    run1 = p1.add_run()
    # Special case for ThanksGiving - show "ThanksGiving Prayers"
    if song_label == "ThanksGiving":
        run1.text = "ThanksGiving Prayers"
    else:
        run1.text = song_label
    run1.font.name = "Tenorite"
    run1.font.size = Pt(48)
    run1.font.bold = True  # Make it bold
    run1.font.color.rgb = RGBColor(255, 255, 255)
    # Add text shadow
    try:
        from pptx.oxml.xmlchemy import OxmlElement
        rPr = run1._r.get_or_add_rPr()
        shadow = OxmlElement('a:effectLst')
        outerShdw = OxmlElement('a:outerShdw')
        outerShdw.set('blurRad', '38100')  # Shadow blur radius
        outerShdw.set('dist', '38100')     # Shadow distance
        outerShdw.set('dir', '2700000')    # Shadow direction (bottom)
        outerShdw.set('algn', 'ctr')       # Alignment
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', '000000')       # Black shadow
        alpha = OxmlElement('a:alpha')
        alpha.set('val', '40000')          # 40% opacity
        srgbClr.append(alpha)
        outerShdw.append(srgbClr)
        shadow.append(outerShdw)
        rPr.append(shadow)
    except Exception:
        pass  # Shadow not critical
    
    # Second line: Song title
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.level = 0
    run2 = p2.add_run()
    
    if hymn_num:
        run2.text = f"Hymn No {hymn_num}"
    elif song_name:
        run2.text = song_name
    else:
        run2.text = ""
    
    run2.font.name = "Tenorite"
    run2.font.size = Pt(36)
    run2.font.bold = False
    run2.font.color.rgb = RGBColor(255, 255, 255)
    
    return slide


def add_message_slide(prs, layout, template_prs=None, service_date=""):
    """Add a Message title slide (no hymn content)."""
    # Use blank layout
    blank_layout = prs.slide_layouts[11] if len(prs.slide_layouts) > 11 else prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Remove placeholders
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.is_placeholder:
            shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    
    # Add background image
    bg_image_path = resolve_image_path("english_title_bg.png")
    if os.path.exists(bg_image_path):
        slide.shapes.add_picture(
            bg_image_path,
            Emu(0), Emu(0),
            prs.slide_width, prs.slide_height
        )
    
    # Add standard header and footer using helper function
    add_header_footer_to_title_slide(slide, prs, service_date)
    
    # Add text in center
    text_left = Emu(1239143)
    text_top = Emu(1654612)
    text_width = Emu(6661177)
    text_height = Emu(2904342)
    
    textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.level = 0
    run = p.add_run()
    run.text = "Message"
    run.font.name = "Tenorite"
    run.font.size = Pt(48)
    run.font.bold = True  # Make it bold
    run.font.color.rgb = RGBColor(255, 255, 255)
    # Add text shadow
    try:
        from pptx.oxml.xmlchemy import OxmlElement
        rPr = run._r.get_or_add_rPr()
        shadow = OxmlElement('a:effectLst')
        outerShdw = OxmlElement('a:outerShdw')
        outerShdw.set('blurRad', '38100')
        outerShdw.set('dist', '38100')
        outerShdw.set('dir', '2700000')
        outerShdw.set('algn', 'ctr')
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', '000000')
        alpha = OxmlElement('a:alpha')
        alpha.set('val', '40000')
        srgbClr.append(alpha)
        outerShdw.append(srgbClr)
        shadow.append(outerShdw)
        rPr.append(shadow)
    except Exception:
        pass  # Shadow not critical
    
    return slide


def add_holy_communion_intro_slide(prs, blank_layout, hymn_num, song_name=""):
    """Add a Holy Communion intro slide - blank slide with title bar and image."""
    slide = prs.slides.add_slide(blank_layout)
    
    # Remove title placeholders
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph = shape.placeholder_format
            if ph.type in [1, 2, 3]:
                shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    
    # Create the pink title bar with hymn info
    if hymn_num:
        title_text = f"Holy Communion - Hymn No {hymn_num}"
    elif song_name:
        title_text = f"Holy Communion - {song_name}"
    else:
        title_text = "Holy Communion"
    create_title_bar(slide, prs, title_text)

    # Add the Holy Communion image below the title bar
    hc_image_path = resolve_image_path("holy_communion.jpg")
    if os.path.exists(hc_image_path):
        slide.shapes.add_picture(
            hc_image_path,
            HC_IMAGE_LEFT, HC_IMAGE_TOP,
            HC_IMAGE_WIDTH, HC_IMAGE_HEIGHT
        )

    return slide


def clone_slides_from_source(source_pptx_path, slide_indices, target_prs,
                              song_label, hymn_num, slide_num_start,
                              blank_layout, song_name="", is_offertory=False):
    """
    Clone slides from source PPT, preserving text and images.
    For Offertory slides with overlapping text, adds QR code.
    Returns the number of slides added.
    """
    source_prs = Presentation(source_pptx_path)
    source_slides = list(source_prs.slides)
    added = 0
    
    # Build title text for content slides (just Hymn No and Title, no section label)
    if song_name:
        slide_title_text = f"Hymn No {hymn_num} - {song_name}"
    else:
        slide_title_text = f"Hymn No {hymn_num}"
    
    # QR code settings for Offertory slides
    QR_LEFT = Inches(6.97)       # Position on right side
    QR_TOP = Inches(1.08)        # Vertical position
    QR_WIDTH = Inches(3.04)      # QR code width
    QR_HEIGHT = Inches(3.13)     # QR code height
    
    # Check if QR code file exists
    qr_code_path = None
    if is_offertory:
        potential_path = resolve_image_path("qr_code.png")
        if os.path.exists(potential_path):
            qr_code_path = potential_path

    for idx in slide_indices:
        if idx >= len(source_slides):
            continue

        src_slide = source_slides[idx]
        
        # Clone the exact slide preserving all formatting
        new_slide = clone_slide_exact(src_slide, target_prs, blank_layout)

        # Always remove QR code and UEN from source slides (in case source was from offertory)
        remove_qr_and_uen(new_slide)
        remove_footer_text(new_slide)

        # Update title bar text based on which section this hymn is being added to
        update_title_bar_text(new_slide, song_label, hymn_num, song_name)

        # Only add QR code if this is an Offertory slide
        if is_offertory and qr_code_path:
            add_qr_code_to_slide(new_slide, qr_code_path, QR_LEFT, QR_TOP, QR_WIDTH, QR_HEIGHT)

        # Update title bar color to match template
        update_title_bar_color(new_slide)
        added += 1

    return added


def add_qr_code_to_slide(slide, qr_code_path, left, top, width, height):
    """Add QR code and UEN text to a slide."""
    try:
        slide.shapes.add_picture(
            qr_code_path,
            left, top,
            width, height
        )
        
        # Add UEN text below QR code
        uen_left = left
        uen_top = top + height + Inches(0.1)
        uen_width = width
        uen_height = Inches(0.3)
        
        uen_box = slide.shapes.add_textbox(uen_left, uen_top, uen_width, uen_height)
        uen_tf = uen_box.text_frame
        uen_tf.word_wrap = False
        uen_p = uen_tf.paragraphs[0]
        uen_p.alignment = PP_ALIGN.CENTER
        
        uen_run = uen_p.add_run()
        uen_run.text = "UEN - S86CC0315K"
        uen_run.font.name = "Arial"
        uen_run.font.size = Pt(10)
        uen_run.font.bold = False
    except Exception as e:
        print(f"    Warning: Could not add QR code: {e}")


def remove_qr_and_uen(slide):
    """Remove existing QR images and UEN text from a slide."""
    for shape in list(slide.shapes):
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.left.inches > 5:
                sp = shape._element
                sp.getparent().remove(sp)
                continue
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "UEN" in text:
                    sp = shape._element
                    sp.getparent().remove(sp)
        except Exception:
            continue


def remove_footer_text(slide):
    """Remove footer text like "B/A: 1 of 2" or "Communion 2: 1 of 7"."""
    footer_pattern = re.compile(r".+?:\s*\d+\s+of\s+\d+", re.IGNORECASE)
    for shape in list(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        if footer_pattern.fullmatch(text) or footer_pattern.search(text):
            sp = shape._element
            sp.getparent().remove(sp)


def clone_slide_exact(src_slide, target_prs, blank_layout):
    """Clone a slide preserving original shapes and formatting."""
    new_slide = target_prs.slides.add_slide(blank_layout)

    # Remove all default shapes/placeholders
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Copy shapes
    for shape in src_slide.shapes:
        el = deepcopy(shape._element)
        new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')

    # Copy relationships (images, etc.) - need to import media parts to target presentation
    for rel in src_slide.part.rels.values():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        
        if rel.is_external:
            # External relationships (e.g., hyperlinks) use target_ref string
            new_slide.part.rels._add_relationship(
                rel.reltype,
                rel.target_ref,
                rel.rId,
            )
        elif "image" in rel.reltype:
            # Image relationships - import the image part to target presentation
            try:
                source_part = rel._target
                # Import image/media part to target presentation
                # Wrap blob bytes in BytesIO to provide file-like interface
                image_stream = BytesIO(source_part.blob)
                image_part = target_prs.part.package.get_or_add_image_part(image_stream)
                # Add relationship from slide to the imported image (without rId to let it auto-generate)
                new_slide.part.relate_to(image_part, rel.reltype)
            except Exception as e:
                # Skip relationships we can't copy
                print(f"    ⚠ Warning: Could not copy image: {e}")
                continue

    return new_slide


def update_title_bar_text(slide, label, hymn_num, song_name=""):
    """Update the title bar text to match the section label and hymn number."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Find title bar by checking if it's at the top and contains hymn-related text
        if shape.top is None or shape.top > Emu(500000):  # Not in title bar area
            continue
        
        text = shape.text_frame.text.strip()
        # Check if this is a title bar (contains "Hymn" or section name)
        if "Hymn" not in text and "Song" not in text and "Offertory" not in text and "Opening" not in text and "Confession" not in text and "Communion" not in text and "Closing" not in text and "ThanksGiving" not in text:
            continue
        
        try:
            # Update the text to show section label with hymn number
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            # Ensure vertical centering
            shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            run = p.add_run()
            # Format: "Opening: Hymn No 40" (no song name if hymn num present)
            if hymn_num:
                run.text = f"{label}: Hymn No {hymn_num}"
            elif song_name:
                run.text = f"{label}: {song_name}"
            else:
                run.text = label
            
            run.font.name = "Segoe UI"
            run.font.size = Pt(28)
            run.font.bold = True  # Make it bold
            run.font.color.rgb = RGBColor(255, 255, 255)  # White text
        except Exception as e:
            pass
        break  # Only update the first matching title bar


def update_title_bar_color(slide, color=TITLE_BAR_COLOR):
    """Update the existing title bar fill color to match the template."""
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        if shape.top > Emu(100000):
            continue
        if shape.height > Emu(700000):
            continue
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
        except Exception:
            pass


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION PROCESSING FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════

def process_opening_song(prs, title_layout, blank_layout, hymn_num, song_name, slide_counter, service_date=None):
    """Process Opening Song section."""
    label = "Opening"
    
    # Find the best source (most slides) across all PPT files
    best_pf, t_idx, c_indices, extracted_title = find_best_song_source(hymn_num, song_name)
    
    if best_pf and c_indices:
        display_title = extracted_title if extracted_title else song_name
        print(f"  ✓ Found in PPT: {os.path.basename(best_pf)} ({len(c_indices)} content slides)")
        add_title_slide(prs, title_layout, label, hymn_num, display_title, TEMPLATE_REFERENCE_PRS, service_date)
        slide_counter += 1
        num_added = clone_slides_from_source(
            best_pf, c_indices, prs, label, hymn_num, slide_counter,
            blank_layout, display_title
        )
        slide_counter += num_added
        return slide_counter, display_title
    
    # Not found - add title only
    print(f"  ⚠ Song not found - adding title slide only")
    add_title_slide(prs, title_layout, label, hymn_num, song_name if song_name else "(Song not found)", TEMPLATE_REFERENCE_PRS, service_date, show_box=False)
    return slide_counter + 1, song_name


def process_thanksgiving_prayers(prs, title_layout, blank_layout, hymn_num, song_name, slide_counter, service_date=None):
    """Process ThanksGiving Prayers section (shown as B/A on summary, ThanksGiving on slides)."""
    label = "ThanksGiving"
    
    best_pf, t_idx, c_indices, extracted_title = find_best_song_source(hymn_num, song_name)
    
    if best_pf and c_indices:
        display_title = extracted_title if extracted_title else song_name
        print(f"  ✓ Found in PPT: {os.path.basename(best_pf)} ({len(c_indices)} content slides)")
        add_title_slide(prs, title_layout, label, hymn_num, display_title, TEMPLATE_REFERENCE_PRS, service_date, show_box=True)
        slide_counter += 1
        num_added = clone_slides_from_source(
            best_pf, c_indices, prs, label, hymn_num, slide_counter,
            blank_layout, display_title
        )
        slide_counter += num_added
        return slide_counter, display_title
    
    print(f"  ⚠ Song not found - adding title slide only")
    add_title_slide(prs, title_layout, label, hymn_num, song_name if song_name else "(Song not found)", TEMPLATE_REFERENCE_PRS, service_date, show_box=True)
    return slide_counter + 1, song_name


def process_offertory(prs, title_layout, blank_layout, hymn_num, song_name, slide_counter, service_date=None):
    """Process Offertory section (with QR code extraction)."""
    label = "Offertory"
    
    best_pf, t_idx, c_indices, extracted_title = find_best_song_source(hymn_num, song_name)
    
    if best_pf and c_indices:
        display_title = extracted_title if extracted_title else song_name
        print(f"  ✓ Found in PPT: {os.path.basename(best_pf)} ({len(c_indices)} content slides)")
        add_title_slide(prs, title_layout, label, hymn_num, display_title, TEMPLATE_REFERENCE_PRS, service_date, show_box=True)
        slide_counter += 1
        num_added = clone_slides_from_source(
            best_pf, c_indices, prs, label, hymn_num, slide_counter,
            blank_layout, display_title, is_offertory=True
        )
        slide_counter += num_added
        return slide_counter, display_title
    
    # Song not found - add title slide and QR code content slide
    print(f"  ⚠ Song not found - adding title slide with QR code")
    display_title = song_name if song_name else "(Song not found)"
    add_title_slide(prs, title_layout, label, hymn_num, display_title, TEMPLATE_REFERENCE_PRS, service_date, show_box=True)
    slide_counter += 1
    
    # Add a content slide with QR code
    content_slide = prs.slides.add_slide(blank_layout)
    
    # Remove any placeholder text boxes (like "Click to add title")
    shapes_to_remove = []
    for shape in content_slide.shapes:
        if shape.is_placeholder:
            shapes_to_remove.append(shape)
    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)
    
    # Create title bar using the standard function
    title_text = f"{label}: Hymn No {hymn_num}" if hymn_num else label
    create_title_bar(content_slide, prs, title_text)
    
    # Add QR code on the right side if available
    qr_code_path = None
    potential_path = os.path.join(IMAGES_DIR, "qr_code.png")
    if os.path.exists(potential_path):
        qr_code_path = potential_path
    
    if qr_code_path:
        # QR code position (right side of slide)
        qr_left = Inches(6.97)
        qr_top = Inches(1.08)
        qr_width = Inches(3.04)
        qr_height = Inches(3.13)
        add_qr_code_to_slide(content_slide, qr_code_path, qr_left, qr_top, qr_width, qr_height)
    
    slide_counter += 1
    return slide_counter, display_title


def process_message(prs, title_layout, slide_counter, service_date=None):
    """Process Message section (title slide only)."""
    add_message_slide(prs, title_layout, TEMPLATE_REFERENCE_PRS, service_date)
    print(f"  ✓ Added Message title slide")
    return slide_counter + 1


def process_confession(prs, title_layout, blank_layout, hymn_num, song_name, slide_counter, service_date=None):
    """Process Confession section."""
    label = "Confession"
    
    best_pf, t_idx, c_indices, extracted_title = find_best_song_source(hymn_num, song_name)
    
    if best_pf and c_indices:
        display_title = extracted_title if extracted_title else song_name
        print(f"  ✓ Found in PPT: {os.path.basename(best_pf)} ({len(c_indices)} content slides)")
        add_title_slide(prs, title_layout, label, hymn_num, display_title, TEMPLATE_REFERENCE_PRS, service_date, show_box=True)
        slide_counter += 1
        num_added = clone_slides_from_source(
            best_pf, c_indices, prs, label, hymn_num, slide_counter,
            blank_layout, display_title
        )
        slide_counter += num_added
        return slide_counter, display_title
    
    print(f"  ⚠ Song not found - adding title slide only")
    add_title_slide(prs, title_layout, label, hymn_num, song_name if song_name else "(Song not found)", TEMPLATE_REFERENCE_PRS, service_date, show_box=True)
    return slide_counter + 1, song_name


def process_holy_communion(prs, title_layout, blank_layout, hymn_num, song_name, slide_counter):
    """Process Holy Communion section - every song gets Holy Communion intro slide with image."""
    label = "Communion"
    
    # Ensure Holy Communion image exists
    ensure_holy_communion_image()
    
    best_pf, t_idx, c_indices, extracted_title = find_best_song_source(hymn_num, song_name)
    
    if best_pf and c_indices:
        display_title = extracted_title if extracted_title else song_name
        print(f"  ✓ Found in PPT: {os.path.basename(best_pf)} ({len(c_indices)} content slides)")
        
        # Add Holy Communion intro slide with image
        add_holy_communion_intro_slide(prs, blank_layout, hymn_num, display_title)
        slide_counter += 1
        
        num_added = clone_slides_from_source(
            best_pf, c_indices, prs, label, hymn_num, slide_counter,
            blank_layout, display_title
        )
        slide_counter += num_added
        return slide_counter, display_title
    
    print(f"  ⚠ Song not found - adding title slide only")
    add_holy_communion_intro_slide(prs, blank_layout, hymn_num, song_name if song_name else "(Song not found)")
    return slide_counter + 1, song_name


def process_closing_hymn(prs, title_layout, blank_layout, hymn_num, song_name, slide_counter, service_date=None):
    """Process Closing Hymn section."""
    label = "Closing"
    
    best_pf, t_idx, c_indices, extracted_title = find_best_song_source(hymn_num, song_name)
    
    if best_pf and c_indices:
        display_title = extracted_title if extracted_title else song_name
        print(f"  ✓ Found in PPT: {os.path.basename(best_pf)} ({len(c_indices)} content slides)")
        add_title_slide(prs, title_layout, label, hymn_num, display_title, TEMPLATE_REFERENCE_PRS, service_date, show_box=True)
        slide_counter += 1
        num_added = clone_slides_from_source(
            best_pf, c_indices, prs, label, hymn_num, slide_counter,
            blank_layout, display_title
        )
        slide_counter += num_added
        return slide_counter, display_title
    
    print(f"  ⚠ Song not found - adding title slide only")
    add_title_slide(prs, title_layout, label, hymn_num, song_name if song_name else "(Song not found)", TEMPLATE_REFERENCE_PRS, service_date, show_box=True)
    return slide_counter + 1, song_name


def process_generic_song(prs, title_layout, blank_layout, label, hymn_num, song_name, slide_counter, service_date=None):
    """Process any generic song section."""
    best_pf, t_idx, c_indices, extracted_title = find_best_song_source(hymn_num, song_name)
    
    if best_pf and c_indices:
        display_title = extracted_title if extracted_title else song_name
        print(f"  ✓ Found in PPT: {os.path.basename(best_pf)} ({len(c_indices)} content slides)")
        add_title_slide(prs, title_layout, label, hymn_num, display_title, TEMPLATE_REFERENCE_PRS, service_date, show_box=True)
        slide_counter += 1
        num_added = clone_slides_from_source(
            best_pf, c_indices, prs, label, hymn_num, slide_counter,
            blank_layout, display_title
        )
        slide_counter += num_added
        return slide_counter, display_title
    
    print(f"  ⚠ Song not found - adding title slide only")
    add_title_slide(prs, title_layout, label, hymn_num, song_name if song_name else "(Song not found)", TEMPLATE_REFERENCE_PRS, service_date, show_box=True)
    return slide_counter + 1, song_name


# ═══════════════════════════════════════════════════════════════════════════════
# SUMMARY SLIDE GENERATION
# ═══════════════════════════════════════════════════════════════════════════════

def create_summary_slide(prs, title_layout, song_list, service_date=None, template_prs=None):
    """
    Create a summary slide listing all songs in the service.
    
    Note: ThanksGiving is displayed as "B/A" on the summary slide only.
    Communion songs are grouped together with the label shown once.
    Uses centered layout with title and song list.
    """
    # Clone template slide to preserve background
    if template_prs:
        summary_slide = clone_slide(prs, 0, template_prs)
    else:
        summary_slide = prs.slides.add_slide(title_layout)
    
    # Remove all placeholder shapes (Click to add title, subtitle, etc.)
    shapes_to_remove = []
    for shape in summary_slide.shapes:
        if shape.is_placeholder:
            shapes_to_remove.append(shape)
    
    # Remove identified placeholders
    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)
    
    # Update the date in the footer if provided
    if service_date:
        import re
        # Find and update the footer text in the slide layout shapes
        for shape in summary_slide.slide_layout.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                # Look for the footer with date pattern "DD Month YYYY"
                date_pattern = r'\d{1,2}\s+(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}'
                if re.search(date_pattern, text) and "Holy Communion Service" in text:
                    # This is the footer - update the date
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if re.search(date_pattern, run.text):
                                run.text = re.sub(date_pattern, service_date, run.text)
                                break
    
    # Add title bar at the top (same as content slides)
    title_bar = summary_slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Emu(0), Emu(0),
        prs.slide_width, TITLE_BAR_HEIGHT
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(144, 85, 98)  # Maroon color
    title_bar.line.fill.background()
    
    # Add "Song list" text on the title bar
    title_textbox = summary_slide.shapes.add_textbox(
        Emu(0), Emu(0),
        prs.slide_width, TITLE_BAR_HEIGHT
    )
    title_tf = title_textbox.text_frame
    title_tf.word_wrap = False
    title_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    title_p = title_tf.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    
    title_run = title_p.add_run()
    title_run.text = "Song list"
    title_run.font.name = "Gill Sans MT"
    title_run.font.size = Pt(32)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(255, 255, 255)  # White text
    
    # Create centered song list textbox below title - start higher and use tighter spacing
    left = Inches(1.5)
    top = Inches(0.8)  # Start right below the title bar
    width = Inches(7)
    height = Inches(5.5)
    
    textbox = summary_slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    # Add custom marker to identify this textbox later
    textbox.name = "SummaryContent"
    
    # Build content with centered formatting - Label on one line, hymn number + title on next line
    last_label = None
    first_para = True
    
    for song in song_list:
        label_lower = song['label'].lower()
        is_communion = label_lower in ("communion", "holy communion")
        is_thanksgiving = label_lower in ("thanksgiving", "thanksgiving prayers", "b/a")
        hymn_num = str(song['hymn_num']) if song.get('hymn_num') else ""
        title_hint = song.get('title_hint', '')
        
        # Show label for new sections (not for communion continuation)
        if not (is_communion and last_label == "communion"):
            # Add label paragraph
            if first_para:
                p_label = tf.paragraphs[0]
                first_para = False
            else:
                p_label = tf.add_paragraph()
            
            p_label.alignment = PP_ALIGN.CENTER
            p_label.space_before = Pt(4)  # Reduced from 8
            p_label.space_after = Pt(1)   # Reduced from 2
            
            display_label = "B/A" if is_thanksgiving else song['label']
            
            run_label = p_label.add_run()
            run_label.text = display_label
            run_label.font.name = "Segoe UI (Body)"
            run_label.font.size = Pt(14)
            run_label.font.bold = True
        
        # Add hymn number and title paragraph (if there's a hymn number or title)
        if hymn_num or title_hint:
            p_hymn = tf.add_paragraph()
            p_hymn.alignment = PP_ALIGN.CENTER
            p_hymn.space_before = Pt(0)
            p_hymn.space_after = Pt(1)  # Reduced from 2
            
            # Hymn number
            if hymn_num:
                run_num = p_hymn.add_run()
                run_num.text = hymn_num
                run_num.font.name = "Segoe UI (Body)"
                run_num.font.size = Pt(14)
                run_num.font.bold = False
            
            # Title
            if title_hint:
                title_hint = title_hint.replace('\x0b', ' ').replace('\x00', '')
                first_line = title_hint.split('\n')[0].strip()
                if len(first_line) > 50:
                    first_line = first_line[:47] + "..."
                
                if hymn_num:
                    run_space = p_hymn.add_run()
                    run_space.text = " "
                    run_space.font.name = "Segoe UI (Body)"
                    run_space.font.size = Pt(14)
                
                run_title = p_hymn.add_run()
                run_title.text = first_line
                run_title.font.name = "Segoe UI (Body)"
                run_title.font.size = Pt(14)
                run_title.font.bold = False
        
        last_label = "communion" if is_communion else label_lower
    
    return summary_slide


def update_summary_slide_from_slides(prs, song_list):
    """Extract first line of English text from each hymn section and update summary."""
    if len(prs.slides) < 2:
        return
    
    # Scan slides to find first English line for each section
    hymn_titles = {}  # {hymn_num: first_line}
    fallback_titles = {}  # {hymn_num: first_line}
    current_hymn = None
    
    for slide_idx in range(1, len(prs.slides)):
        slide = prs.slides[slide_idx]
        
        # First pass: check if this slide has a section header (not a footer)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            
            # Skip lines that look like footers (contain " of " pattern)
            if " of " in text and ":" in text:
                continue
            
            # Look for section title pattern with hymn number
            # Patterns: "Opening Hymn No 231", "ThanksGiving Hymn- 236", "Communion Hymn- 242"
            section_match = re.search(r"(Opening|Thanksgiving|ThanksGiving|Offertory|Confession|Communion|Closing|Dedication|B/A)\s+(?:(?:Hymn|Song).*?\s+)?(\d+)", text)
            if section_match:
                hymn_num = section_match.group(2)
                # Accept hymn numbers with 2 or more digits (excludes single-digit footer numbers like "6")
                if len(hymn_num) >= 2:
                    current_hymn = hymn_num
                    if current_hymn not in hymn_titles:
                        hymn_titles[current_hymn] = ""
                else:
                    # Skip single-digit numbers like "6" from corrupted footers
                    pass
        
        # Second pass: if we're tracking a hymn, look for first line of English
        if current_hymn and hymn_titles[current_hymn] == "":
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text or len(text) < 5:
                    continue
                
                # Skip section headers and footers
                if re.search(r"(Opening|Thanksgiving|ThanksGiving|Offertory|Confession|Communion|Closing|B/A|Dedication)", text, re.IGNORECASE):
                    continue
                if re.match(r"^\d+\s*:\s*\d+\s+of\s+\d+", text):  # Footer like "1: 1 of 6"
                    continue
                if re.match(r"^(Trinity|Message|Holy|Mar Thoma)", text):
                    continue
                if "Hymn" in text or "Song" in text:
                    continue
                
                # Get first line (handle both \n and \x0b vertical tab)
                first_line = text.replace('\x0b', '\n').split('\n')[0].strip()
                
                if len(first_line) < 5 or len(first_line) > 120:
                    continue
                
                # Check if it's English (mostly Latin characters)
                basic_latin = sum(1 for c in first_line if (c >= 'a' and c <= 'z') or (c >= 'A' and c <= 'Z'))
                total_alpha = sum(1 for c in first_line if c.isalpha())
                
                if total_alpha < 10:
                    continue
                if basic_latin < total_alpha * 0.8:  # At least 80% Latin characters
                    continue
                
                # Extract only first 2-3 words for summary
                words = first_line.split()
                title_words = ' '.join(words[:3]) if len(words) >= 3 else ' '.join(words[:2])
                
                hymn_titles[current_hymn] = title_words
                break  # Found it for this hymn, move to next slide

        # Fallback: store a readable line even if English rules fail
        if current_hymn and hymn_titles[current_hymn] == "" and current_hymn not in fallback_titles:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text or len(text) < 5:
                    continue
                if re.search(r"(Opening|Thanksgiving|ThanksGiving|Offertory|Confession|Communion|Closing|Dedication|B/A)", text, re.IGNORECASE):
                    continue
                if re.match(r"^\d+\s*:\s*\d+\s+of\s+\d+", text):
                    continue
                if re.match(r"^(Trinity|Message|Holy|Mar Thoma)", text):
                    continue
                if "Hymn" in text or "Song" in text:
                    continue
                first_line = text.replace('\x0b', '\n').split('\n')[0].strip()
                if len(first_line) < 5:
                    continue
                words = first_line.split()
                title_words = ' '.join(words[:3]) if len(words) >= 3 else ' '.join(words[:2])
                fallback_titles[current_hymn] = title_words
                break
    
    # Update song_list with extracted titles (only if not already set)
    for song in song_list:
        hymn_num = str(song.get('hymn_num', ''))
        # Skip if title_hint already exists (was set during processing)
        if song.get('title_hint'):
            continue
        # Try to extract from slides
        if hymn_num in hymn_titles and hymn_titles[hymn_num]:
            song['title_hint'] = hymn_titles[hymn_num]
        elif hymn_num in fallback_titles and fallback_titles[hymn_num]:
            song['title_hint'] = fallback_titles[hymn_num]
    
    # Now rebuild summary with titles
    summary_slide = prs.slides[0]
    
    # Find the summary content textbox (marked with name "SummaryContent")
    summary_textbox = None
    for shape in summary_slide.shapes:
        if shape.name == "SummaryContent" and shape.has_text_frame:
            summary_textbox = shape
            break
    
    if not summary_textbox:
        return  # Can't update if textbox not found
    
    # Clear and rebuild content
    tf = summary_textbox.text_frame
    tf.clear()
    
    # Rebuild with extracted titles using the correct format: Label on one line, hymn+title on next
    last_label = None
    first_para = True
    
    for song in song_list:
        label = song.get('label', '')
        label_lower = label.lower()
        is_communion = label_lower in ('communion', 'holy communion')
        is_thanksgiving = label_lower in ('thanksgiving', 'thanksgiving prayers', 'b/a')
        hymn_num = str(song.get('hymn_num', ''))
        title_hint = song.get('title_hint', '')
        
        # Show label for new sections (not for communion continuation)
        if not (is_communion and last_label == "communion"):
            # Add label paragraph
            if first_para:
                p_label = tf.paragraphs[0]
                first_para = False
            else:
                p_label = tf.add_paragraph()
            
            p_label.alignment = PP_ALIGN.CENTER
            p_label.space_before = Pt(4)  # Reduced spacing
            p_label.space_after = Pt(1)
            
            display_label = "B/A" if is_thanksgiving else label
            
            run_label = p_label.add_run()
            run_label.text = display_label
            run_label.font.name = "Segoe UI (Body)"
            run_label.font.size = Pt(14)
            run_label.font.bold = True
        
        # Add hymn number and title paragraph
        if hymn_num or title_hint:
            p_hymn = tf.add_paragraph()
            p_hymn.alignment = PP_ALIGN.CENTER
            p_hymn.space_before = Pt(0)
            p_hymn.space_after = Pt(1)  # Reduced spacing
            
            # Hymn number
            if hymn_num:
                run_num = p_hymn.add_run()
                run_num.text = hymn_num
                run_num.font.name = "Segoe UI (Body)"
                run_num.font.size = Pt(14)
                run_num.font.bold = False
            
            # Title
            if title_hint:
                title_hint = title_hint.replace('\x0b', ' ').replace('\x00', '')
                first_line = title_hint.split('\n')[0].strip()
                if len(first_line) > 50:
                    first_line = first_line[:47] + "..."
                
                if hymn_num:
                    run_space = p_hymn.add_run()
                    run_space.text = " "
                    run_space.font.name = "Segoe UI (Body)"
                    run_space.font.size = Pt(14)
                
                run_title = p_hymn.add_run()
                run_title.text = first_line
                run_title.font.name = "Segoe UI (Body)"
                run_title.font.size = Pt(14)
                run_title.font.bold = False
        
        last_label = "communion" if is_communion else label_lower


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN GENERATION LOGIC
# ═══════════════════════════════════════════════════════════════════════════════

def generate_presentation(song_list, output_filename=None, service_date=None):
    """
    Generate a PPT presentation from a list of songs.
    
    song_list: list of dicts with keys:
        - hymn_num: int or str (the hymn number, or empty for Message)
        - label: str (Opening, ThanksGiving, Offertory, Message, Confession, Communion, Closing)
        - title_hint: str (optional song title hint)
    service_date: str (optional service date in format "DD Month YYYY", e.g., "08 February 2026")
    """
    # Reset used slide ranges for this generation
    global USED_SLIDE_RANGES
    USED_SLIDE_RANGES = {}
    
    # Refresh search directories based on current working directory
    global ENGLISH_SEARCH_DIRS
    ENGLISH_SEARCH_DIRS = get_search_dirs("English")
    
    if output_filename is None:
        today = datetime.now().strftime("%d %b %Y")
        output_filename = f"{today} - Generated English HCS.pptx"

    # Check if output_filename is already an absolute path (from web app)
    if os.path.isabs(output_filename):
        output_path = output_filename
    else:
        output_path = os.path.join(BASE_DIR, output_filename)
    
    print(f"  Language: English")
    print(f"  Search directories: {get_search_dirs('English')}")

    def normalize_service_date(date_text):
        if not date_text:
            return ""
        formats = [
            "%d %b %Y",
            "%d %B %Y",
            "%d-%b-%Y",
            "%d-%B-%Y",
            "%d/%m/%Y",
            "%d-%m-%Y",
        ]
        for fmt in formats:
            try:
                parsed = datetime.strptime(date_text, fmt)
                return parsed.strftime("%d %B %Y")
            except ValueError:
                continue
        return date_text

    # Load template and create presentation
    template_path = find_template_ppt()
    if not template_path:
        raise FileNotFoundError(
            "English Template PPT not found.\n\n"
            "Please ensure the source folder contains '8 Feb 2026.pptx' under the English HCS folder."
        )

    print(f"  Template: {template_path}")
    prs = Presentation(template_path)
    
    # Keep a reference template presentation for cloning slides
    global TEMPLATE_REFERENCE_PRS
    TEMPLATE_REFERENCE_PRS = Presentation(template_path)

    # Update date in slide masters/layouts if provided
    normalized_date = normalize_service_date(service_date)
    if normalized_date:
        date_pattern = re.compile(r"\d{1,2}\s+(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}")
        for master in prs.slide_masters:
            for shape in master.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if date_pattern.search(run.text):
                            run.text = date_pattern.sub(normalized_date, run.text)
        for layout in prs.slide_layouts:
            for shape in layout.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if date_pattern.search(run.text):
                            run.text = date_pattern.sub(normalized_date, run.text)

    # Remove all existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # Get layouts
    title_layout = None
    blank_layout = None
    
    for layout in prs.slide_layouts:
        # Look for English-specific layout names or use first/last as fallback
        if "English" in layout.name or "HC" in layout.name:
            title_layout = layout
        elif layout.name == "1_Blank" and blank_layout is None:
            blank_layout = layout

    if title_layout is None:
        title_layout = prs.slide_layouts[0]
    if blank_layout is None:
        blank_layout = prs.slide_layouts[-1]
    
    print(f"  Using title layout: '{title_layout.name}'")
    print(f"  Using content layout: '{blank_layout.name}'")

    slide_counter = 1

    # Create summary slide
    print("\n📋 Creating summary slide...")
    create_summary_slide(prs, title_layout, song_list, normalized_date)
    slide_counter += 1

    # Process each song section
    print("\n🎵 Processing songs...\n")

    for song_info in song_list:
        hymn_num = str(song_info["hymn_num"]) if song_info["hymn_num"] else ""
        label = song_info["label"]
        title_hint = song_info.get("title_hint", "")

        print(f"── {label}: Hymn No {hymn_num} {title_hint} ──" if hymn_num else f"── {label} ──")

        # Route to appropriate processor based on label
        label_lower = label.lower()
        
        # Each processor now returns (slide_counter, extracted_title)
        if label_lower == "opening":
            slide_counter, extracted_title = process_opening_song(prs, title_layout, blank_layout, hymn_num, title_hint, slide_counter, normalized_date)
            if extracted_title:
                song_info["title_hint"] = extracted_title
        elif label_lower in ("thanksgiving", "thanksgiving prayers", "b/a"):
            slide_counter, extracted_title = process_thanksgiving_prayers(prs, title_layout, blank_layout, hymn_num, title_hint, slide_counter, normalized_date)
            if extracted_title:
                song_info["title_hint"] = extracted_title
        elif label_lower == "offertory":
            slide_counter, extracted_title = process_offertory(prs, title_layout, blank_layout, hymn_num, title_hint, slide_counter, normalized_date)
            if extracted_title:
                song_info["title_hint"] = extracted_title
        elif label_lower == "message":
            slide_counter = process_message(prs, title_layout, slide_counter, normalized_date)
        elif label_lower == "confession":
            slide_counter, extracted_title = process_confession(prs, title_layout, blank_layout, hymn_num, title_hint, slide_counter, normalized_date)
            if extracted_title:
                song_info["title_hint"] = extracted_title
        elif label_lower in ("communion", "holy communion"):
            slide_counter, extracted_title = process_holy_communion(prs, title_layout, blank_layout, hymn_num, title_hint, slide_counter)
            if extracted_title:
                song_info["title_hint"] = extracted_title
        elif label_lower == "closing":
            slide_counter, extracted_title = process_closing_hymn(prs, title_layout, blank_layout, hymn_num, title_hint, slide_counter, normalized_date)
            if extracted_title:
                song_info["title_hint"] = extracted_title
        elif label_lower == "dedication":
            slide_counter, extracted_title = process_generic_song(prs, title_layout, blank_layout, "Dedication", hymn_num, title_hint, slide_counter, normalized_date)
            if extracted_title:
                song_info["title_hint"] = extracted_title
        else:
            slide_counter, extracted_title = process_generic_song(prs, title_layout, blank_layout, label, hymn_num, title_hint, slide_counter, normalized_date)
            if extracted_title:
                song_info["title_hint"] = extracted_title

    # Now update the summary slide with extracted titles from the created slides
    update_summary_slide_from_slides(prs, song_list)

    # Update PowerPoint document properties to set the correct title
    # This prevents the browser from using an embedded title when downloading
    try:
        core_props = prs.core_properties
        # Extract just the filename without path and extension for the title
        filename_only = os.path.splitext(os.path.basename(output_path))[0]
        core_props.title = filename_only
        core_props.subject = "English Holy Communion Service"
        core_props.author = "Mar Thoma Syrian Church, Singapore"
    except Exception as e:
        print(f"  Warning: Could not update document properties: {e}")

    # Save
    prs.save(output_path)
    print(f"\n{'═' * 60}")
    print(f"✅ Presentation saved: {output_path}")
    print(f"   Total slides: {slide_counter - 1}")
    print(f"{'═' * 60}")

    return output_path


# ═══════════════════════════════════════════════════════════════════════════════
# USER INTERFACE
# ═══════════════════════════════════════════════════════════════════════════════

def get_song_list_from_user():
    """Interactive prompt to get the song list from the user."""
    print("╔══════════════════════════════════════════════════════════╗")
    print("║       Church Songs PPT Generator                         ║")
    print("║       (English Holy Communion Service)                   ║")
    print("╚══════════════════════════════════════════════════════════╝")
    print()
    print("Enter songs one by one. For each song, provide:")
    print("  1. The hymn number (e.g., 91, or leave empty for Message)")
    print("  2. The label (Opening, ThanksGiving, Offertory, Message, Confession, Communion, Closing)")
    print("  3. Optionally, a title hint")
    print()
    print("Type 'done' when finished, or 'example' to see a sample.")
    print()

    songs = []
    while True:
        try:
            entry = input(f"Song {len(songs) + 1} (hymn_num label [title_hint]) or 'done': ").strip()
        except (EOFError, KeyboardInterrupt):
            print()
            break

        if entry.lower() == "done":
            break

        if entry.lower() == "example":
            print("  Example: 91 Opening Come Thou Fount")
            print("  Example: 110 ThanksGiving")
            print("  Example: 420 Offertory")
            print("  Example: Message")
            print("  Example: 211 Confession")
            print("  Example: 313 Communion")
            print("  Example: 427 Closing")
            continue

        if not entry:
            continue

        # Special case: Message (no hymn number)
        if entry.lower() == "message":
            songs.append({"hymn_num": "", "label": "Message", "title_hint": ""})
            print(f"  ✓ Added: Message")
            continue

        parts = entry.split(None, 2)
        if len(parts) < 2:
            print("  ⚠ Please provide at least: hymn_number label (or just 'Message')")
            continue

        hymn_num = parts[0]
        label = parts[1]
        title_hint = parts[2] if len(parts) > 2 else ""

        try:
            int(hymn_num)
        except ValueError:
            print(f"  ⚠ '{hymn_num}' is not a valid hymn number")
            continue

        songs.append({
            "hymn_num": hymn_num,
            "label": label.capitalize(),
            "title_hint": title_hint,
        })
        print(f"  ✓ Added: {label} - Hymn No {hymn_num} {title_hint}")

    return songs


def main():
    """Main entry point."""
    if len(sys.argv) > 1 and sys.argv[1] == "--batch":
        songs = []
        language = "English"  # Default language
        service_date = None  # Service date from batch file
        if len(sys.argv) > 2:
            batch_file = sys.argv[2]
            print(f"Reading songs from: {batch_file}")
            with open(batch_file, "r") as f:
                for line in f:
                    line = line.strip()
                    if not line:
                        continue
                    # Parse language directive
                    if line.startswith("# Language:"):
                        language = line.split(":", 1)[1].strip()
                        print(f"  Language set to: {language}")
                        continue
                    # Parse date directive
                    if line.lower().startswith("# date:") or line.lower().startswith("date:"):
                        service_date = line.split(":", 1)[1].strip()
                        print(f"  Service date set to: {service_date}")
                        continue
                    if line.startswith("#"):
                        continue
                    parts = line.split("|")
                    if len(parts) >= 2:
                        songs.append({
                            "hymn_num": parts[0].strip(),
                            "label": parts[1].strip(),
                            "title_hint": parts[2].strip() if len(parts) > 2 else "",
                        })
                    elif parts[0].strip().lower() == "message":
                        songs.append({"hymn_num": "", "label": "Message", "title_hint": ""})
        else:
            print("Usage: python3 generate_english_hcs_ppt.py --batch songs.txt [output.pptx]")
            return
    else:
        songs = get_song_list_from_user()
        language = "English"  # Default for interactive mode
        service_date = None

    if not songs:
        print("No songs specified. Exiting.")
        return

    print(f"\n📝 Song list ({len(songs)} songs):")
    for i, s in enumerate(songs, 1):
        if s['hymn_num']:
            print(f"  {i}. {s['label']}: Hymn No {s['hymn_num']} {s['title_hint']}")
        else:
            print(f"  {i}. {s['label']}")

    output_name = None
    is_batch = len(sys.argv) > 1 and sys.argv[1] == "--batch"
    
    if not is_batch:
        try:
            custom_name = input("\nOutput filename (press Enter for default): ").strip()
            if custom_name:
                output_name = custom_name
        except (EOFError, KeyboardInterrupt):
            pass
    
    if len(sys.argv) > 3:
        output_name = sys.argv[3]

    # Pass service_date if it was set from batch file
    if 'service_date' in locals() and service_date:
        generate_presentation(songs, output_name, service_date)
    else:
        generate_presentation(songs, output_name)


if __name__ == "__main__":
    main()
