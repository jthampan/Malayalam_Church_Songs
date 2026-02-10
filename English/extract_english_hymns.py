#!/usr/bin/env python3
"""
Extract English hymns from PowerPoint files and create an Excel report.
Lists hymn number, title, slide location, and slide count.
"""

import os
import re
from pptx import Presentation
import openpyxl
from openpyxl.styles import Font, Alignment, PatternFill
from openpyxl.utils import get_column_letter

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PARENT_DIR = os.path.dirname(BASE_DIR)

# Only search English directories
ENGLISH_SEARCH_DIRS = [
    os.path.join(PARENT_DIR, "onedrive_git_local", "Holy Communion Services - Slides", "English HCS"),
]


def find_all_pptx_files():
    """Find all PowerPoint files in English directories."""
    pptx_files = []
    for d in ENGLISH_SEARCH_DIRS:
        if os.path.isdir(d):
            for root, dirs, files in os.walk(d):
                for f in files:
                    if f.endswith(".pptx") and not f.startswith("~$"):
                        pptx_files.append(os.path.join(root, f))
    return pptx_files


def extract_hymn_number(text):
    """Extract hymn number from text."""
    patterns = [
        # Explicit "Hymn No" patterns with optional parentheses (1-4 digits)
        # Supports: "Hymn No 22", "Hymn No. 22", "Hymn No: 22", "Hymn No – 22"
        r'Hymn\s*No\.?[\:\-–]?\s*[\(\[]?(\d{1,4})[\)\]]?',
        r'Hymn\s*[-–]\s*(\d{1,4})\b',
        # Song No patterns - includes "Song. No:" format with period between words
        r'Song\.\s*No\.?\s*[\:\-–]\s*(\d{1,4})',  # "Song. No: 524" or "Song. No: 524 v1"
        r'Song\s*No\.?\s*[-–]\s*(\d{1,4})',  # "Song No. – 297" or "Song No: 1134"
        r'Song\s*No\.?\s*[\:\-–]?\s*(\d{1,4})',  # "Song No: 522, v1" or "Song no 208"
        r'Song\s+no\.?\s*[\:\-–]?\s*(\d{1,4})',  # lowercase "Song no"
        # Section with song number like "Offertory – 896 I Will Sing"
        r'(?:Offertory|Confession)\s*[-–]\s*(\d{1,4})\b',
        # Holy Communion with number like "Holy Communion – Song no. 650"
        r'(?:Holy\s+)?Communion\s*[-–]\s*(?:(?:Song|Hymn)\s+no\.?\s*)?(\d{1,4})(?:\s|$)',
        # Pattern like "Holy Communion – And Can It Be (42)"
        r'Holy\s+Communion\s*[-–]\s*[^(]*[\(\[](\d{1,4})[\)\]]',
        # Offertory/Communion as section title marker
        r'(?:Offertory|Confession)\s*[-–]\s*(?:Song\s+no\.?\s*)?(\d{1,4})(?:\s|$)',
        # General pattern: numbers in parentheses/brackets like "Song Title (42)" or "Prayer (105)"
        # NOTE: Skip if number is 1-4, as these are typically verse/part indicators
        (True, r'(?:\(|\[)(\d{1,4})(?:\)|\])'),  # Tuple: (skip_small_nums, pattern)
    ]
    
    for pattern_item in patterns:
        # Handle both tuple (for special patterns) and string (for normal patterns)
        if isinstance(pattern_item, tuple):
            skip_small = pattern_item[0]
            pattern = pattern_item[1]
        else:
            skip_small = False
            pattern = pattern_item
        
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            num = match.group(1)
            num_int = int(num)
            
            # Skip very small numbers in parentheses (verse/part indicators)
            if skip_small and num_int < 5:
                continue
            
            # Avoid false matches like slide numbers from "2 of 4"
            # Only skip if this number appears in a "X of Y" slide number pattern
            if re.search(r'\b\d+\s+of\s+\d+.*?\b' + re.escape(num) + r'\b', text):
                continue
            return num
    
    return None


def extract_hymn_number_from_slide(slide):
    """Extract hymn number from a slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        
        hymn_num = extract_hymn_number(text)
        if hymn_num:
            return hymn_num
    
    return None


def extract_title_from_slide_heading(slide):
    """
    Extract hymn title from slide heading/title bar.
    Looks for patterns like:
    - "Offertory – When I survey the wondrous cross"
    - "Opening Hymn – Amazing grace"
    - "Communion - O sacred head"
    - "Part 1: Daivame en daivame (140)"
    """
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text.strip()
                    
                    # Skip very short text
                    if len(text) < 10:
                        continue
                    
                    # Look for pattern: "Part X: Title (hymn_num)"
                    part_pattern = r'^Part\s+\d+\s*:\s*(.+?)\s*\(\d+\)$'
                    match = re.search(part_pattern, text, re.IGNORECASE)
                    if match:
                        title = match.group(1).strip()
                        if len(title) > 5:
                            return title
                    
                    # Look for pattern: SectionLabel – Title or SectionLabel - Title
                    section_patterns = [
                        r'^(?:Offertory|Opening\s+Hymn?|Closing\s+Hymn?|Communion|Holy\s+Communion|Confession|Thanksgiving)\s*[-–]\s*(.+)$',
                        r'^Hymn\s+(?:No\.?)?\s*\d+\s*[-–]\s*(.+)$',
                        r'^(?:Offertory|Opening\s+Hymn?|Closing\s+Hymn?|Communion|Holy\s+Communion|Confession|Thanksgiving)\s*[-–]\s*Song\.?\s*(?:No\.?)?\s*\d+\s*[-–]\s*(.+)$',
                    ]
                    
                    for pattern in section_patterns:
                        match = re.search(pattern, text, re.IGNORECASE)
                        if match:
                            title = match.group(1).strip()
                            # Clean up the title
                            # Remove trailing punctuation, slide numbers, etc.
                            title = re.sub(r'\s*\d+\s*:\s*\d+\s+of\s+\d+.*$', '', title)  # Remove "1:2 of 3"
                            title = re.sub(r'\s+\d+\s*$', '', title)  # Remove trailing numbers
                            title = re.sub(r'[\.!]+$', '', title)  # Remove trailing punctuation
                            
                            # Skip if title is just a hymn number indicator
                            if any(x in title.lower() for x in ['song.no:', 'song no:', 'hymn.no:', 'hymn no:']):
                                continue
                            
                            # Check if this looks like a valid title
                            if len(title) > 5 and len(title) < 100:
                                # Count letters
                                letters = sum(1 for c in title if c.isalpha())
                                if letters > 10:  # At least 10 letters
                                    return title
    
    # Fallback: try to extract from any text that looks like a title
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            # Look for lines that have title characteristics
            lines = text.split('\n')
            for line in lines:
                line = line.strip()
                
                # Check for "Part X: Title (num)" pattern across multiple lines
                part_match = re.search(r'Part\s+\d+\s*:\s*(.+?)\s*\(\d+\)', line, re.IGNORECASE)
                if part_match:
                    title = part_match.group(1).strip()
                    if len(title) > 5:
                        return title
                
                # Skip if it's just hymn numbers, dates, or very short
                if len(line) < 10 or re.match(r'^\d+$', line):
                    continue
                # Skip footer text
                if re.search(r'\d+\s*:\s*\d+\s+of\s+\d+', line):
                    continue
                # If it has mostly letters and some spaces, might be a title
                letters = sum(1 for c in line if c.isalpha())
                spaces = sum(1 for c in line if c.isspace())
                if letters > 15 and spaces >= 2:
                    # Clean and return
                    line = re.sub(r'\s+\d+\s*$', '', line)
                    if len(line) > 5:
                        return line[:60]  # Cap at 60 chars
    
    return ""


def extract_title_from_content(slide):
    """Extract title from slide content (first 3-4 words of lyrics)."""
    content_lines = []
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text or len(text) < 3:
                    continue
                
                # Skip if it's just a hymn number
                if re.match(r'^\d{1,3}$', text.strip()):
                    continue
                
                # Skip section headers and labels
                skip_patterns = [
                    r'^(hymn|offertory|communion|confession|thanksgiving|opening|closing)',
                    r'^holy\s+communion',
                    r'^(message|prayer|scripture|reading)',
                    r'^\d+\s*$',
                    r'^slide\s+\d+',
                    r'^uen\s*[-–]',
                    r'^song\s+no',
                    r'^hymn\s+no',
                    r'^song\.?\s+no',  # Additional: catch "Song. No:" format
                    r'^song\.?\s+no\.?[\:\-]',  # Catch "Song. No:", "Song No:"
                    r'^hymn\.?\s+no\.?[\:\-]',  # Catch "Hymn. No:", "Hymn No:"
                    r'^song\.no:',  # Catch "Song.no:" without space (compact format)
                    r'^hymn\.no:',  # Catch "Hymn.no:" without space (compact format)
                    r'^\s*–\s+(?:song|hymn)',  # Catch "– Song no:" or "– Hymn no:"
                    r'^\s*–\s+hymn\s+no:',  # Catch "– Hymn No:"
                    r'^\s*–\s+song\s*\.?\s*no:',  # Catch "– Song. No:" or "– Song No:"
                    r'(?:song|hymn)\s+no\.?[\:\-]\s*\d+',  # General hymn number indicators
                    r'^theme:',
                    r'^\d+\s+[a-z]+\s+\d{4}',
                    r'order\s+of\s+worship',
                    r'^dedication\s*[-–]',
                ]
                
                should_skip = any(
                    re.search(pattern, text.lower())
                    for pattern in skip_patterns
                )
                
                if should_skip:
                    continue
                
                # Check for Malayalam characters (reject those)
                malayalam_chars = sum(1 for c in text if '\u0D00' <= c <= '\u0D7F')
                if malayalam_chars > 0:
                    continue
                
                # Check for other non-Latin scripts
                non_ascii = sum(1 for c in text if ord(c) > 127)
                if len(text) > 0 and non_ascii / len(text) > 0.1:  # More than 10% non-ASCII
                    continue
                
                # This looks like English hymn content
                if len(text) > 3:
                    content_lines.append(text)
                    break
        
        if content_lines:
            break
    
    # Extract full first line as title (instead of just 3-4 words)
    if content_lines:
        full_text = content_lines[0]
        
        # Remove trailing semicolon, comma, or punctuation marks
        title = re.sub(r'[;,:.!?]+$', '', full_text).strip()
        
        # If it's too long (>80 chars), try to break at punctuation
        if len(title) > 80:
            for punct in [',', ';', ':', '—', '–', '-']:
                if punct in title:
                    title = title.split(punct)[0].strip()
                    break
        
        # Cap at 100 characters for Excel
        if len(title) > 100:
            title = title[:97] + "..."
        
        return title if len(title) > 3 else ""
    
    return ""


def extract_date_from_filename(filename):
    """Extract and parse date from filename for sorting.
    
    Handles formats like:
    - "8 Feb 2026.pptx"
    - "[13 Oct 2024] Eng HCS.pptx"
    - "23 Mar 2025 Eng HCS.pptx"
    - "HCS 23-02-25.pptx"
    - "25 Dec2024_HCS.pptx"
    - "24Aug 2025 Eng HCS.pptx"
    - "26May2024"
    - "22 Dec2024"
    - "17th Apr" (no year)
    """
    from datetime import datetime
    
    # Remove file extension and brackets
    name_without_ext = filename.replace('.pptx', '').strip('[]')
    
    # Try to find date patterns
    # Pattern 1: "DD Month YYYY" with proper spacing - "25 Feb 2024", "8 Feb 2026"
    month_pattern_spaces = r'(\d{1,2})\s+(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\.?\s+(\d{4})'
    match = re.search(month_pattern_spaces, name_without_ext, re.IGNORECASE)
    if match:
        day, month_str, year = match.groups()
        try:
            date_obj = datetime.strptime(f"{day} {month_str} {year}", "%d %b %Y")
            return date_obj
        except:
            pass
    
    # Pattern 2: Month and year without space between - "Dec2024", "May2024", "Feb2024"
    # Handles patterns like "22 Dec2024" or "25Feb2024"
    month_pattern_compact = r'(\d{1,2})\s*(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\.?\s*(\d{4})'
    match = re.search(month_pattern_compact, name_without_ext, re.IGNORECASE)
    if match:
        day, month_str, year = match.groups()
        try:
            date_obj = datetime.strptime(f"{day} {month_str} {year}", "%d %b %Y")
            return date_obj
        except:
            pass
    
    # Pattern 3: "DD-MM-YY" or "DD-MM-YYYY"
    dash_pattern = r'(\d{1,2})[\-_](\d{1,2})[\-_](\d{2,4})'
    match = re.search(dash_pattern, name_without_ext)
    if match:
        day, month, year = match.groups()
        # Handle 2-digit years
        if len(year) == 2:
            year = '20' + year
        try:
            # Try DD-MM-YY format first
            date_obj = datetime.strptime(f"{day}-{month}-{year}", "%d-%m-%Y")
            return date_obj
        except:
            pass
    
    # If no date found, return a very old date to sort to the beginning
    return datetime(1900, 1, 1)


def clean_text_for_excel(text):
    """Remove illegal characters for Excel cells."""
    if not text:
        return ""
    cleaned = ''.join(char for char in text if ord(char) >= 32 or char in '\t\n\r')
    if len(cleaned) > 200:
        cleaned = cleaned[:200] + "..."
    return cleaned


def is_section_title_slide(slide):
    """
    Check if slide is a section title slide (like "Opening Hymn O Day of rest").
    Returns (True, section, title) if match found, otherwise (False, None, None).
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        
        # Look for patterns like:
        # "Opening Hymn O Day of rest" (same line)
        # or "Opening Hymn\nO Day of rest" (multi-line)
        section_patterns = [
            (r'^(Opening\s+Hymn?)\s+(.{5,})$', re.IGNORECASE | re.MULTILINE),
            (r'^(Closing\s+Hymn?)\s+(.{5,})$', re.IGNORECASE | re.MULTILINE),
            (r'^(Offertory)\s+(.{5,})$', re.IGNORECASE | re.MULTILINE),
            (r'^(Holy\s+Communion)\s+(.{5,})$', re.IGNORECASE | re.MULTILINE),
            (r'^(Communion)\s+(.{5,})$', re.IGNORECASE | re.MULTILINE),
            (r'^(Confession)\s+(.{5,})$', re.IGNORECASE | re.MULTILINE),
            (r'^(Thanksgiving)\s+(.{5,})$', re.IGNORECASE | re.MULTILINE),
        ]
        
        for pattern, flags in section_patterns:
            match = re.search(pattern, text, flags)
            if match:
                section = match.group(1).strip()
                title = match.group(2).strip()
                
                # Clean up title - remove footer info
                title = re.sub(r'\s*\d+\s*:\s*\d+\s+of\s+\d+.*$', '', title)
                title = re.sub(r'\d{1,2}\s+[A-Z][a-z]+\s+\d{4}.*$', '', title)  # Remove dates
                
                # Check if this looks like a valid title
                letters = sum(1 for c in title if c.isalpha())
                if letters > 5 and len(title) < 80:
                    return True, section, title
        
        # Also check for multi-line pattern: "Opening Hymn" on one line, title on next
        lines = text.split('\n')
        for i in range(len(lines) - 1):
            line1 = lines[i].strip()
            line2 = lines[i + 1].strip()
            
            section_labels = [
                'Opening Hymn', 'Opening Hymn', 'Closing Hymn', 'Closing Hymn',
                'Offertory', 'Holy Communion', 'Communion', 'Confession', 'Thanksgiving'
            ]
            
            for label in section_labels:
                if re.match(f'^{re.escape(label)}$', line1, re.IGNORECASE):
                    # Next line should be the title
                    if len(line2) > 5:
                        # Clean up title
                        title = re.sub(r'\s*\d+\s*:\s*\d+\s+of\s+\d+.*$', '', line2)
                        title = re.sub(r'\d{1,2}\s+[A-Z][a-z]+\s+\d{4}.*$', '', title)
                        letters = sum(1 for c in title if c.isalpha())
                        if letters > 5 and len(title) < 80:
                            return True, label, title
    
    return False, None, None


def analyze_pptx_file(pptx_path):
    """
    Analyze a PowerPoint file and extract English hymn information.
    Returns a list of hymn dictionaries.
    """
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"  ⚠ Could not open {pptx_path}: {e}")
        return []
    
    file_name = os.path.basename(pptx_path)
    hymns = []
    current_hymn = None
    
    for slide_idx, slide in enumerate(prs.slides, start=1):
        # Extract all text from slide
        all_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += " " + shape.text_frame.text.strip()
        
        # Skip slides with very little content
        if len(all_text.strip()) < 5:
            if current_hymn and len(current_hymn['content_slides']) > 0:
                hymns.append(current_hymn)
                current_hymn = None
            continue
        
        # Skip "Song list" / "Order of Worship" metadata slides
        # These list songs but are NOT the actual hymn slides
        metadata_patterns = [
            r'^\s*song\s+list\b',
            r'order\s+of\s+(?:worship|service)',
            r'^\s*opening\s*:\s*\d+\s+',  # "Opening: 1 of 3"
        ]
        is_metadata = any(re.search(pattern, all_text, re.IGNORECASE) for pattern in metadata_patterns)
        if is_metadata:
            if current_hymn and len(current_hymn['content_slides']) > 0:
                hymns.append(current_hymn)
                current_hymn = None
            continue
        
        # Look for hymn number on this slide FIRST (before filtering)
        hymn_num = extract_hymn_number_from_slide(slide)
        
        # Only skip non-hymn slides if they DON'T have a hymn number
        # This ensures "Thanksgiving Prayers Hymn No 306" is recognized as a hymn
        if not hymn_num:
            non_hymn_patterns = [
                r'^prayer\b',
                r'^prayers?\b',
                r'thanksgiving\s+prayer',
                r'intercessory\s+prayer',
                r'\bprayer\b.*\d+\s*:\s*\d+\s+of\s+\d+',  # "Prayer 1 : 1 of 3"
                r'\bresponse\b',
                r'\b[LC]\s*[-–]\s*for\b',
                r'\bleader\b',
                r'\bcongregation\b',
                r'\bdedication\b.*\d+\s+of\s+\d+',
                r'\breading\b',
                r'\bscripture\b',
                r'(?:hymns?|song)\s+list',  # "hymn list", "hymns list", or "song list"
                r'^b/?a:',
                r'^youtube\.be',
                r'^theme:',
                r'^announcements?\b',
                r'^message\b',
            ]
            
            is_non_hymn = any(re.search(pattern, all_text, re.IGNORECASE) for pattern in non_hymn_patterns)
            if is_non_hymn:
                if current_hymn and len(current_hymn['content_slides']) > 0:
                    hymns.append(current_hymn)
                    current_hymn = None
                continue
        
        if hymn_num:
            # Check if this is the same hymn number as current_hymn (continuation)
            if current_hymn and current_hymn.get('hymn_number') == hymn_num:
                # Same hymn, just continue adding slides
                current_hymn['content_slides'].append(slide_idx)
                current_hymn['total_slides'] += 1
                
                # If we don't have a title yet, try to extract from this slide
                if not current_hymn['title']:
                    title = extract_title_from_content(slide)
                    if title:
                        current_hymn['title'] = title
            else:
                # Beginning of a NEW hymn WITH hymn number
                if current_hymn and len(current_hymn['content_slides']) > 0:
                    hymns.append(current_hymn)
                
                # Extract title from slide content or hymn header
                title = extract_title_from_content(slide)
                
                # Also check for title in the same text as hymn number
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        # Look for pattern like "Hymn 171 - Title" or "Title - Hymn 171"
                        title_match = re.search(
                            r'Hymn\s+(?:No\.?)?\s*\d+\s*[-–]\s*([A-Za-z][A-Za-z\s]+?)(?:\s*$|\s+\d)',
                            text,
                            re.IGNORECASE
                        )
                        if title_match and not title:
                            title = title_match.group(1).strip()
                            if len(title) > 30:
                                title = title[:30]
                
                current_hymn = {
                    'hymn_number': hymn_num,
                    'title': title,
                    'file_name': file_name,
                    'title_slide': slide_idx,
                    'content_slides': [slide_idx] if title or extract_title_from_content(slide) else [],
                    'total_slides': 1,
                    'slides_to_check_for_title': [],  # Track slides to check if title is missing
                }
        else:
            # No hymn number, but check if this is a section title slide
            is_section_title, section, title_from_section = is_section_title_slide(slide)
            
            if is_section_title:
                # Section title slide found - always start a new hymn
                # First, save any current hymn
                if current_hymn and len(current_hymn['content_slides']) > 0:
                    hymns.append(current_hymn)
                
                # Start a new hymn WITHOUT hymn number (will look for it in next slides)
                current_hymn = {
                    'hymn_number': '',  # Will try to extract from content slides
                    'title': title_from_section,
                    'file_name': file_name,
                    'title_slide': slide_idx,
                    'content_slides': [slide_idx],
                    'total_slides': 1,
                    'slides_to_check_for_title': [],
                    'needs_hymn_number': True,  # Flag to extract number from content
                }
            elif current_hymn:
                # Continuation of current hymn
                # Check if this slide still has hymn content (not a new section)
                section_keywords = [
                    'Opening', 'Closing', 'Offertory', 'Communion',
                    'Confession', 'Thanksgiving', 'Message', 'Prayer'
                ]
                has_section_keyword = any(
                    re.search(f'^{kw}', all_text, re.IGNORECASE)
                    for kw in section_keywords
                )
                
                if has_section_keyword and not hymn_num:
                    # New section without hymn number - end current hymn
                    if len(current_hymn['content_slides']) > 0:
                        hymns.append(current_hymn)
                    current_hymn = None
                else:
                    # Continue collecting slides for current hymn
                    current_hymn['content_slides'].append(slide_idx)
                    current_hymn['total_slides'] += 1
                    
                    # If we need a hymn number, try to extract from this content slide
                    if current_hymn.get('needs_hymn_number') and not current_hymn['hymn_number']:
                        # Look for patterns like "Part 1: Title (140)" or just "(140)"
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text = shape.text_frame.text
                                # Pattern: (hymn_num) or [hymn_num]
                                num_match = re.search(r'[\(\[](\d{1,3})[\)\]]', text)
                                if num_match:
                                    current_hymn['hymn_number'] = num_match.group(1)
                                    current_hymn['needs_hymn_number'] = False
                                    break
                    
                    # If we don't have a title yet, try to extract from this content slide heading
                    if not current_hymn['title']:
                        current_hymn['slides_to_check_for_title'].append(slide_idx - 1)  # Store 0-indexed
    
    # Don't forget the last hymn
    if current_hymn and len(current_hymn['content_slides']) > 0:
        hymns.append(current_hymn)
    
    # Post-process: Extract titles from content slides if title is missing
    for hymn in hymns:
        if not hymn['title'] and hymn['slides_to_check_for_title']:
            try:
                prs_reopen = Presentation(pptx_path)
                for slide_idx_0 in hymn['slides_to_check_for_title']:
                    if slide_idx_0 < len(prs_reopen.slides):
                        slide = prs_reopen.slides[slide_idx_0]
                        # Try to extract title from slide heading
                        title = extract_title_from_slide_heading(slide)
                        if title:
                            hymn['title'] = title
                            break
            except Exception:
                pass
    
    # Clean up titles: remove hymn indicators and keep only real titles
    for hymn in hymns:
        title = hymn.get('title', '').strip()
        if title:
            # Skip entries that are just hymn indicators
            if any(x in title.lower() for x in ['song.no:', 'song no:', 'hymn.no:', 'hymn no:']):
                # Try to extract the real title after the indicator
                # Example: "Song.no: 522 (v1)" -> look for content afterthe hymn number
                match = re.search(
                    r'(?:song|hymn)\.?\s*no\.?[\:\-]\s*\d+[,\s]*(.+)',
                    title,
                    re.IGNORECASE
                )
                if match:
                    real_title = match.group(1).strip()
                    # Clean up version markers like "(v1)"
                    real_title = re.sub(r'\s*\(v\d+\)\s*', '', real_title)
                    if real_title and len(real_title) > 3:
                        hymn['title'] = real_title
                    else:
                        # No real title found, clear it
                        hymn['title'] = ''
                else:
                    # Just hymn indicator with no content after
                    hymn['title'] = ''
    
    # Deduplicate: Keep only first occurrence of each hymn
    # - If hymn has a number, deduplicate by number
    # - If hymn has no number, deduplicate by title (for multi-part hymns like "God of mercy" parts 1-4)
    seen_hymns = set()
    seen_normalized_titles = set()  # Track normalized titles to catch duplicates with different hymn numbers
    deduplicated = []
    for hymn in hymns:
        hymn_num = hymn.get('hymn_number', '')
        hymn_title = hymn.get('title', '').strip()
        
        # Create a dedup key: use number if available, otherwise use title
        if hymn_num:
            dedup_key = f"num:{hymn_num}"
        elif hymn_title:
            # Normalize title for case-insensitive, dash-insensitive comparison:
            # - Remove version markers like "v1", "v2"
            normalized_title = re.sub(r'\s+v\d+\s*$', '', hymn_title, flags=re.IGNORECASE)
            # - Normalize all dash characters (-, –, —) to a single dash
            normalized_title = re.sub(r'[-–—]+', '-', normalized_title)
            # - Remove quotes after dashes using Unicode class for quotes
            # Matches: -", –", -", etc. (with regular and curly quotes)
            normalized_title = re.sub(r'[-–]\s*["\u201c\u201d]', '- ', normalized_title)
            # - Remove leading dashes, spaces, and all quote types (regular and curly)
            normalized_title = re.sub(r'^[-–\s"\u201c\u201d\'\'\'\']+', '', normalized_title)
            # - Remove trailing quotes (regular and curly)
            normalized_title = re.sub(r'["\u201c\u201d\'\'\'\']+$', '', normalized_title)
            # - Normalize whitespace (multiple spaces to single space)
            normalized_title = re.sub(r'\s+', ' ', normalized_title)
            # - Lowercase for case-insensitive comparison
            normalized_title = normalized_title.strip().lower()
            dedup_key = f"title:{normalized_title}"
        else:
            dedup_key = None
        
        # Skip if we've already seen this hymn by number or title
        if dedup_key and dedup_key in seen_hymns:
            continue
        
        # Also skip if we've already seen this normalized title (catches duplicates with different hymn numbers)
        # For example: one entry "Hymn 45: O for a thousand tongues" and another "– O for a thousand tongues"
        if hymn_title:
            # Extract normalized title for cross-check
            normalized_title = re.sub(r'\s+v\d+\s*$', '', hymn_title, flags=re.IGNORECASE)
            normalized_title = re.sub(r'[-–—]+', '-', normalized_title)
            normalized_title = re.sub(r'[-–]\s*["\u201c\u201d]', '- ', normalized_title)
            normalized_title = re.sub(r'^[-–\s"\u201c\u201d\'\'\'\']+', '', normalized_title)
            normalized_title = re.sub(r'["\u201c\u201d\'\'\'\']+$', '', normalized_title)
            normalized_title = re.sub(r'\s+', ' ', normalized_title)
            normalized_title = normalized_title.strip().lower()
            
            if normalized_title in seen_normalized_titles:
                continue  # Skip this duplicate
            
            seen_normalized_titles.add(normalized_title)
        
        if dedup_key:
            seen_hymns.add(dedup_key)
        
        deduplicated.append(hymn)
    
    return deduplicated


def create_excel_report(all_hymns, output_file):
    """Create an Excel report from the hymn data with multiple sort views."""
    # Create workbook with multiple sheets
    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # Remove default sheet
    
    # Sheet 1: Sort by Hymn Number
    hymns_by_num = sorted(all_hymns, key=lambda x: (
        int(x['hymn_number']) if x['hymn_number'] else 999999,
        x['title'] if not x['hymn_number'] else ''
    ))
    _add_hymns_sheet(wb, hymns_by_num, "By Hymn Number")
    
    # Sheet 2: Sort by Date (extracted from filename) then Hymn Number
    hymns_by_file = sorted(all_hymns, key=lambda x: (
        extract_date_from_filename(x['file_name']),
        int(x['hymn_number']) if x['hymn_number'] else 999999,
    ))
    _add_hymns_sheet(wb, hymns_by_file, "By Date")
    
    # Save workbook
    wb.save(output_file)
    print(f"\n✓ Excel report saved to: {output_file}")
    print(f"  - Sheet 1: 'By Hymn Number' (sorted by hymn number)")
    print(f"  - Sheet 2: 'By Date' (sorted by date from filename)")


def _add_hymns_sheet(wb, hymns, sheet_name):
    """Add a sheet with hymn data to the workbook."""
    ws = wb.create_sheet(sheet_name)
    
    # Define headers
    headers = ['Hymn Number', 'Title', 'File Name']
    ws.append(headers)
    
    # Style headers
    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF", size=12)
    
    for col in range(1, len(headers) + 1):
        cell = ws.cell(row=1, column=col)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal='center', vertical='center')
    
    # Add data
    for hymn in hymns:
        row_data = [
            int(hymn['hymn_number']) if hymn['hymn_number'] else '',
            clean_text_for_excel(hymn['title']),
            clean_text_for_excel(hymn['file_name'])
        ]
        ws.append(row_data)
    
    # Auto-adjust column widths
    for col in range(1, len(headers) + 1):
        column_letter = get_column_letter(col)
        max_length = 0
        
        for cell in ws[column_letter]:
            try:
                if len(str(cell.value)) > max_length:
                    max_length = len(str(cell.value))
            except:
                pass
        
        adjusted_width = min(max_length + 2, 60)
        ws.column_dimensions[column_letter].width = adjusted_width
    
    # Center align hymn number column
    for row in ws.iter_rows(min_row=2, max_row=ws.max_row):
        row[0].alignment = Alignment(horizontal='center')


def main():
    print("=" * 70)
    print("Extracting English Hymn Information from PowerPoint Files")
    print("=" * 70)
    
    # Find all PowerPoint files
    print("\n🔍 Searching for English PowerPoint files...")
    pptx_files = find_all_pptx_files()
    print(f"  Found {len(pptx_files)} PowerPoint files")
    
    if not pptx_files:
        print(f"\n⚠ No PowerPoint files found in English HCS directories:")
        for d in ENGLISH_SEARCH_DIRS:
            print(f"  - {d}")
        return
    
    # Extract hymn information from each file
    all_hymns = []
    for pptx_file in pptx_files:
        print(f"\n📄 Analyzing: {os.path.basename(pptx_file)}")
        hymns = analyze_pptx_file(pptx_file)
        print(f"  Found {len(hymns)} hymns")
        all_hymns.extend(hymns)
    
    print(f"\n📊 Total hymn entries found across all files: {len(all_hymns)}")
    
    if not all_hymns:
        print("\n⚠ No hymns found!")
        return
    
    # Create Excel report
    output_file = os.path.join(BASE_DIR, "english_hymns_report.xlsx")
    create_excel_report(all_hymns, output_file)
    
    # Count unique hymn numbers
    unique_hymn_nums = len(set(h['hymn_number'] for h in all_hymns if h['hymn_number']))
    
    # Print summary
    print("\n" + "=" * 70)
    print("Summary:")
    print("=" * 70)
    print(f"Total PowerPoint files analyzed: {len(pptx_files)}")
    print(f"Total hymn entries found: {len(all_hymns)}")
    print(f"Unique hymn numbers: {unique_hymn_nums}")
    
    if all_hymns:
        hymn_nums = [int(h['hymn_number']) for h in all_hymns if h['hymn_number']]
        if hymn_nums:
            print(f"Hymn numbers range: {min(hymn_nums)} - {max(hymn_nums)}")
    
    print(f"Output file: {output_file}")
    print("=" * 70)


if __name__ == "__main__":
    main()
